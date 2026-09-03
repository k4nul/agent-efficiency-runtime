from __future__ import annotations

import io
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from aer.errors import AerError
from aer.limits import (
    MAX_ARTIFACT_ELEMENTS,
    MAX_IMAGE_PIXELS,
    MAX_PRESENTATION_SLIDES,
)
from aer.paths import atomic_write_bytes, ensure_regular_input

SLIDE_W = 40 / 3
SLIDE_H = 7.5
COLORS = {
    "navy": RGBColor(25, 42, 67),
    "blue": RGBColor(39, 105, 193),
    "teal": RGBColor(25, 145, 140),
    "light": RGBColor(243, 246, 250),
    "white": RGBColor(255, 255, 255),
    "text": RGBColor(32, 39, 49),
    "muted": RGBColor(91, 103, 118),
    "accent": RGBColor(238, 151, 43),
}
SUPPORTED_LAYOUTS = {
    "title",
    "section",
    "bullets",
    "two-column",
    "comparison",
    "metrics",
    "table",
    "image",
    "image-with-caption",
    "chart",
    "quote",
    "timeline",
    "closing",
}


def _position(
    block: dict[str, Any], default: tuple[float, float, float, float]
) -> tuple[float, float, float, float]:
    override = block.get("position")
    if not isinstance(override, dict):
        return default
    return tuple(
        float(override.get(key, value))
        for key, value in zip(("x", "y", "w", "h"), default, strict=True)
    )  # type: ignore[return-value]


def _text_box(
    slide: Any,
    slide_id: str,
    shape_id: str,
    text: str,
    box: tuple[float, float, float, float],
    *,
    size: float = 24,
    bold: bool = False,
    color: RGBColor = COLORS["text"],
    align: PP_ALIGN = PP_ALIGN.LEFT,
    fill: RGBColor | None = None,
) -> Any:
    x, y, width, height = box
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    shape.name = f"aer:{slide_id}/{shape_id}"
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(text)
    paragraph.alignment = align
    font = run.font
    font.name = "Noto Sans CJK KR"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
    return shape


def _add_title(slide: Any, slide_id: str, title: str) -> Any:
    return _text_box(
        slide,
        slide_id,
        "title",
        title,
        (0.7, 0.35, 11.9, 0.8),
        size=28,
        bold=True,
        color=COLORS["navy"],
    )


def _bullet_box(
    slide: Any,
    slide_id: str,
    shape_id: str,
    items: list[Any],
    box: tuple[float, float, float, float],
) -> Any:
    shape = _text_box(slide, slide_id, shape_id, "", box, size=20)
    frame = shape.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        if isinstance(item, dict):
            paragraph.text = str(item.get("text", ""))
            paragraph.level = int(item.get("level", 0))
        else:
            paragraph.text = str(item)
            paragraph.level = 0
        paragraph.text = f"• {paragraph.text}"
        paragraph.font.name = "Noto Sans CJK KR"
        paragraph.font.size = Pt(20 if paragraph.level == 0 else 17)
        paragraph.font.color.rgb = COLORS["text"]
        paragraph.space_after = Pt(9)
    return shape


def _add_picture(
    slide: Any, slide_id: str, shape_id: str, source: Path, box: tuple[float, float, float, float]
) -> Any:
    source = ensure_regular_input(source, operation="presentation.build")
    try:
        with Image.open(source) as image:
            if image.width * image.height > MAX_IMAGE_PIXELS:
                raise AerError(
                    "LIMIT_EXCEEDED",
                    "Image pixel count exceeds the safety limit.",
                    "presentation.build",
                    str(source),
                )
            image.load()
            ratio = image.width / image.height
    except AerError:
        raise
    except Exception as exc:
        raise AerError(
            "CORRUPT_FILE",
            f"Cannot open image: {exc}",
            "presentation.build",
            str(source),
        ) from exc
    x, y, width, height = box
    box_ratio = width / height
    if ratio >= box_ratio:
        fitted_w = width
        fitted_h = width / ratio
        fitted_x = x
        fitted_y = y + (height - fitted_h) / 2
    else:
        fitted_h = height
        fitted_w = height * ratio
        fitted_x = x + (width - fitted_w) / 2
        fitted_y = y
    shape = slide.shapes.add_picture(
        str(source), Inches(fitted_x), Inches(fitted_y), Inches(fitted_w), Inches(fitted_h)
    )
    shape.name = f"aer:{slide_id}/{shape_id}"
    return shape


def _add_table(slide: Any, slide_id: str, block: dict[str, Any]) -> Any:
    headers = [str(value) for value in block.get("headers", [])]
    rows = block.get("rows", [])
    column_count = max([len(headers), *(len(row) for row in rows)], default=0)
    if not headers and rows:
        headers = [str(index + 1) for index in range(column_count)]
    if not headers:
        raise AerError(
            "INVALID_SPEC",
            "Table requires headers or rows.",
            "presentation.build",
            f"slide:id={slide_id}",
        )
    headers.extend("" for _ in range(column_count - len(headers)))
    normalized_rows = [[str(value) for value in row] for row in rows]
    shape = slide.shapes.add_table(
        len(normalized_rows) + 1,
        len(headers),
        Inches(0.8),
        Inches(1.4),
        Inches(11.7),
        Inches(5.3),
    )
    shape.name = f"aer:{slide_id}/table"
    table = shape.table
    for column, value in enumerate(headers):
        table.cell(0, column).text = value
    for row_index, row in enumerate(normalized_rows, start=1):
        for column in range(len(headers)):
            table.cell(row_index, column).text = row[column] if column < len(row) else ""
    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["navy"] if row_index == 0 else COLORS["light"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Noto Sans CJK KR"
                paragraph.font.size = Pt(14)
                paragraph.font.bold = row_index == 0
                paragraph.font.color.rgb = COLORS["white"] if row_index == 0 else COLORS["text"]
    return shape


def _chart_type(value: str) -> XL_CHART_TYPE:
    mapping = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "horizontal-bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
    }
    try:
        return mapping[value]
    except KeyError as exc:
        raise AerError(
            "INVALID_SPEC",
            f"Unsupported presentation chart type: {value}",
            "presentation.build",
            "/chart_type",
            {"supported": sorted(mapping)},
        ) from exc


def _add_chart(slide: Any, slide_id: str, block: dict[str, Any]) -> Any:
    raw_categories = block.get("categories", [])
    series = block.get("series", [])
    if not raw_categories or not series:
        raise AerError(
            "INVALID_SPEC",
            "Chart requires categories and series.",
            "presentation.build",
            f"slide:id={slide_id}",
        )
    chart_type = str(block.get("chart_type", "bar"))
    if chart_type == "scatter":
        try:
            categories = [float(value) for value in raw_categories]
        except (TypeError, ValueError) as exc:
            raise AerError(
                "INVALID_SPEC",
                "Scatter chart categories must be numeric x values.",
                "presentation.build",
                f"slide:id={slide_id}/categories",
            ) from exc
        xy_data = XyChartData()
        for series_index, item in enumerate(series):
            values = item.get("values", [])
            if len(values) != len(categories):
                raise AerError(
                    "INVALID_SPEC",
                    "Scatter chart series must have one y value for each x value.",
                    "presentation.build",
                    f"slide:id={slide_id}/series/{series_index}/values",
                    {"x_values": len(categories), "y_values": len(values)},
                )
            xy_series = xy_data.add_series(str(item.get("name", "Series")))
            try:
                for x_value, y_value in zip(categories, values, strict=True):
                    xy_series.add_data_point(x_value, float(y_value))
            except (TypeError, ValueError) as exc:
                raise AerError(
                    "INVALID_SPEC",
                    "Scatter chart series values must be numeric.",
                    "presentation.build",
                    f"slide:id={slide_id}/series/{series_index}/values",
                ) from exc
        data: ChartData | XyChartData = xy_data
    else:
        category_data = ChartData()
        category_data.categories = [str(value) for value in raw_categories]
        for series_index, item in enumerate(series):
            values = item.get("values", [])
            if len(values) != len(raw_categories):
                raise AerError(
                    "INVALID_SPEC",
                    "Chart series must have one value for each category.",
                    "presentation.build",
                    f"slide:id={slide_id}/series/{series_index}/values",
                    {"categories": len(raw_categories), "values": len(values)},
                )
            try:
                category_data.add_series(
                    str(item.get("name", "Series")), [float(value) for value in values]
                )
            except (TypeError, ValueError) as exc:
                raise AerError(
                    "INVALID_SPEC",
                    "Chart series values must be numeric.",
                    "presentation.build",
                    f"slide:id={slide_id}/series/{series_index}/values",
                ) from exc
        data = category_data
    chart_shape = slide.shapes.add_chart(
        _chart_type(chart_type),
        Inches(0.9),
        Inches(1.4),
        Inches(11.5),
        Inches(5.2),
        data,
    )
    chart_shape.name = f"aer:{slide_id}/chart"
    chart = chart_shape.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.has_title = False
    return chart_shape


def _validate_slide_limits(block: dict[str, Any], *, index: int) -> None:
    layout = str(block.get("layout", "bullets"))
    collections: list[tuple[str, Any]] = []
    if layout in {"bullets", "timeline"}:
        collections.append(("items", block.get("items", [])))
    elif layout in {"two-column", "comparison"}:
        for side in ("left", "right"):
            column = block.get(side, {})
            collections.append(
                (f"{side}/items", column.get("items", []) if isinstance(column, dict) else column)
            )
    elif layout == "metrics":
        collections.append(("metrics", block.get("metrics", [])))
    elif layout == "table":
        headers = block.get("headers", [])
        rows = block.get("rows", [])
        if not isinstance(headers, list) or not isinstance(rows, list):
            raise AerError(
                "INVALID_SPEC",
                "Presentation table headers and rows must be arrays.",
                "presentation.build",
                f"/content/{index}",
            )
        width = len(headers)
        for row_index, row in enumerate(rows):
            if not isinstance(row, list):
                raise AerError(
                    "INVALID_SPEC",
                    "Presentation table rows must be arrays.",
                    "presentation.build",
                    f"/content/{index}/rows/{row_index}",
                )
            width = max(width, len(row))
        cells = (len(rows) + 1) * width
        if cells > MAX_ARTIFACT_ELEMENTS:
            raise AerError(
                "LIMIT_EXCEEDED",
                "Presentation table exceeds the cell-count limit.",
                "presentation.build",
                f"/content/{index}",
                {"cells": cells, "limit": MAX_ARTIFACT_ELEMENTS},
            )
    elif layout == "chart":
        categories = block.get("categories", [])
        series = block.get("series", [])
        if not isinstance(categories, list) or not isinstance(series, list):
            raise AerError(
                "INVALID_SPEC",
                "Presentation chart categories and series must be arrays.",
                "presentation.build",
                f"/content/{index}",
            )
        values = 0
        for series_index, item in enumerate(series):
            if not isinstance(item, dict) or not isinstance(item.get("values", []), list):
                raise AerError(
                    "INVALID_SPEC",
                    "Presentation chart series and values must be objects and arrays.",
                    "presentation.build",
                    f"/content/{index}/series/{series_index}",
                )
            values += len(item.get("values", []))
        if len(categories) + len(series) + values > MAX_ARTIFACT_ELEMENTS:
            raise AerError(
                "LIMIT_EXCEEDED",
                "Presentation chart exceeds the data-point limit.",
                "presentation.build",
                f"/content/{index}",
                {"items": len(categories) + len(series) + values, "limit": MAX_ARTIFACT_ELEMENTS},
            )
    for name, value in collections:
        if not isinstance(value, list):
            raise AerError(
                "INVALID_SPEC",
                "Presentation item collection must be an array.",
                "presentation.build",
                f"/content/{index}/{name}",
            )
        if len(value) > MAX_ARTIFACT_ELEMENTS:
            raise AerError(
                "LIMIT_EXCEEDED",
                "Presentation item collection exceeds the count limit.",
                "presentation.build",
                f"/content/{index}/{name}",
                {"items": len(value), "limit": MAX_ARTIFACT_ELEMENTS},
            )


def _render_slide(
    slide: Any, slide_id: str, block: dict[str, Any], spec_dir: Path
) -> list[dict[str, Any]]:
    layout = str(block.get("layout", "bullets"))
    if layout not in SUPPORTED_LAYOUTS:
        raise AerError(
            "INVALID_SPEC",
            f"Unsupported presentation layout: {layout}",
            "presentation.build",
            f"slide:id={slide_id}/layout",
            {"supported": sorted(SUPPORTED_LAYOUTS)},
        )
    title = str(block.get("title", ""))
    density = len(title) + sum(len(str(value)) for value in block.values())
    if density > 2400:
        raise AerError(
            "TEXT_OVERFLOW",
            "Slide text exceeds the estimated capacity.",
            "presentation.build",
            f"slide:id={slide_id}",
            suggested_action="Split the content across slides.",
        )
    elements: list[dict[str, Any]] = []
    if layout in {"title", "section", "closing"}:
        _text_box(
            slide,
            slide_id,
            "title",
            title,
            (1.0, 2.0, 11.3, 1.3),
            size=38,
            bold=True,
            color=COLORS["white"],
            align=PP_ALIGN.CENTER,
        )
        subtitle = str(block.get("subtitle", block.get("message", "")))
        if subtitle:
            _text_box(
                slide,
                slide_id,
                "subtitle",
                subtitle,
                (1.4, 3.35, 10.5, 0.8),
                size=21,
                color=COLORS["white"],
                align=PP_ALIGN.CENTER,
            )
            elements.append(
                {
                    "id": "subtitle",
                    "type": "text",
                    "selector": f"slide:id={slide_id}/shape:id=subtitle",
                }
            )
        elements.append(
            {"id": "title", "type": "text", "selector": f"slide:id={slide_id}/shape:id=title"}
        )
    else:
        if title:
            _add_title(slide, slide_id, title)
            elements.append(
                {"id": "title", "type": "text", "selector": f"slide:id={slide_id}/shape:id=title"}
            )
        if layout == "bullets":
            _bullet_box(
                slide,
                slide_id,
                "body",
                list(block.get("items", [])),
                _position(block, (1.0, 1.5, 11.2, 5.2)),
            )
            elements.append(
                {"id": "body", "type": "text", "selector": f"slide:id={slide_id}/shape:id=body"}
            )
        elif layout in {"two-column", "comparison"}:
            for key, x in (("left", 0.8), ("right", 6.75)):
                column = block.get(key, {})
                if not isinstance(column, dict):
                    column = {"items": column}
                _text_box(
                    slide,
                    slide_id,
                    f"{key}-title",
                    str(column.get("title", key.title())),
                    (x, 1.45, 5.5, 0.65),
                    size=21,
                    bold=True,
                    color=COLORS["white"],
                    fill=COLORS["blue"] if key == "left" else COLORS["teal"],
                )
                _bullet_box(
                    slide,
                    slide_id,
                    f"{key}-body",
                    list(column.get("items", [])),
                    (x, 2.2, 5.5, 4.35),
                )
                elements.extend(
                    [
                        {
                            "id": f"{key}-title",
                            "type": "text",
                            "selector": f"slide:id={slide_id}/shape:id={key}-title",
                        },
                        {
                            "id": f"{key}-body",
                            "type": "text",
                            "selector": f"slide:id={slide_id}/shape:id={key}-body",
                        },
                    ]
                )
        elif layout == "metrics":
            metrics = list(block.get("metrics", []))
            if not metrics:
                raise AerError(
                    "INVALID_SPEC",
                    "Metrics layout requires metrics.",
                    "presentation.build",
                    f"slide:id={slide_id}/metrics",
                )
            width = min(3.7, 11.4 / len(metrics))
            total = width * len(metrics)
            start = (SLIDE_W - total) / 2
            for index, metric in enumerate(metrics):
                metric_id = str(metric.get("id", f"metric-{index + 1}"))
                x = start + index * width
                _text_box(
                    slide,
                    slide_id,
                    f"{metric_id}-value",
                    str(metric.get("value", "")),
                    (x, 2.1, width - 0.2, 1.35),
                    size=35,
                    bold=True,
                    color=COLORS["blue"],
                    align=PP_ALIGN.CENTER,
                    fill=COLORS["light"],
                )
                _text_box(
                    slide,
                    slide_id,
                    f"{metric_id}-label",
                    str(metric.get("label", "")),
                    (x, 3.5, width - 0.2, 1.0),
                    size=16,
                    align=PP_ALIGN.CENTER,
                )
                elements.extend(
                    [
                        {
                            "id": f"{metric_id}-value",
                            "type": "text",
                            "selector": f"slide:id={slide_id}/shape:id={metric_id}-value",
                        },
                        {
                            "id": f"{metric_id}-label",
                            "type": "text",
                            "selector": f"slide:id={slide_id}/shape:id={metric_id}-label",
                        },
                    ]
                )
        elif layout == "table":
            _add_table(slide, slide_id, block)
            elements.append(
                {"id": "table", "type": "table", "selector": f"slide:id={slide_id}/shape:id=table"}
            )
        elif layout in {"image", "image-with-caption"}:
            raw_source = block.get("source") or block.get("image")
            if not raw_source:
                raise AerError(
                    "INVALID_SPEC",
                    "Image layout requires source.",
                    "presentation.build",
                    f"slide:id={slide_id}/source",
                )
            source = (
                spec_dir / str(raw_source)
                if not Path(str(raw_source)).is_absolute()
                else Path(str(raw_source))
            )
            _add_picture(
                slide,
                slide_id,
                "image",
                source,
                (1.0, 1.35, 11.3, 4.9 if layout == "image-with-caption" else 5.7),
            )
            elements.append(
                {"id": "image", "type": "image", "selector": f"slide:id={slide_id}/shape:id=image"}
            )
            if layout == "image-with-caption":
                _text_box(
                    slide,
                    slide_id,
                    "caption",
                    str(block.get("caption", "")),
                    (1.0, 6.35, 11.3, 0.45),
                    size=13,
                    color=COLORS["muted"],
                    align=PP_ALIGN.CENTER,
                )
                elements.append(
                    {
                        "id": "caption",
                        "type": "text",
                        "selector": f"slide:id={slide_id}/shape:id=caption",
                    }
                )
        elif layout == "chart":
            _add_chart(slide, slide_id, block)
            elements.append(
                {"id": "chart", "type": "chart", "selector": f"slide:id={slide_id}/shape:id=chart"}
            )
        elif layout == "quote":
            _text_box(
                slide,
                slide_id,
                "quote",
                f"“{block.get('quote', block.get('text', ''))}”",
                (1.3, 1.8, 10.7, 2.9),
                size=30,
                color=COLORS["navy"],
                align=PP_ALIGN.CENTER,
                fill=COLORS["light"],
            )
            _text_box(
                slide,
                slide_id,
                "attribution",
                str(block.get("attribution", "")),
                (4.0, 5.0, 5.3, 0.55),
                size=15,
                color=COLORS["muted"],
                align=PP_ALIGN.CENTER,
            )
            elements.extend(
                [
                    {
                        "id": "quote",
                        "type": "text",
                        "selector": f"slide:id={slide_id}/shape:id=quote",
                    },
                    {
                        "id": "attribution",
                        "type": "text",
                        "selector": f"slide:id={slide_id}/shape:id=attribution",
                    },
                ]
            )
        elif layout == "timeline":
            items = list(block.get("items", []))
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.4), Inches(11.2), Inches(0.06)
            )
            line.name = f"aer:{slide_id}/timeline-line"
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS["blue"]
            line.line.fill.background()
            spacing = 10.5 / max(1, len(items) - 1)
            for index, item in enumerate(items):
                item_id = (
                    str(item.get("id", f"event-{index + 1}"))
                    if isinstance(item, dict)
                    else f"event-{index + 1}"
                )
                label = (
                    str(item.get("label", item.get("title", "")))
                    if isinstance(item, dict)
                    else str(item)
                )
                x = 1.0 + index * spacing
                marker = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(x), Inches(3.15), Inches(0.5), Inches(0.5)
                )
                marker.name = f"aer:{slide_id}/{item_id}-marker"
                marker.fill.solid()
                marker.fill.fore_color.rgb = COLORS["teal"]
                marker.line.fill.background()
                _text_box(
                    slide,
                    slide_id,
                    item_id,
                    label,
                    (x - 0.65, 3.8 if index % 2 == 0 else 2.0, 1.8, 0.9),
                    size=14,
                    align=PP_ALIGN.CENTER,
                )
                elements.append(
                    {
                        "id": item_id,
                        "type": "text",
                        "selector": f"slide:id={slide_id}/shape:id={item_id}",
                    }
                )
    return elements


def build_presentation(
    spec: dict[str, Any], output: Path, *, spec_dir: Path
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    theme = str(spec.get("theme", "business-clean"))
    if theme != "business-clean":
        raise AerError(
            "INVALID_SPEC",
            "Unknown presentation theme.",
            "presentation.build",
            "/theme",
            {"supported": ["business-clean"]},
        )
    metadata = spec.get("metadata", {})
    if not isinstance(metadata, dict):
        raise AerError(
            "INVALID_SPEC",
            "Presentation metadata must be an object.",
            "presentation.build",
            "/metadata",
        )
    ratio = metadata.get("ratio", "16:9")
    if ratio != "16:9":
        raise AerError(
            "INVALID_SPEC",
            "Unsupported presentation aspect ratio.",
            "presentation.build",
            "/metadata/ratio",
            {"supported": ["16:9"]},
        )
    content = spec.get("content")
    if not isinstance(content, list) or not content:
        raise AerError(
            "INVALID_SPEC",
            "Presentation content must be a non-empty array.",
            "presentation.build",
            "/content",
        )
    if len(content) > MAX_PRESENTATION_SLIDES:
        raise AerError(
            "LIMIT_EXCEEDED",
            "Presentation exceeds the slide-count limit.",
            "presentation.build",
            "/content",
            {"slides": len(content), "limit": MAX_PRESENTATION_SLIDES},
        )
    for index, block in enumerate(content):
        if isinstance(block, dict):
            _validate_slide_limits(block, index=index)
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    presentation.core_properties.title = str(metadata.get("title", ""))
    presentation.core_properties.subject = "Built by Agent Efficiency Runtime"
    fixed = datetime(2000, 1, 1, tzinfo=UTC)
    presentation.core_properties.created = fixed
    presentation.core_properties.modified = fixed
    elements: list[dict[str, Any]] = []
    seen_ids: set[str] = set()
    footer = str(spec.get("footer", ""))
    for index, block in enumerate(content, start=1):
        if not isinstance(block, dict):
            raise AerError(
                "INVALID_SPEC",
                "Each slide must be an object.",
                "presentation.build",
                f"/content/{index - 1}",
            )
        slide_id = str(block.get("id", f"slide-{index}"))
        if not slide_id or slide_id in seen_ids or "/" in slide_id:
            raise AerError(
                "INVALID_SPEC",
                "Slide IDs must be unique and cannot contain '/'.",
                "presentation.build",
                f"/content/{index - 1}/id",
            )
        seen_ids.add(slide_id)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.element.cSld.set("name", f"aer:{slide_id}")
        background = slide.background.fill
        background.solid()
        layout = str(block.get("layout", "bullets"))
        background.fore_color.rgb = (
            COLORS["navy"] if layout in {"title", "section", "closing"} else COLORS["white"]
        )
        slide_elements = _render_slide(slide, slide_id, block, spec_dir)
        if footer and layout not in {"title", "closing"}:
            _text_box(
                slide,
                slide_id,
                "footer",
                footer,
                (0.7, 7.02, 10.6, 0.26),
                size=9,
                color=COLORS["muted"],
            )
            slide_elements.append(
                {"id": "footer", "type": "text", "selector": f"slide:id={slide_id}/shape:id=footer"}
            )
        _text_box(
            slide,
            slide_id,
            "slide-number",
            str(index),
            (11.8, 7.0, 0.75, 0.26),
            size=9,
            color=COLORS["muted"],
            align=PP_ALIGN.RIGHT,
        )
        slide_elements.append(
            {
                "id": "slide-number",
                "type": "text",
                "selector": f"slide:id={slide_id}/shape:id=slide-number",
            }
        )
        seen_selectors: set[str] = set()
        for element in slide_elements:
            selector = str(element["selector"])
            if selector in seen_selectors:
                raise AerError(
                    "INVALID_SPEC",
                    "Shape IDs must be unique within each slide.",
                    "presentation.build",
                    f"slide:id={slide_id}",
                    {"duplicate_selector": selector},
                )
            seen_selectors.add(selector)
        elements.append(
            {
                "id": slide_id,
                "type": "slide",
                "index": index,
                "selector": f"slide:id={slide_id}",
                "children": slide_elements,
            }
        )
    buffer = io.BytesIO()
    presentation.save(buffer)
    atomic_write_bytes(output, buffer.getvalue())
    return elements, []
