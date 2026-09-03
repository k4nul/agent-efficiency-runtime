"""Executable, persisted token-efficiency benchmark scenarios."""

from __future__ import annotations

import csv
import io
import json
import shutil
import sqlite3
import subprocess
import sys
import tempfile
import textwrap
import time
import uuid
from collections.abc import Callable
from datetime import UTC, datetime
from pathlib import Path
from typing import Final

import yaml
from pptx import Presentation

from aer.archive import verify_archive
from aer.artifacts import build_artifact
from aer.benchmark.models import BenchmarkRun, ScenarioResult, VariantMeasurement
from aer.config import Settings
from aer.data import query_data
from aer.errors import AerError
from aer.hashing import sha256_bytes, sha256_file
from aer.patch import apply_patch
from aer.paths import atomic_write_bytes, atomic_write_text
from aer.recipes import RecipeEngine
from aer.runner import run_command
from aer.store import ObjectStore

SCENARIOS: Final[tuple[str, ...]] = (
    "log-compaction",
    "data-query",
    "presentation-spec",
    "json-patch",
    "presentation-patch",
    "recipe-package",
)

DESCRIPTIONS: Final[dict[str, str]] = {
    "log-compaction": "Full 5,000-line command log versus compact failure context.",
    "data-query": "Full 10,000-row CSV transfer versus local filtered preview.",
    "presentation-spec": "Direct python-pptx boilerplate versus semantic presentation spec.",
    "json-patch": "Whole JSON rewrite versus an RFC 6901 pointer patch.",
    "presentation-patch": "Whole presentation rebuild versus one stable-ID element patch.",
    "recipe-package": "Repeated packaging script versus a validated capability recipe.",
}


def _timestamp() -> str:
    return datetime.now(UTC).isoformat(timespec="microseconds").replace("+00:00", "Z")


def _elapsed_ms(started_ns: int) -> float:
    return (time.perf_counter_ns() - started_ns) / 1_000_000


def _compact_json(value: object) -> bytes:
    return json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        default=str,
    ).encode("utf-8")


def _yaml_bytes(value: object) -> bytes:
    return yaml.safe_dump(value, allow_unicode=True, sort_keys=False).encode("utf-8")


def _argv_bytes(argv: list[str]) -> int:
    return len(_compact_json(argv))


def _pptx_text(path: Path) -> list[str]:
    presentation = Presentation(str(path))
    return [
        str(shape.text)
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]


class BenchmarkEngine:
    """Execute benchmark workloads and persist measured results to SQLite."""

    def __init__(self, settings: Settings | None = None) -> None:
        self.settings = settings or Settings.load()
        self.settings.ensure()
        self.store = ObjectStore(self.settings)
        self._initialize_database()

    @staticmethod
    def scenario_names() -> tuple[str, ...]:
        return SCENARIOS

    def run(self, *, scenario: str | None = None) -> BenchmarkRun:
        if scenario is not None and scenario not in SCENARIOS:
            raise AerError(
                "INVALID_ARGUMENT",
                "Unknown benchmark scenario.",
                operation="benchmark.run",
                target=scenario,
                details={"available": list(SCENARIOS)},
            )
        selected = SCENARIOS if scenario is None else (scenario,)
        methods: dict[str, Callable[[Path], ScenarioResult]] = {
            "log-compaction": self._log_compaction,
            "data-query": self._data_query,
            "presentation-spec": self._presentation_spec,
            "json-patch": self._json_patch,
            "presentation-patch": self._presentation_patch,
            "recipe-package": self._recipe_package,
        }
        run_started = time.perf_counter_ns()
        results: list[ScenarioResult] = []
        with tempfile.TemporaryDirectory(prefix="aer-benchmark-") as temporary:
            root = Path(temporary)
            for name in selected:
                workdir = root / name
                workdir.mkdir(mode=0o700)
                try:
                    result = methods[name](workdir)
                except AerError as exc:
                    result = ScenarioResult.failed(
                        name,
                        DESCRIPTIONS[name],
                        code=exc.code,
                        message=exc.message,
                    )
                except Exception as exc:  # benchmark must persist compact failure evidence
                    result = ScenarioResult.failed(
                        name,
                        DESCRIPTIONS[name],
                        code="INTERNAL_ERROR",
                        message=f"{type(exc).__name__}: {exc}",
                    )
                results.append(result)
        run = BenchmarkRun(
            run_id=uuid.uuid4().hex,
            timestamp=_timestamp(),
            duration_ms=_elapsed_ms(run_started),
            success=all(result.success for result in results),
            scenarios=tuple(results),
        )
        self._persist(run)
        return run

    def report(self, *, limit: int = 10) -> list[BenchmarkRun]:
        if limit < 1 or limit > 1000:
            raise AerError(
                "INVALID_ARGUMENT",
                "Benchmark report limit must be between 1 and 1000.",
                operation="benchmark.report",
                target=str(limit),
            )
        with self.settings.connect() as connection:
            rows = connection.execute(
                """
                SELECT run_id, timestamp, duration_ms, success, result_json
                FROM aer_benchmark_runs
                ORDER BY timestamp DESC, rowid DESC
                LIMIT ?
                """,
                (limit,),
            ).fetchall()
        return [self._row_to_run(row) for row in rows]

    def _log_compaction(self, workdir: Path) -> ScenarioResult:
        del workdir
        script = textwrap.dedent(
            """
            import sys
            for index in range(5000):
                print(f"test_{index}: passed")
            print("ERROR benchmark sentinel failure", file=sys.stderr)
            raise SystemExit(1)
            """
        ).strip()
        argv = [sys.executable, "-c", script]

        direct_started = time.perf_counter_ns()
        completed = subprocess.run(
            argv,
            capture_output=True,
            timeout=30,
            check=False,
            shell=False,
        )
        direct_ms = _elapsed_ms(direct_started)
        direct_payload = completed.stdout + completed.stderr
        direct_valid = completed.returncode == 1 and b"benchmark sentinel failure" in direct_payload
        direct = VariantMeasurement.measured(
            "full-log",
            input_bytes=_argv_bytes(argv),
            output_bytes=len(direct_payload),
            context_bytes=len(direct_payload),
            wall_time_ms=direct_ms,
            valid=direct_valid,
            sha256=sha256_bytes(direct_payload),
            details={"exit_code": completed.returncode, "lines": len(direct_payload.splitlines())},
        )

        aer_started = time.perf_counter_ns()
        result = run_command(argv, timeout=30, store=self.store)
        aer_ms = _elapsed_ms(aer_started)
        compact = _compact_json(result.to_dict())
        aer_valid = (
            result.exit_code == 1
            and result.raw_ref is not None
            and any("benchmark sentinel failure" in line for line in result.failure_context)
            and len(compact) < len(direct_payload)
        )
        aer = VariantMeasurement.measured(
            "aer-run",
            input_bytes=_argv_bytes(argv),
            output_bytes=len(compact),
            context_bytes=len(compact),
            wall_time_ms=aer_ms,
            valid=aer_valid,
            sha256=sha256_bytes(compact),
            details={
                "exit_code": result.exit_code,
                "raw_ref": result.raw_ref,
                "failure_context_lines": len(result.failure_context),
                "captured_bytes": result.bytes_captured,
            },
        )
        return ScenarioResult.compared(
            "log-compaction", DESCRIPTIONS["log-compaction"], direct, aer
        )

    def _data_query(self, workdir: Path) -> ScenarioResult:
        source = workdir / "orders.csv"
        stream = io.StringIO(newline="")
        writer = csv.writer(stream, lineterminator="\n")
        writer.writerow(["id", "status", "total"])
        for index in range(10_000):
            writer.writerow([f"O-{index:05d}", "pending" if index % 2 == 0 else "paid", index * 17])
        atomic_write_text(source, stream.getvalue())
        source_size = source.stat().st_size

        direct_started = time.perf_counter_ns()
        direct_payload = source.read_bytes()
        direct_ms = _elapsed_ms(direct_started)
        direct = VariantMeasurement.measured(
            "full-csv",
            input_bytes=source_size,
            output_bytes=len(direct_payload),
            context_bytes=len(direct_payload),
            wall_time_ms=direct_ms,
            valid=direct_payload.count(b"\n") == 10_001,
            sha256=sha256_bytes(direct_payload),
            details={"source_rows": 10_000},
        )

        aer_started = time.perf_counter_ns()
        result = query_data(
            source,
            where="status == pending",
            select="id,total",
            store=self.store,
        )
        aer_ms = _elapsed_ms(aer_started)
        compact = _compact_json(result.to_dict())
        aer_valid = (
            result.source_rows == 10_000
            and result.matched_rows == 5_000
            and len(result.preview) <= 20
            and result.result_ref is not None
            and len(compact) < source_size
        )
        aer = VariantMeasurement.measured(
            "aer-data-query",
            input_bytes=source_size,
            output_bytes=len(compact),
            context_bytes=len(compact),
            wall_time_ms=aer_ms,
            valid=aer_valid,
            sha256=sha256_bytes(compact),
            details={
                "source_rows": result.source_rows,
                "matched_rows": result.matched_rows,
                "preview_rows": len(result.preview),
                "result_ref": result.result_ref,
            },
        )
        return ScenarioResult.compared("data-query", DESCRIPTIONS["data-query"], direct, aer)

    def _presentation_spec(self, workdir: Path) -> ScenarioResult:
        direct_output = workdir / "direct.pptx"
        direct_code = textwrap.dedent(
            """
            import sys
            from pptx import Presentation
            from pptx.util import Inches, Pt

            presentation = Presentation()
            slides = [
                ("Agent Efficiency Runtime", ["Deterministic local runtime"]),
                ("Problem", ["Repeated Python boilerplate", "Whole-file rewrites"]),
                ("Result", ["Compact selectors", "Content-addressed artifacts"]),
            ]
            for title, bullets in slides:
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                title_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(0.5), Inches(11), Inches(0.8)
                )
                title_box.text = title
                title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
                body = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(10.5), Inches(4.5))
                body.text = "\\n".join(f"• {item}" for item in bullets)
                for paragraph in body.text_frame.paragraphs:
                    paragraph.runs[0].font.size = Pt(20)
            presentation.save(sys.argv[1])
            """
        ).strip()
        direct_started = time.perf_counter_ns()
        completed = subprocess.run(
            [sys.executable, "-c", direct_code, str(direct_output)],
            capture_output=True,
            timeout=30,
            check=False,
            shell=False,
        )
        direct_ms = _elapsed_ms(direct_started)
        direct_valid = (
            completed.returncode == 0 and len(Presentation(str(direct_output)).slides) == 3
        )
        direct = VariantMeasurement.measured(
            "direct-python-pptx",
            input_bytes=len(direct_code.encode("utf-8")),
            output_bytes=direct_output.stat().st_size,
            context_bytes=len(direct_code.encode("utf-8")),
            wall_time_ms=direct_ms,
            valid=direct_valid,
            sha256=sha256_file(direct_output),
            details={"slides": 3, "return_code": completed.returncode},
        )

        spec = {
            "version": 1,
            "kind": "presentation",
            "metadata": {"title": "Agent Efficiency Runtime", "locale": "en-US"},
            "theme": "business-clean",
            "content": [
                {
                    "id": "cover",
                    "layout": "title",
                    "title": "Agent Efficiency Runtime",
                    "subtitle": "Deterministic local runtime",
                },
                {
                    "id": "problem",
                    "layout": "bullets",
                    "title": "Problem",
                    "items": ["Repeated Python boilerplate", "Whole-file rewrites"],
                },
                {
                    "id": "result",
                    "layout": "bullets",
                    "title": "Result",
                    "items": ["Compact selectors", "Content-addressed artifacts"],
                },
            ],
        }
        spec_bytes = _yaml_bytes(spec)
        spec_path = workdir / "presentation.yaml"
        output = workdir / "semantic.pptx"
        atomic_write_bytes(spec_path, spec_bytes)
        aer_started = time.perf_counter_ns()
        build_result = build_artifact(spec_path, output)
        aer_ms = _elapsed_ms(aer_started)
        aer_valid = (
            len(Presentation(str(output)).slides) == 3 and int(build_result["element_count"]) > 0
        )
        compact_result = _compact_json(build_result)
        aer = VariantMeasurement.measured(
            "semantic-spec",
            input_bytes=len(spec_bytes),
            output_bytes=output.stat().st_size,
            context_bytes=len(spec_bytes),
            wall_time_ms=aer_ms,
            valid=aer_valid,
            sha256=sha256_file(output),
            details={"slides": 3, "compact_result_bytes": len(compact_result)},
        )
        return ScenarioResult.compared(
            "presentation-spec", DESCRIPTIONS["presentation-spec"], direct, aer
        )

    def _json_patch(self, workdir: Path) -> ScenarioResult:
        payload = {
            "metadata": {"title": "Benchmark", "owner": "aer"},
            "items": [
                {
                    "id": f"item-{index:04d}",
                    "status": "pending",
                    "values": list(range(index % 20)),
                    "preserve": f"stable-{index:04d}",
                }
                for index in range(1_000)
            ],
        }
        original = (json.dumps(payload, ensure_ascii=False, indent=2) + "\n").encode("utf-8")
        direct_path = workdir / "direct.json"
        patch_path = workdir / "patched.json"
        atomic_write_bytes(direct_path, original)
        atomic_write_bytes(patch_path, original)

        direct_started = time.perf_counter_ns()
        direct_value = json.loads(direct_path.read_text(encoding="utf-8"))
        direct_value["items"][500]["status"] = "complete"
        rewritten = (json.dumps(direct_value, ensure_ascii=False, indent=2) + "\n").encode("utf-8")
        atomic_write_bytes(direct_path, rewritten)
        direct_ms = _elapsed_ms(direct_started)
        direct_valid = (
            direct_value["items"][500]["status"] == "complete"
            and direct_value["items"][501]["preserve"] == "stable-0501"
        )
        direct = VariantMeasurement.measured(
            "whole-json-rewrite",
            input_bytes=len(original),
            output_bytes=len(rewritten),
            context_bytes=len(rewritten),
            wall_time_ms=direct_ms,
            valid=direct_valid,
            sha256=sha256_file(direct_path),
            details={"items": 1_000},
        )

        patch_spec = {
            "version": 1,
            "operations": [{"op": "json.set", "target": "/items/500/status", "value": "complete"}],
        }
        patch_bytes = _yaml_bytes(patch_spec)
        patch_spec_path = workdir / "patch.yaml"
        atomic_write_bytes(patch_spec_path, patch_bytes)
        aer_started = time.perf_counter_ns()
        patch_result = apply_patch(patch_path, patch_spec_path, settings=self.settings)
        aer_ms = _elapsed_ms(aer_started)
        patched_value = json.loads(patch_path.read_text(encoding="utf-8"))
        aer_valid = (
            patched_value == direct_value
            and patched_value["items"][501]["preserve"] == "stable-0501"
            and patch_result["before_sha256"] != patch_result["after_sha256"]
        )
        aer = VariantMeasurement.measured(
            "json-pointer-patch",
            input_bytes=len(original) + len(patch_bytes),
            output_bytes=patch_path.stat().st_size,
            context_bytes=len(patch_bytes),
            wall_time_ms=aer_ms,
            valid=aer_valid,
            sha256=sha256_file(patch_path),
            details={"operations": 1, "preserved_equal": patched_value == direct_value},
        )
        return ScenarioResult.compared("json-patch", DESCRIPTIONS["json-patch"], direct, aer)

    def _presentation_patch(self, workdir: Path) -> ScenarioResult:
        base_spec = {
            "version": 1,
            "kind": "presentation",
            "metadata": {"title": "Patch benchmark"},
            "theme": "business-clean",
            "content": [
                {
                    "id": "cover",
                    "layout": "title",
                    "title": "Patch benchmark",
                    "subtitle": "Preserve unrelated content",
                },
                {
                    "id": "metrics",
                    "layout": "metrics",
                    "title": "Measured reduction",
                    "metrics": [
                        {"id": "token", "value": "50%", "label": "Token reduction"},
                        {"id": "calls", "value": "8", "label": "Model calls"},
                    ],
                },
                {
                    "id": "closing",
                    "layout": "closing",
                    "title": "Keep the stable ending",
                },
            ],
        }
        base_path = workdir / "base.yaml"
        base_output = workdir / "base.pptx"
        atomic_write_bytes(base_path, _yaml_bytes(base_spec))
        build_artifact(base_path, base_output)

        modified_spec = json.loads(json.dumps(base_spec))
        modified_spec["content"][1]["metrics"][0]["value"] = "63%"
        modified_bytes = _yaml_bytes(modified_spec)
        modified_path = workdir / "modified.yaml"
        rebuilt_output = workdir / "rebuilt.pptx"
        atomic_write_bytes(modified_path, modified_bytes)
        direct_started = time.perf_counter_ns()
        build_artifact(modified_path, rebuilt_output)
        direct_ms = _elapsed_ms(direct_started)
        direct_text = _pptx_text(rebuilt_output)
        direct_valid = "63%" in direct_text and "Keep the stable ending" in direct_text
        direct = VariantMeasurement.measured(
            "full-presentation-rebuild",
            input_bytes=len(modified_bytes),
            output_bytes=rebuilt_output.stat().st_size,
            context_bytes=len(modified_bytes),
            wall_time_ms=direct_ms,
            valid=direct_valid,
            sha256=sha256_file(rebuilt_output),
            details={"slides": len(Presentation(str(rebuilt_output)).slides)},
        )

        patched_output = workdir / "patched.pptx"
        shutil.copyfile(base_output, patched_output)
        patch_spec = {
            "version": 1,
            "operations": [
                {
                    "op": "pptx.set_text",
                    "target": "slide:id=metrics/shape:id=token-value",
                    "value": "63%",
                }
            ],
        }
        patch_bytes = _yaml_bytes(patch_spec)
        patch_path = workdir / "presentation-patch.yaml"
        atomic_write_bytes(patch_path, patch_bytes)
        aer_started = time.perf_counter_ns()
        patch_result = apply_patch(patched_output, patch_path, settings=self.settings)
        aer_ms = _elapsed_ms(aer_started)
        patched_text = _pptx_text(patched_output)
        aer_valid = (
            "63%" in patched_text
            and "Keep the stable ending" in patched_text
            and len(Presentation(str(patched_output)).slides)
            == len(Presentation(str(base_output)).slides)
            and patch_result["before_sha256"] != patch_result["after_sha256"]
        )
        aer = VariantMeasurement.measured(
            "stable-id-element-patch",
            input_bytes=len(patch_bytes),
            output_bytes=patched_output.stat().st_size,
            context_bytes=len(patch_bytes),
            wall_time_ms=aer_ms,
            valid=aer_valid,
            sha256=sha256_file(patched_output),
            details={"operations": 1, "unrelated_text_preserved": True},
        )
        return ScenarioResult.compared(
            "presentation-patch", DESCRIPTIONS["presentation-patch"], direct, aer
        )

    def _recipe_package(self, workdir: Path) -> ScenarioResult:
        source = workdir / "delivery"
        source.mkdir()
        for index in range(12):
            atomic_write_text(
                source / f"file-{index:02d}.txt",
                f"deterministic benchmark file {index}\n" * (index + 1),
            )
        direct_output = workdir / "direct.zip"
        direct_script = textwrap.dedent(
            """
            import hashlib
            import json
            import stat
            import sys
            import zipfile
            from pathlib import Path

            source = Path(sys.argv[1])
            output = Path(sys.argv[2])
            contents = []
            manifest_files = []
            for path in sorted(item for item in source.rglob("*") if item.is_file()):
                name = path.relative_to(source).as_posix()
                data = path.read_bytes()
                contents.append((name, data))
                manifest_files.append({
                    "path": name,
                    "sha256": hashlib.sha256(data).hexdigest(),
                    "size": len(data),
                })
            manifest = (json.dumps({
                "version": 1,
                "deterministic_timestamp": "1980-01-01T00:00:00Z",
                "files": manifest_files,
            }, ensure_ascii=False, indent=2, sort_keys=True) + "\\n").encode("utf-8")

            def info(name):
                value = zipfile.ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
                value.compress_type = zipfile.ZIP_DEFLATED
                value.create_system = 3
                value.external_attr = (stat.S_IFREG | 0o644) << 16
                return value

            with zipfile.ZipFile(
                output,
                "w",
                compression=zipfile.ZIP_DEFLATED,
                compresslevel=9,
                strict_timestamps=True,
            ) as archive:
                for name, data in contents:
                    archive.writestr(
                        info(name), data, compress_type=zipfile.ZIP_DEFLATED, compresslevel=9
                    )
                archive.writestr(
                    info("manifest.json"),
                    manifest,
                    compress_type=zipfile.ZIP_DEFLATED,
                    compresslevel=9,
                )
            """
        ).strip()
        direct_started = time.perf_counter_ns()
        completed = subprocess.run(
            [sys.executable, "-c", direct_script, str(source), str(direct_output)],
            capture_output=True,
            timeout=30,
            check=False,
            shell=False,
        )
        direct_verification = verify_archive(direct_output)
        direct_ms = _elapsed_ms(direct_started)
        direct_valid = completed.returncode == 0 and bool(direct_verification["valid"])
        direct = VariantMeasurement.measured(
            "one-off-package-script",
            input_bytes=len(direct_script.encode("utf-8")),
            output_bytes=direct_output.stat().st_size,
            context_bytes=len(direct_script.encode("utf-8")),
            wall_time_ms=direct_ms,
            valid=direct_valid,
            sha256=sha256_file(direct_output),
            details={
                "entries": direct_verification["entries"],
                "return_code": completed.returncode,
            },
        )

        recipe = {
            "version": 1,
            "name": "benchmark-package",
            "description": "Package and verify a benchmark directory.",
            "inputs": {
                "source": {"type": "path"},
                "output": {"type": "path"},
            },
            "steps": [
                {
                    "id": "package",
                    "uses": "archive.create",
                    "with": {
                        "source": "${{ inputs.source }}",
                        "output": "${{ inputs.output }}",
                    },
                },
                {
                    "id": "verify",
                    "uses": "archive.verify",
                    "with": {"target": "${{ steps.package.output }}"},
                },
            ],
        }
        recipe_bytes = _yaml_bytes(recipe)
        recipe_path = workdir / "recipe.yaml"
        recipe_output = workdir / "recipe.zip"
        atomic_write_bytes(recipe_path, recipe_bytes)
        engine = RecipeEngine(self.settings)
        aer_started = time.perf_counter_ns()
        recipe_result = engine.run(
            recipe_path,
            variables={"source": str(source), "output": str(recipe_output)},
            trust=True,
        )
        aer_ms = _elapsed_ms(aer_started)
        recipe_verification = verify_archive(recipe_output)
        aer_valid = (
            bool(recipe_result["success"])
            and bool(recipe_verification["valid"])
            and sha256_file(recipe_output) == sha256_file(direct_output)
        )
        aer = VariantMeasurement.measured(
            "capability-recipe",
            input_bytes=len(recipe_bytes),
            output_bytes=recipe_output.stat().st_size,
            context_bytes=len(recipe_bytes),
            wall_time_ms=aer_ms,
            valid=aer_valid,
            sha256=sha256_file(recipe_output),
            details={
                "steps": len(recipe_result["steps"]),
                "log_ref": recipe_result["log_ref"],
                "byte_identical_to_direct": sha256_file(recipe_output)
                == sha256_file(direct_output),
            },
        )
        return ScenarioResult.compared(
            "recipe-package", DESCRIPTIONS["recipe-package"], direct, aer
        )

    def _initialize_database(self) -> None:
        with self.settings.connect() as connection:
            connection.execute(
                """
                CREATE TABLE IF NOT EXISTS aer_benchmark_runs (
                    run_id TEXT PRIMARY KEY,
                    timestamp TEXT NOT NULL,
                    duration_ms REAL NOT NULL CHECK(duration_ms >= 0),
                    success INTEGER NOT NULL CHECK(success IN (0, 1)),
                    result_json TEXT NOT NULL
                )
                """
            )
            connection.execute(
                """
                CREATE INDEX IF NOT EXISTS aer_benchmark_timestamp
                ON aer_benchmark_runs(timestamp)
                """
            )

    def _persist(self, run: BenchmarkRun) -> None:
        payload = json.dumps(
            run.as_dict(),
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        )
        with self.settings.connect() as connection:
            connection.execute(
                """
                INSERT INTO aer_benchmark_runs (
                    run_id, timestamp, duration_ms, success, result_json
                ) VALUES (?, ?, ?, ?, ?)
                """,
                (run.run_id, run.timestamp, run.duration_ms, int(run.success), payload),
            )

    @staticmethod
    def _row_to_run(row: sqlite3.Row) -> BenchmarkRun:
        try:
            payload = json.loads(str(row["result_json"]))
            scenario_values = payload["scenarios"]
            if not isinstance(scenario_values, list):
                raise TypeError("scenarios is not a list")
            scenarios = tuple(
                BenchmarkEngine._scenario_from_dict(value) for value in scenario_values
            )
        except (KeyError, TypeError, ValueError, json.JSONDecodeError) as exc:
            raise AerError(
                "CORRUPT_FILE",
                "Stored benchmark result is invalid.",
                operation="benchmark.report",
                target=str(row["run_id"]),
            ) from exc
        return BenchmarkRun(
            run_id=str(row["run_id"]),
            timestamp=str(row["timestamp"]),
            duration_ms=float(row["duration_ms"]),
            success=bool(row["success"]),
            scenarios=scenarios,
        )

    @staticmethod
    def _scenario_from_dict(value: object) -> ScenarioResult:
        if not isinstance(value, dict):
            raise TypeError("scenario is not an object")

        def measurement(raw: object) -> VariantMeasurement | None:
            if raw is None:
                return None
            if not isinstance(raw, dict):
                raise TypeError("measurement is not an object")
            details = raw.get("details", {})
            if not isinstance(details, dict):
                raise TypeError("measurement details is not an object")
            return VariantMeasurement(
                name=str(raw["name"]),
                input_bytes=int(raw["input_bytes"]),
                output_bytes=int(raw["output_bytes"]),
                context_bytes=int(raw["context_bytes"]),
                estimated_tokens=int(raw["estimated_tokens"]),
                estimation_method=str(raw["estimation_method"]),
                not_provider_billed_tokens=bool(raw["not_provider_billed_tokens"]),
                wall_time_ms=float(raw["wall_time_ms"]),
                retries=int(raw.get("retries", 0)),
                valid=bool(raw["valid"]),
                sha256=str(raw["sha256"]),
                details={str(key): item for key, item in details.items()},
            )

        error = value.get("error")
        if error is not None and not isinstance(error, dict):
            raise TypeError("scenario error is not an object")
        context_saved = value.get("context_bytes_saved")
        token_saved = value.get("estimated_tokens_saved")
        return ScenarioResult(
            scenario=str(value["scenario"]),
            description=str(value["description"]),
            direct=measurement(value.get("direct")),
            aer=measurement(value.get("aer")),
            success=bool(value["success"]),
            context_bytes_saved=None if context_saved is None else int(context_saved),
            estimated_tokens_saved=None if token_saved is None else int(token_saved),
            error=None if error is None else {str(key): item for key, item in error.items()},
        )
