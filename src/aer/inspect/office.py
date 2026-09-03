"""Bounded structural inspection for Office Open XML and PDF artifacts."""

from __future__ import annotations

import io
import re
import zipfile
from datetime import date, datetime, time
from pathlib import Path
from typing import Any

from docx import Document
from docx.document import Document as DocumentType
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
from docx.oxml.ns import qn
from openpyxl import load_workbook
from openpyxl.cell.cell import Cell
from openpyxl.utils import get_column_letter, range_boundaries
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
from pypdf import PdfReader
from pypdf.errors import PdfReadError

from aer.artifacts.workbook.selectors import stable_cell_name, stable_sheet_name
from aer.errors import AerError
from aer.inspect.common import RawSink, compact_line, parse_inclusive_range, preserve_overflow
from aer.pdf.safety import (
    bounded_pdf_page_count,
    ensure_bounded_pdf_input,
    extract_pdf_page_text,
    pdf_attachment_names,
    search_pdf_text,
)
from aer.zip_safety import enforce_zip_expansion_limits

_SLIDE_ID_RE = re.compile(r"(?:^|/)slide(?::id)?[=:]([^/]+)", re.IGNORECASE)
_SHAPE_ID_RE = re.compile(r"(?:^|/)shape(?::id)?[=:]([^/]+)", re.IGNORECASE)
_AER_ID_RE = re.compile(r"(?:^|[\[/])aer:id=([^\]/]+)", re.IGNORECASE)
_AER_SHAPE_PATH_RE = re.compile(r"^aer:([^/]+)/(.+)$", re.IGNORECASE)
_HEADING_RE = re.compile(r"^Heading\s+(\d+)$", re.IGNORECASE)
_MAX_RANGE_CELLS = 100_000
_MAX_PREVIEW_COLUMNS = 20


def inspect_xlsx(
    path: Path,
    *,
    sheet: str | None,
    cell_range: str | None,
    rows: str | None,
    selector: str | None,
    formulas: bool,
    max_items: int,
    raw_sink: RawSink | None,
) -> dict[str, Any]:
    _check_ooxml_package(path, required_part="xl/workbook.xml")
    selector_sheet, selector_range, selector_cell_id = _parse_xlsx_selector(selector)
    if selector_sheet is not None:
        if sheet is not None and sheet != selector_sheet:
            raise _invalid_selector(selector or "", "Selector and --sheet disagree.")
        sheet = selector_sheet
    if selector_range is not None:
        if cell_range is not None and cell_range != selector_range:
            raise _invalid_selector(selector or "", "Selector and --range disagree.")
        cell_range = selector_range
    if selector_cell_id is not None and cell_range is not None:
        raise _invalid_selector(selector or "", "Stable cell selector and --range disagree.")
    try:
        workbook = load_workbook(path, read_only=False, data_only=False, keep_links=False)
    except (OSError, ValueError, KeyError, zipfile.BadZipFile) as exc:
        raise _corrupt(path, "XLSX could not be opened.", exc) from exc
    try:
        sheet_summaries: list[dict[str, Any]] = []
        all_formulas: list[dict[str, Any]] = []
        for worksheet in workbook.worksheets:
            worksheet_formulas = [
                {"sheet": worksheet.title, "cell": cell.coordinate, "formula": cell.value}
                for row_cells in worksheet.iter_rows()
                for cell in row_cells
                if cell.data_type == "f"
            ]
            all_formulas.extend(worksheet_formulas)
            sheet_summaries.append(
                {
                    "name": worksheet.title,
                    "range": worksheet.calculate_dimension(),
                    "rows": worksheet.max_row,
                    "columns": worksheet.max_column,
                    "formula_count": len(worksheet_formulas),
                    "merged_cells": sorted(str(value) for value in worksheet.merged_cells.ranges),
                    "chart_count": len(worksheet._charts),
                }
            )
        result: dict[str, Any] = {
            "type": "xlsx",
            "sheet_count": len(workbook.sheetnames),
            "sheets": sheet_summaries,
            "named_ranges": sorted(str(name) for name in workbook.defined_names),
        }
        if formulas:
            result["formula_count"] = len(all_formulas)
            result["formulas"] = all_formulas[:max_items]
            if len(all_formulas) > max_items:
                result["truncated"] = True
                result["raw_ref"] = preserve_overflow(
                    all_formulas, raw_sink=raw_sink, name=f"{path.name}.formulas.json"
                )
        if sheet is not None or cell_range is not None or rows is not None or selector is not None:
            selected_name = sheet or workbook.sheetnames[0]
            if selector_sheet is not None:
                defined_name = stable_sheet_name(selector_sheet)
                if defined_name in workbook.defined_names:
                    destinations = list(workbook.defined_names[defined_name].destinations)
                    if len(destinations) != 1 or destinations[0][0] not in workbook.sheetnames:
                        raise _invalid_selector(
                            selector or selector_sheet,
                            "Stable sheet ID does not resolve to exactly one sheet.",
                        )
                    selected_name = destinations[0][0]
                elif selector_sheet in workbook.sheetnames:
                    selected_name = selector_sheet
            elif selected_name not in workbook.sheetnames:
                defined_name = stable_sheet_name(selected_name)
                if defined_name in workbook.defined_names:
                    destinations = list(workbook.defined_names[defined_name].destinations)
                    if len(destinations) == 1:
                        selected_name = destinations[0][0]
            if selected_name not in workbook.sheetnames:
                normalized = selected_name.casefold()
                selected_name = next(
                    (
                        name
                        for name in workbook.sheetnames
                        if name.casefold().replace(" ", "-") == normalized
                    ),
                    selected_name,
                )
            if selected_name not in workbook.sheetnames:
                raise _invalid_selector(selected_name, "Workbook sheet was not found.")
            if selector_cell_id is not None:
                assert selector_sheet is not None
                defined_name = stable_cell_name(selector_sheet, selector_cell_id)
                if defined_name not in workbook.defined_names:
                    raise _invalid_selector(
                        selector or selector_cell_id, "Stable cell ID was not found."
                    )
                destinations = list(workbook.defined_names[defined_name].destinations)
                if len(destinations) != 1 or destinations[0][0] != selected_name:
                    raise _invalid_selector(
                        selector or selector_cell_id,
                        "Stable cell ID does not resolve to the selected sheet.",
                    )
                cell_range = destinations[0][1].replace("$", "")
            worksheet = workbook[selected_name]
            result["selection"] = _xlsx_selection(
                worksheet,
                cell_range=cell_range,
                rows=rows,
                max_items=max_items,
                raw_sink=raw_sink,
                name=path.name,
            )
        return result
    finally:
        workbook.close()


def _xlsx_selection(
    worksheet: Any,
    *,
    cell_range: str | None,
    rows: str | None,
    max_items: int,
    raw_sink: RawSink | None,
    name: str,
) -> dict[str, Any]:
    if cell_range is not None and rows is not None:
        raise _invalid_selector(cell_range, "Use either a cell range or row range, not both.")
    if cell_range is not None:
        try:
            min_col, min_row, max_col, max_row = range_boundaries(cell_range)
        except ValueError as exc:
            raise _invalid_selector(cell_range, "Cell range is invalid.") from exc
    elif rows is not None:
        min_row, max_row = parse_inclusive_range(rows, target=rows)
        min_col, max_col = 1, max(1, worksheet.max_column)
    else:
        min_row, max_row = 1, min(max(1, worksheet.max_row), max_items)
        min_col, max_col = 1, max(1, worksheet.max_column)
    cell_count = (max_row - min_row + 1) * (max_col - min_col + 1)
    if cell_count > _MAX_RANGE_CELLS:
        raise AerError(
            "LIMIT_EXCEEDED",
            "Requested worksheet range exceeds the inspection cell limit.",
            operation="inspect",
            target=f"{worksheet.title}!{cell_range or rows}",
            details={"cells": cell_count, "limit": _MAX_RANGE_CELLS},
        )
    full_rows: list[dict[str, Any]] = []
    for row_number in range(min_row, max_row + 1):
        cells = []
        for column_number in range(min_col, max_col + 1):
            cell: Cell = worksheet.cell(row=row_number, column=column_number)
            cells.append(
                {
                    "cell": f"{get_column_letter(column_number)}{row_number}",
                    "value": _cell_value(cell.value),
                    "type": cell.data_type,
                }
            )
        full_rows.append({"row": row_number, "cells": cells})
    preview_rows = full_rows[:max_items]
    preview = [{**row, "cells": row["cells"][:_MAX_PREVIEW_COLUMNS]} for row in preview_rows]
    truncated = len(full_rows) > max_items or max_col - min_col + 1 > _MAX_PREVIEW_COLUMNS
    selection: dict[str, Any] = {
        "sheet": worksheet.title,
        "range": f"{get_column_letter(min_col)}{min_row}:{get_column_letter(max_col)}{max_row}",
        "rows": preview,
        "truncated": truncated,
    }
    if truncated:
        selection["raw_ref"] = preserve_overflow(
            full_rows, raw_sink=raw_sink, name=f"{name}.{worksheet.title}.range.json"
        )
    return selection


def inspect_pptx(
    path: Path,
    *,
    slide: int | None,
    selector: str | None,
    query: str | None,
    max_items: int,
    raw_sink: RawSink | None,
    full: bool = False,
) -> dict[str, Any]:
    _check_ooxml_package(path, required_part="ppt/presentation.xml")
    try:
        presentation = Presentation(str(path))
    except (OSError, ValueError, KeyError, zipfile.BadZipFile, PptxPackageNotFoundError) as exc:
        raise _corrupt(path, "PPTX could not be opened.", exc) from exc
    exact_summaries = [
        _slide_summary(item, index=index, full_text=True)
        for index, item in enumerate(presentation.slides, start=1)
    ]
    summaries = (
        exact_summaries
        if full
        else [
            _slide_summary(item, index=index)
            for index, item in enumerate(presentation.slides, start=1)
        ]
    )
    summary_text_truncated = any(item.get("title_truncated", False) for item in summaries)
    result: dict[str, Any] = {
        "type": "pptx",
        "slide_count": len(presentation.slides),
        "width_emu": presentation.slide_width,
        "height_emu": presentation.slide_height,
        "slides": summaries[:max_items],
    }
    if len(summaries) > max_items or summary_text_truncated:
        result["truncated"] = True
        result["raw_ref"] = preserve_overflow(
            exact_summaries, raw_sink=raw_sink, name=f"{path.name}.slides.json"
        )

    selected_slide, shape_selector, title_only = _select_pptx(
        presentation, summaries, slide=slide, selector=selector
    )
    if selected_slide is not None:
        selected_index = _slide_index(presentation, selected_slide)
        exact_slide_summary = _slide_summary(selected_slide, index=selected_index, full_text=True)
        if title_only:
            exact_title = _slide_title(selected_slide)
            exact_selection = {
                "slide": exact_slide_summary,
                "title": exact_title,
            }
            if full:
                result["selection"] = exact_selection
            else:
                title, title_truncated = _preview_text(exact_title)
                result["selection"] = {
                    "slide": _slide_summary(selected_slide, index=selected_index),
                    "title": title,
                }
                if title_truncated or result["selection"]["slide"].get("title_truncated"):
                    result["selection"]["text_truncated"] = True
                    result["selection"]["raw_ref"] = preserve_overflow(
                        [exact_selection],
                        raw_sink=raw_sink,
                        name=f"{path.name}.slide-{selected_index}-title.json",
                    )
        elif shape_selector is not None:
            shape = _find_shape(selected_slide, shape_selector)
            exact_selection = _shape_record(shape, full_text=True)
            result["selection"] = exact_selection if full else _shape_record(shape)
            if result["selection"].get("text_truncated"):
                result["selection"]["raw_ref"] = preserve_overflow(
                    [exact_selection],
                    raw_sink=raw_sink,
                    name=f"{path.name}.slide-{selected_index}-shape.json",
                )
        else:
            exact_shapes = [_shape_record(shape, full_text=True) for shape in selected_slide.shapes]
            shapes = (
                exact_shapes if full else [_shape_record(shape) for shape in selected_slide.shapes]
            )
            shape_text_truncated = any(item.get("text_truncated", False) for item in shapes)
            result["selection"] = {
                **(
                    exact_slide_summary
                    if full
                    else _slide_summary(selected_slide, index=selected_index)
                ),
                "shapes": shapes[:max_items],
                "truncated": len(shapes) > max_items or shape_text_truncated,
            }
            if len(shapes) > max_items or shape_text_truncated:
                result["selection"]["raw_ref"] = preserve_overflow(
                    exact_shapes,
                    raw_sink=raw_sink,
                    name=f"{path.name}.slide-{selected_index}.json",
                )
    if query is not None:
        folded = query.casefold()
        exact_matches: list[dict[str, Any]] = []
        for index, item in enumerate(presentation.slides, start=1):
            slide_id = _stable_slide_id(item) or str(item.slide_id)
            for shape in item.shapes:
                text = _shape_text(shape)
                if text and folded in text.casefold():
                    exact_matches.append(
                        {
                            "slide": index,
                            "slide_id": slide_id,
                            "shape_id": _stable_shape_id(shape) or shape.name,
                            "text": text,
                        }
                    )
        matches = exact_matches if full else [_compact_text_record(item) for item in exact_matches]
        text_truncated = any(item.get("text_truncated", False) for item in matches)
        result["query"] = query
        result["match_count"] = len(matches)
        result["matches"] = matches[:max_items]
        if len(matches) > max_items or text_truncated:
            result["truncated"] = True
            result["raw_ref"] = preserve_overflow(
                exact_matches, raw_sink=raw_sink, name=f"{path.name}.matches.json"
            )
    return result


def inspect_docx(
    path: Path,
    *,
    outline: bool,
    selector: str | None,
    query: str | None,
    max_items: int,
    raw_sink: RawSink | None,
    full: bool = False,
) -> dict[str, Any]:
    _check_ooxml_package(path, required_part="word/document.xml")
    try:
        document = Document(str(path))
    except (OSError, ValueError, KeyError, zipfile.BadZipFile, DocxPackageNotFoundError) as exc:
        raise _corrupt(path, "DOCX could not be opened.", exc) from exc
    exact_paragraphs = [
        _paragraph_record(item, index, full_text=True)
        for index, item in enumerate(document.paragraphs, 1)
    ]
    paragraphs = (
        exact_paragraphs
        if full
        else [_paragraph_record(item, index) for index, item in enumerate(document.paragraphs, 1)]
    )
    exact_headings = [record for record in exact_paragraphs if record["heading_level"] is not None]
    headings = [record for record in paragraphs if record["heading_level"] is not None]
    heading_text_truncated = any(item.get("text_truncated", False) for item in headings)
    result: dict[str, Any] = {
        "type": "docx",
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "image_count": len(document.inline_shapes),
        "section_count": len(document.sections),
        "headings": headings[:max_items],
    }
    if outline:
        result["outline"] = headings[:max_items]
    if len(headings) > max_items or heading_text_truncated:
        result["truncated"] = True
        result["raw_ref"] = preserve_overflow(
            exact_headings, raw_sink=raw_sink, name=f"{path.name}.headings.json"
        )
    if selector is not None:
        result["selector"] = selector
        exact_selection = _select_docx(document, exact_paragraphs, selector, max_items=max_items)
        if full:
            result["selection"] = exact_selection
        else:
            selection, selection_text_truncated = _compact_docx_selection(exact_selection)
            result["selection"] = selection
            if selection_text_truncated and isinstance(selection, dict):
                selection["truncated"] = True
                selection["raw_ref"] = preserve_overflow(
                    [exact_selection],
                    raw_sink=raw_sink,
                    name=f"{path.name}.selection.json",
                )
    if query is not None:
        folded = query.casefold()
        exact_matches = [
            record
            for record in exact_paragraphs
            if folded in str(record.get("text", "")).casefold()
        ]
        for table_index, table in enumerate(document.tables, start=1):
            for row_index, row in enumerate(table.rows, start=1):
                values = [cell.text for cell in row.cells]
                if any(folded in value.casefold() for value in values):
                    exact_matches.append(
                        {"table": table_index, "row": row_index, "values": values[:20]}
                    )
        if full:
            matches = exact_matches
            text_truncated = False
        else:
            matches = []
            text_truncated = False
            for item in exact_matches:
                compacted, item_truncated = _compact_docx_selection(item)
                matches.append(compacted)
                text_truncated = text_truncated or item_truncated
        result["query"] = query
        result["match_count"] = len(matches)
        result["matches"] = matches[:max_items]
        if len(matches) > max_items or text_truncated:
            result["truncated"] = True
            result["raw_ref"] = preserve_overflow(
                exact_matches, raw_sink=raw_sink, name=f"{path.name}.matches.json"
            )
    return result


def inspect_pdf(
    path: Path,
    *,
    page: int | None,
    selector: str | None,
    query: str | None,
    max_items: int,
    raw_sink: RawSink | None,
    full: bool = False,
) -> dict[str, Any]:
    source = ensure_bounded_pdf_input(path, operation="inspect")
    try:
        reader = PdfReader(source, strict=False)
    except (OSError, PdfReadError, ValueError) as exc:
        raise _corrupt(source, "PDF could not be opened.", exc) from exc
    encrypted = bool(reader.is_encrypted)
    if encrypted:
        return {
            "type": "pdf",
            "encrypted": True,
            "page_count": None,
            "content_omitted": True,
        }
    page_count = bounded_pdf_page_count(reader, path=source, operation="inspect")
    pages = reader.pages
    page_summaries = [
        {
            "page": index,
            "width_points": round(float(item.mediabox.width), 3),
            "height_points": round(float(item.mediabox.height), 3),
        }
        for index, item in enumerate(pages, start=1)
    ]
    metadata = {str(key): str(value) for key, value in (reader.metadata or {}).items()}
    attachment_info = pdf_attachment_names(
        reader, path=source, operation="inspect", max_items=max_items
    )
    result: dict[str, Any] = {
        "type": "pdf",
        "encrypted": False,
        "page_count": page_count,
        "metadata": metadata,
        "pages": page_summaries[:max_items],
        "attachments": attachment_info["names"],
        "attachment_count": attachment_info["count"],
        "attachments_truncated": attachment_info["truncated"],
    }
    selected_page = page or _pdf_selector_page(selector)
    if selected_page is not None:
        if selected_page < 1 or selected_page > len(pages):
            raise _invalid_selector(str(selected_page), "PDF page is outside the document.")
        extracted = extract_pdf_page_text(source, selected_page, operation="inspect")
        text = str(extracted["text"])
        preview_lines = _first_pdf_lines(text, max_items)
        if not full:
            preview_lines = [compact_line(line) for line in preview_lines]
        line_count = int(extracted["line_count"])
        text_truncated = any(
            preview != exact.rstrip("\r\n")
            for preview, exact in zip(preview_lines, _first_pdf_lines(text, max_items), strict=True)
        )
        extraction_truncated = bool(extracted.get("extraction_truncated"))
        result["selection"] = {
            **page_summaries[selected_page - 1],
            "text": preview_lines,
            "line_count": line_count,
            "extracted_text_bytes": int(extracted["text_bytes"]),
            "truncated": line_count > max_items or text_truncated or extraction_truncated,
        }
        if text_truncated:
            result["selection"]["text_truncated"] = True
        if extraction_truncated:
            result["selection"]["extraction_truncated"] = True
        if line_count > max_items or text_truncated or extraction_truncated:
            result["selection"]["raw_ref"] = (
                raw_sink(text.encode("utf-8"), f"{path.name}.page-{selected_page}.txt")
                if raw_sink is not None
                else None
            )
            result["selection"]["raw_content_complete"] = not extraction_truncated
    if query is not None:
        searched = search_pdf_text(source, query, operation="inspect")
        raw_matches = searched.get("matches")
        if not isinstance(raw_matches, list) or not all(
            isinstance(item, dict) for item in raw_matches
        ):
            raise AerError(
                "CORRUPT_FILE",
                "PDF text search worker returned invalid matches.",
                "inspect",
                str(source),
            )
        exact_matches = [dict(item) for item in raw_matches]
        matches = exact_matches if full else [_compact_text_record(item) for item in exact_matches]
        text_truncated = any(item.get("text_truncated", False) for item in matches)
        extraction_truncated = bool(searched.get("truncated"))
        match_count = int(searched.get("match_count", len(exact_matches)))
        result["query"] = query
        result["match_count"] = match_count
        result["matches"] = matches[:max_items]
        if match_count > max_items or text_truncated or extraction_truncated:
            result["truncated"] = True
            result["raw_ref"] = preserve_overflow(
                exact_matches, raw_sink=raw_sink, name=f"{path.name}.matches.json"
            )
        if extraction_truncated:
            result["extraction_truncated"] = True
            result["raw_content_complete"] = False
    return result


def _first_pdf_lines(value: str, limit: int) -> list[str]:
    lines: list[str] = []
    for line in io.StringIO(value):
        lines.append(line.rstrip("\r\n"))
        if len(lines) >= limit:
            break
    return lines


def _parse_xlsx_selector(
    selector: str | None,
) -> tuple[str | None, str | None, str | None]:
    if selector is None:
        return None, None, None
    if "!" in selector:
        sheet, cell_range = selector.rsplit("!", 1)
        return sheet, cell_range, None
    stable_cell_match = re.fullmatch(r"sheet:(?:id=)?([^/]+)/cell:id=(.+)", selector)
    if stable_cell_match is not None:
        return stable_cell_match.group(1), None, stable_cell_match.group(2)
    match = re.fullmatch(r"sheet:(?:id=)?([^/]+)(?:/(?:range|cell):(.+))?", selector)
    if match is None:
        raise _invalid_selector(
            selector,
            "Use SHEET!A1:F20, sheet:id=ID/range:A1:F20, or sheet:id=ID/cell:id=ID.",
        )
    return match.group(1), match.group(2), None


def _cell_value(value: Any) -> Any:
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    return value


def _preview_text(value: str) -> tuple[str, bool]:
    flattened = value.replace("\n", " ")
    preview = compact_line(flattened)
    return preview, preview != value


def _compact_text_record(record: dict[str, Any]) -> dict[str, Any]:
    exact = str(record["text"])
    preview, truncated = _preview_text(exact)
    result = {**record, "text": preview}
    if truncated:
        result["text_truncated"] = True
    return result


def _compact_docx_selection(value: Any) -> tuple[Any, bool]:
    if isinstance(value, list):
        result: list[Any] = []
        truncated = False
        for item in value:
            compacted, item_truncated = _compact_docx_selection(item)
            result.append(compacted)
            truncated = truncated or item_truncated
        return result, truncated
    if isinstance(value, dict):
        result_dict: dict[str, Any] = {}
        truncated = False
        text_truncated = False
        for key, item in value.items():
            if key == "text" and isinstance(item, str):
                compacted, item_truncated = _preview_text(item)
            else:
                compacted, item_truncated = _compact_docx_selection(item)
            result_dict[key] = compacted
            truncated = truncated or item_truncated
            text_truncated = text_truncated or (key == "text" and item_truncated)
        if text_truncated:
            result_dict["text_truncated"] = True
        return result_dict, truncated
    return value, False


def _slide_summary(slide: Any, *, index: int, full_text: bool = False) -> dict[str, Any]:
    shapes = list(slide.shapes)
    stable_id = _stable_slide_id(slide)
    exact_title = _slide_title(slide)
    title, title_truncated = (exact_title, False) if full_text else _preview_text(exact_title)
    record = {
        "index": index,
        "id": stable_id or str(slide.slide_id),
        "id_source": "stable" if stable_id else "native",
        "native_id": slide.slide_id,
        "title": title,
        "shape_count": len(shapes),
        "table_count": sum(1 for shape in shapes if getattr(shape, "has_table", False)),
        "chart_count": sum(1 for shape in shapes if getattr(shape, "has_chart", False)),
        "image_count": sum(1 for shape in shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE),
        "warnings": [] if exact_title.strip() else ["empty_title"],
    }
    if title_truncated:
        record["title_truncated"] = True
    return record


def _shape_record(shape: Any, *, full_text: bool = False) -> dict[str, Any]:
    stable_id = _stable_shape_id(shape)
    record: dict[str, Any] = {
        "id": stable_id or shape.name,
        "id_source": "stable" if stable_id else "native_name",
        "native_name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }
    text = _shape_text(shape)
    if text:
        preview, text_truncated = (text, False) if full_text else _preview_text(text)
        record["text"] = preview
        if text_truncated:
            record["text_truncated"] = True
    return record


def _shape_text(shape: Any) -> str:
    return str(shape.text) if getattr(shape, "has_text_frame", False) else ""


def _slide_title(slide: Any) -> str:
    title = slide.shapes.title
    if title is not None and getattr(title, "has_text_frame", False):
        return str(title.text)
    for shape in slide.shapes:
        if shape.name.endswith("/title") and getattr(shape, "has_text_frame", False):
            return str(shape.text)
    return ""


def _stable_slide_id(slide: Any) -> str | None:
    native_name = slide.element.cSld.get("name")
    if native_name and native_name.startswith("aer:"):
        return str(native_name)[4:]
    for shape in slide.shapes:
        aer_match = _AER_SHAPE_PATH_RE.fullmatch(shape.name)
        if aer_match:
            return aer_match.group(1)
        match = _SLIDE_ID_RE.search(shape.name)
        if match:
            return match.group(1)
    try:
        notes_text = "\n".join(
            shape.text for shape in slide.notes_slide.shapes if shape.has_text_frame
        )
    except (AttributeError, KeyError):
        return None
    marker = re.search(r"AER_SLIDE_ID:\s*([^\s]+)", notes_text)
    return marker.group(1) if marker else None


def _stable_shape_id(shape: Any) -> str | None:
    aer_match = _AER_SHAPE_PATH_RE.fullmatch(shape.name)
    if aer_match:
        return aer_match.group(2)
    for expression in (_SHAPE_ID_RE, _AER_ID_RE):
        match = expression.search(shape.name)
        if match:
            return match.group(1)
    return None


def _select_pptx(
    presentation: Any,
    summaries: list[dict[str, Any]],
    *,
    slide: int | None,
    selector: str | None,
) -> tuple[Any | None, str | None, bool]:
    slide_key: str | None = None
    shape_key: str | None = None
    title_only = False
    if selector is not None:
        match = re.fullmatch(r"slide:(?:id=)?([^/]+)(?:/(.+))?", selector)
        if match is None:
            raise _invalid_selector(selector, "PPTX selector must start with slide: or slide:id=.")
        slide_key = match.group(1)
        remainder = match.group(2)
        if remainder == "title":
            title_only = True
        elif remainder is not None:
            shape_match = re.fullmatch(r"shape:(?:id=)?(.+)", remainder)
            if shape_match is None:
                raise _invalid_selector(selector, "Shape selector must use shape:id=ID.")
            shape_key = shape_match.group(1)
    if slide is not None:
        if slide < 1 or slide > len(presentation.slides):
            raise _invalid_selector(str(slide), "Slide number is outside the presentation.")
        if slide_key is not None and str(summaries[slide - 1]["id"]) != slide_key:
            raise _invalid_selector(selector or "", "Selector and --slide disagree.")
        return presentation.slides[slide - 1], shape_key, title_only
    if slide_key is None:
        return None, None, False
    for index, summary in enumerate(summaries):
        if slide_key in {str(summary["id"]), str(summary["native_id"]), str(index + 1)}:
            return presentation.slides[index], shape_key, title_only
    raise _invalid_selector(selector or slide_key, "Slide ID was not found.")


def _find_shape(slide: Any, key: str) -> Any:
    for shape in slide.shapes:
        if key in {shape.name, _stable_shape_id(shape)}:
            return shape
    raise _invalid_selector(key, "Shape ID or native name was not found.")


def _slide_index(presentation: Any, selected: Any) -> int:
    for index, slide in enumerate(presentation.slides, start=1):
        if slide is selected:
            return index
    raise _invalid_selector(str(selected.slide_id), "Slide no longer belongs to the presentation.")


def _paragraph_record(paragraph: Any, index: int, *, full_text: bool = False) -> dict[str, Any]:
    heading_match = _HEADING_RE.match(paragraph.style.name or "")
    stable_id = paragraph._p.get(qn("w14:paraId"))
    for bookmark in paragraph._p.xpath(".//w:bookmarkStart"):
        name = bookmark.get(qn("w:name"))
        if name and not name.startswith("_"):
            stable_id = name[4:].replace("_", "-") if name.startswith("aer_") else name
            break
    exact_text = paragraph.text
    text, text_truncated = (exact_text, False) if full_text else _preview_text(exact_text)
    record = {
        "index": index,
        "id": stable_id or f"paragraph:{index}",
        "style": paragraph.style.name,
        "heading_level": int(heading_match.group(1)) if heading_match else None,
        "text": text,
    }
    if text_truncated:
        record["text_truncated"] = True
    return record


def _table_record(table: Any, index: int, *, max_items: int) -> dict[str, Any]:
    rows = [[cell.text for cell in row.cells] for row in table.rows]
    caption = table._tbl.tblPr.find(qn("w:tblCaption"))
    caption_value = caption.get(qn("w:val"), "") if caption is not None else ""
    stable_id = caption_value[4:] if caption_value.startswith("aer:") else None
    return {
        "table": index,
        "id": stable_id or f"table:{index}",
        "id_source": "stable" if stable_id else "native",
        "rows": rows[:max_items],
        "row_count": len(rows),
        "truncated": len(rows) > max_items,
    }


def _select_docx(
    document: DocumentType,
    paragraphs: list[dict[str, Any]],
    selector: str,
    *,
    max_items: int,
) -> Any:
    paragraph_match = re.fullmatch(r"paragraph:(\d+)", selector)
    if paragraph_match:
        index = int(paragraph_match.group(1))
        if index < 1 or index > len(paragraphs):
            raise _invalid_selector(selector, "Paragraph is outside the document.")
        return paragraphs[index - 1]
    table_match = re.fullmatch(r"table:(\d+)", selector)
    if table_match:
        index = int(table_match.group(1))
        if index < 1 or index > len(document.tables):
            raise _invalid_selector(selector, "Table is outside the document.")
        return _table_record(document.tables[index - 1], index, max_items=max_items)
    block_match = re.fullmatch(r"(?:block|section):id=(.+)", selector)
    if block_match:
        key = block_match.group(1)
        for index, paragraph in enumerate(paragraphs):
            slug = re.sub(r"[^\w]+", "-", str(paragraph["text"]).casefold()).strip("-")
            if key in {str(paragraph["id"]), slug}:
                if selector.startswith("block:"):
                    return paragraph
                level = paragraph["heading_level"]
                selected = [paragraph]
                for following in paragraphs[index + 1 :]:
                    following_level = following["heading_level"]
                    if (
                        level is not None
                        and following_level is not None
                        and following_level <= level
                    ):
                        break
                    selected.append(following)
                    if len(selected) >= max_items:
                        break
                return {"section_id": key, "paragraphs": selected}
        if selector.startswith("block:"):
            for index, table in enumerate(document.tables, start=1):
                record = _table_record(table, index, max_items=max_items)
                if record["id_source"] == "stable" and record["id"] == key:
                    return record
        raise _invalid_selector(selector, "Stable block or section ID was not found.")
    raise _invalid_selector(
        selector,
        "Use paragraph:N, table:N, block:id=ID, or section:id=ID.",
    )


def _pdf_selector_page(selector: str | None) -> int | None:
    if selector is None:
        return None
    match = re.fullmatch(r"(?:/pages/|page:)(\d+)(?:/text)?", selector)
    if match is None:
        raise _invalid_selector(selector, "PDF selector must use page:N or /pages/N/text.")
    return int(match.group(1))


def _invalid_selector(target: str, message: str) -> AerError:
    return AerError("INVALID_SELECTOR", message, operation="inspect", target=target)


def _check_ooxml_package(path: Path, *, required_part: str) -> None:
    enforce_zip_expansion_limits(path, operation="inspect", target=str(path))
    try:
        with zipfile.ZipFile(path) as package:
            entries = package.infolist()
    except (OSError, zipfile.BadZipFile) as exc:
        raise _corrupt(path, "OOXML package is not a readable ZIP archive.", exc) from exc
    names = {entry.filename for entry in entries}
    unsafe = [
        name
        for name in names
        if name.startswith("/") or ".." in Path(name.replace("\\", "/")).parts
    ]
    if unsafe:
        raise AerError(
            "CORRUPT_FILE",
            "OOXML package contains an unsafe part path.",
            operation="inspect",
            target=str(path),
            details={"part": sorted(unsafe)[0]},
        )
    if "[Content_Types].xml" not in names or required_part not in names:
        raise AerError(
            "CORRUPT_FILE",
            "OOXML package is missing a required part.",
            operation="inspect",
            target=str(path),
            details={"required_part": required_part},
        )


def _corrupt(path: Path, message: str, error: BaseException) -> AerError:
    return AerError(
        "CORRUPT_FILE",
        message,
        operation="inspect",
        target=str(path),
        details={"error": f"{type(error).__name__}: {error}"[:300]},
    )
