from __future__ import annotations

import csv
import json
import mimetypes
import posixpath
import shutil
import subprocess
import tempfile
import zipfile
from collections.abc import Callable
from pathlib import Path, PurePosixPath
from typing import Any
from urllib.parse import unquote
from xml.etree import ElementTree

from docx import Document
from openpyxl import load_workbook
from openpyxl.utils.cell import range_boundaries
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from aer.artifacts.common.spec import manifest_path
from aer.artifacts.workbook.selectors import normalize_stable_selector, unescape_sheet_name
from aer.config import Settings
from aer.errors import AerError
from aer.hashing import sha256_file
from aer.limits import (
    MAX_IMAGE_PIXELS,
    MAX_SVG_DEPTH,
    MAX_SVG_ELEMENTS,
    MAX_TABULAR_CELLS,
    MAX_TEXT_FILE_BYTES,
)
from aer.paths import ensure_regular_input
from aer.pdf.safety import (
    MAX_PDF_TEXT_VALIDATION_PAGES,
    bounded_pdf_page_count,
    ensure_bounded_pdf_input,
    pdf_text_presence,
)
from aer.store import ObjectStore
from aer.yaml_safety import load_yaml_safely
from aer.zip_safety import enforce_zip_expansion_limits, inspect_zip_active_content

OOXML_REQUIRED = {
    ".pptx": {"[Content_Types].xml", "ppt/presentation.xml"},
    ".docx": {"[Content_Types].xml", "word/document.xml"},
    ".xlsx": {"[Content_Types].xml", "xl/workbook.xml"},
}


def _relationship_target(relationship_part: str, target: str) -> str | None:
    target = unquote(target.replace("\\", "/")).split("#", 1)[0]
    if not target:
        return None
    if target.startswith("/"):
        normalized = posixpath.normpath(target.lstrip("/"))
    else:
        relationship_path = PurePosixPath(relationship_part)
        base = (
            PurePosixPath()
            if relationship_part == "_rels/.rels"
            else relationship_path.parent.parent
        )
        normalized = posixpath.normpath((base / target).as_posix())
    if normalized == ".." or normalized.startswith("../"):
        return None
    return normalized


def _zip_checks(
    path: Path, suffix: str
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    stats = enforce_zip_expansion_limits(path, operation="artifact.validate", target=str(path))
    metadata: dict[str, Any] = {
        "zip_entries": stats.entries,
        "uncompressed_bytes": stats.uncompressed_bytes,
    }
    try:
        with zipfile.ZipFile(path) as archive:
            entries = archive.infolist()
            entry_names = [item.filename for item in entries]
            duplicate_names = sorted(
                name for name in set(entry_names) if entry_names.count(name) > 1
            )
            if duplicate_names:
                errors.append(
                    {
                        "code": "CORRUPT_FILE",
                        "message": "OOXML package contains duplicate part names.",
                        "details": {"parts": duplicate_names[:20]},
                    }
                )
            names = {item.filename for item in entries}
            for name in names:
                pure = PurePosixPath(name)
                if pure.is_absolute() or ".." in pure.parts:
                    errors.append(
                        {
                            "code": "PATH_OUTSIDE_ROOT",
                            "message": "Archive contains an unsafe entry.",
                            "target": name,
                        }
                    )
            bad = archive.testzip()
            if bad:
                errors.append(
                    {"code": "CORRUPT_FILE", "message": "ZIP CRC check failed.", "target": bad}
                )
            missing = OOXML_REQUIRED.get(suffix, set()) - names
            if missing:
                errors.append(
                    {
                        "code": "CORRUPT_FILE",
                        "message": "Required OOXML parts are missing.",
                        "details": {"missing": sorted(missing)},
                    }
                )
            embedded = sorted(
                name
                for name in names
                if "/embeddings/" in name or name.lower().endswith("vbaproject.bin")
            )
            if embedded:
                warnings.append(
                    {
                        "code": "SUSPICIOUS_EMBEDDED_OBJECT",
                        "message": "Embedded object or macro payload is present; AER never executes it.",
                        "details": {"parts": embedded[:20]},
                    }
                )
            external: list[str] = []
            for name in names:
                if not name.endswith(".rels"):
                    continue
                try:
                    root = ElementTree.fromstring(archive.read(name))
                except ElementTree.ParseError:
                    errors.append(
                        {
                            "code": "CORRUPT_FILE",
                            "message": "Relationship XML is malformed.",
                            "target": name,
                        }
                    )
                    continue
                for relationship in root:
                    target = relationship.attrib.get("Target", "")
                    if relationship.attrib.get("TargetMode") == "External":
                        external.append(target)
                        continue
                    resolved_target = _relationship_target(name, target)
                    if resolved_target is None or resolved_target not in names:
                        errors.append(
                            {
                                "code": "CORRUPT_FILE",
                                "message": "OOXML relationship target is missing or unsafe.",
                                "target": f"{name} -> {target}",
                            }
                        )
            metadata["external_links"] = external[:100]
    except (OSError, zipfile.BadZipFile, RuntimeError) as exc:
        raise AerError(
            "CORRUPT_FILE", f"Cannot read ZIP/OOXML package: {exc}", "artifact.validate", str(path)
        ) from exc
    return errors, warnings, metadata


def _manifest_check(
    path: Path,
    warnings: list[dict[str, Any]],
    errors: list[dict[str, Any]],
    *,
    actual_selectors: set[str] | None = None,
    normalize_selector: Callable[[str], str] | None = None,
) -> None:
    sidecar = manifest_path(path)
    if not sidecar.is_file():
        warnings.append(
            {"code": "MANIFEST_MISSING", "message": "No AER mapping manifest was found."}
        )
        return
    try:
        payload = json.loads(sidecar.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        errors.append(
            {
                "code": "CORRUPT_FILE",
                "message": "AER mapping manifest cannot be parsed.",
                "target": str(sidecar),
            }
        )
        return
    actual = sha256_file(path)
    if payload.get("version") != 1 or payload.get("artifact") != path.name:
        errors.append(
            {
                "code": "CORRUPT_FILE",
                "message": "AER manifest version or artifact name is invalid.",
                "target": str(sidecar),
            }
        )
    if payload.get("artifact_sha256") != actual:
        errors.append(
            {
                "code": "HASH_MISMATCH",
                "message": "AER manifest hash does not match the artifact.",
                "details": {"expected": payload.get("artifact_sha256"), "actual": actual},
            }
        )
    elements = payload.get("elements")
    if not isinstance(elements, list):
        errors.append(
            {
                "code": "CORRUPT_FILE",
                "message": "AER manifest elements must be an array.",
                "target": str(sidecar),
            }
        )
        return

    def selectors(records: list[Any]) -> list[str]:
        values: list[str] = []
        for record in records:
            if not isinstance(record, dict):
                continue
            if isinstance(record.get("selector"), str):
                values.append(record["selector"])
            if isinstance(record.get("children"), list):
                values.extend(selectors(record["children"]))
        return values

    expected = selectors(elements)
    if len(expected) != len(set(expected)):
        errors.append(
            {
                "code": "VALIDATION_FAILED",
                "message": "AER manifest contains duplicate stable selectors.",
                "target": str(sidecar),
            }
        )
    if actual_selectors is not None:
        normalizer = normalize_selector or (lambda value: value)
        missing = sorted(
            selector
            for selector in expected
            if normalizer(selector) not in {normalizer(value) for value in actual_selectors}
        )
        if missing:
            errors.append(
                {
                    "code": "VALIDATION_FAILED",
                    "message": "AER manifest selectors do not match the artifact.",
                    "details": {"missing": missing[:20]},
                }
            )


def _pptx_checks(path: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise AerError(
            "CORRUPT_FILE", f"Cannot reopen PPTX: {exc}", "artifact.validate", str(path)
        ) from exc
    if presentation.slide_width is None or presentation.slide_height is None:
        raise AerError(
            "CORRUPT_FILE",
            "Presentation slide dimensions are missing.",
            "artifact.validate",
            str(path),
        )
    slide_width, slide_height = int(presentation.slide_width), int(presentation.slide_height)
    if not presentation.slides:
        errors.append(
            {
                "code": "VALIDATION_FAILED",
                "message": "Presentation has no slides.",
            }
        )
    stable_ids: set[str] = set()
    actual_selectors: set[str] = set()
    overlap_candidates = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_name = slide.element.cSld.get("name", "")
        if slide_name.startswith("aer:"):
            if slide_name in stable_ids:
                errors.append(
                    {
                        "code": "VALIDATION_FAILED",
                        "message": "Duplicate stable slide ID.",
                        "target": slide_name,
                    }
                )
            stable_ids.add(slide_name)
            actual_selectors.add(f"slide:id={slide_name[4:]}")
        title_found = False
        total_text = 0
        shape_ids: set[str] = set()
        shape_boxes: list[tuple[str, int, int, int, int]] = []
        for shape in slide.shapes:
            if shape.name.startswith("aer:"):
                if shape.name in shape_ids:
                    errors.append(
                        {
                            "code": "VALIDATION_FAILED",
                            "message": "Duplicate stable shape ID.",
                            "target": shape.name,
                        }
                    )
                shape_ids.add(shape.name)
                if "/" in shape.name[4:]:
                    shape_slide_id, shape_id = shape.name[4:].split("/", 1)
                    actual_selectors.add(f"slide:id={shape_slide_id}/shape:id={shape_id}")
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > slide_width
                or shape.top + shape.height > slide_height
            ):
                errors.append(
                    {
                        "code": "VALIDATION_FAILED",
                        "message": "Shape extends outside the slide.",
                        "target": f"slide:{slide_index}/{shape.name}",
                    }
                )
            if shape.width > 0 and shape.height > 0:
                shape_boxes.append(
                    (
                        shape.name,
                        int(shape.left),
                        int(shape.top),
                        int(shape.left + shape.width),
                        int(shape.top + shape.height),
                    )
                )
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                total_text += len(text)
                if shape.name.endswith("/title"):
                    title_found = bool(text)
                    if not text:
                        warnings.append(
                            {
                                "code": "EMPTY_TITLE",
                                "message": "Slide title is empty.",
                                "target": f"slide:{slide_index}",
                            }
                        )
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None and run.font.size.pt < 8:
                            warnings.append(
                                {
                                    "code": "SMALL_FONT",
                                    "message": "Text smaller than 8pt was found.",
                                    "target": f"slide:{slide_index}/{shape.name}",
                                }
                            )
        for first_index, first in enumerate(shape_boxes):
            for second in shape_boxes[first_index + 1 :]:
                intersection_width = max(0, min(first[3], second[3]) - max(first[1], second[1]))
                intersection_height = max(0, min(first[4], second[4]) - max(first[2], second[2]))
                intersection = intersection_width * intersection_height
                if not intersection:
                    continue
                first_area = (first[3] - first[1]) * (first[4] - first[2])
                second_area = (second[3] - second[1]) * (second[4] - second[2])
                if intersection / max(1, min(first_area, second_area)) < 0.2:
                    continue
                overlap_candidates += 1
                if overlap_candidates <= 20:
                    warnings.append(
                        {
                            "code": "OVERLAP_CANDIDATE",
                            "message": "Two slide shapes overlap substantially.",
                            "target": f"slide:{slide_index}/{first[0]}|{second[0]}",
                        }
                    )
        if total_text > 2400:
            warnings.append(
                {
                    "code": "TEXT_DENSITY",
                    "message": "Slide text density exceeds the recommended limit.",
                    "target": f"slide:{slide_index}",
                }
            )
        if not title_found and slide_name.startswith("aer:"):
            warnings.append(
                {
                    "code": "TITLE_NOT_IDENTIFIED",
                    "message": "No stable title shape was identified.",
                    "target": f"slide:{slide_index}",
                }
            )
    _manifest_check(path, warnings, errors, actual_selectors=actual_selectors)
    return (
        errors,
        warnings,
        {
            "slide_count": len(presentation.slides),
            "slide_width": slide_width,
            "slide_height": slide_height,
            "stable_slide_ids": len(stable_ids),
            "overlap_candidates": overlap_candidates,
        },
    )


def _docx_checks(path: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    try:
        document = Document(str(path))
    except Exception as exc:
        raise AerError(
            "CORRUPT_FILE", f"Cannot reopen DOCX: {exc}", "artifact.validate", str(path)
        ) from exc
    if not document.paragraphs and not document.tables:
        errors.append({"code": "VALIDATION_FAILED", "message": "Document is empty."})
    bookmark_names: set[str] = set()
    actual_selectors: set[str] = set()
    last_heading = 0
    for paragraph in document.paragraphs:
        for bookmark in paragraph._p.findall(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}bookmarkStart"
        ):
            name = bookmark.get(
                "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}name", ""
            )
            if name.startswith("aer_"):
                if name in bookmark_names:
                    errors.append(
                        {
                            "code": "VALIDATION_FAILED",
                            "message": "Duplicate stable block ID.",
                            "target": name,
                        }
                    )
                bookmark_names.add(name)
                actual_selectors.add(f"block:id={name[4:]}")
        style_name = paragraph.style.name if paragraph.style else ""
        if style_name.startswith("Heading "):
            try:
                level = int(style_name.split()[-1])
            except ValueError:
                continue
            if last_heading and level > last_heading + 1:
                warnings.append(
                    {
                        "code": "HEADING_HIERARCHY",
                        "message": "Heading level skips a level.",
                        "target": paragraph.text[:80],
                    }
                )
            last_heading = level
    for index, section in enumerate(document.sections, start=1):
        if not section.page_width or not section.page_height:
            errors.append(
                {
                    "code": "VALIDATION_FAILED",
                    "message": "Document section has invalid page dimensions.",
                    "target": f"section:{index}",
                }
            )
        for margin_name in ("top_margin", "right_margin", "bottom_margin", "left_margin"):
            margin = getattr(section, margin_name)
            if margin is not None and margin < 0:
                errors.append(
                    {
                        "code": "VALIDATION_FAILED",
                        "message": "Document section has a negative margin.",
                        "target": f"section:{index}/{margin_name}",
                    }
                )
    for table in document.tables:
        caption = table._tbl.tblPr.find(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblCaption"
        )
        if caption is not None:
            value = caption.get(
                "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", ""
            )
            if value.startswith("aer:"):
                actual_selectors.add(f"block:id={value[4:]}")
    _manifest_check(
        path,
        warnings,
        errors,
        actual_selectors=actual_selectors,
        normalize_selector=lambda value: value.replace("-", "_"),
    )
    return (
        errors,
        warnings,
        {
            "paragraph_count": len(document.paragraphs),
            "table_count": len(document.tables),
            "section_count": len(document.sections),
            "stable_block_ids": len(bookmark_names),
            "inline_shapes": len(document.inline_shapes),
        },
    )


def _formula_parentheses_balanced(formula: str) -> bool:
    depth = 0
    in_double_quote = False
    in_single_quote = False
    index = 0
    while index < len(formula):
        character = formula[index]
        if character == '"' and not in_single_quote:
            if in_double_quote and index + 1 < len(formula) and formula[index + 1] == '"':
                index += 2
                continue
            in_double_quote = not in_double_quote
        elif character == "'" and not in_double_quote:
            if in_single_quote and index + 1 < len(formula) and formula[index + 1] == "'":
                index += 2
                continue
            in_single_quote = not in_single_quote
        elif not in_double_quote and not in_single_quote and character == "(":
            depth += 1
        elif not in_double_quote and not in_single_quote and character == ")":
            depth -= 1
            if depth < 0:
                return False
        index += 1
    return depth == 0 and not in_double_quote and not in_single_quote


def _xlsx_checks(path: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    try:
        workbook = load_workbook(path, data_only=False, read_only=False)
    except Exception as exc:
        raise AerError(
            "CORRUPT_FILE", f"Cannot reopen XLSX: {exc}", "artifact.validate", str(path)
        ) from exc
    if not workbook.sheetnames:
        errors.append({"code": "VALIDATION_FAILED", "message": "Workbook has no worksheets."})
    formula_count = 0
    invalid_formulas: list[str] = []
    invalid_cell_types: list[str] = []
    invalid_merges: list[str] = []
    materialized_cells = 0
    for sheet in workbook.worksheets:
        merge_bounds: list[tuple[str, tuple[int, int, int, int]]] = []
        for merged in sheet.merged_cells.ranges:
            try:
                bounds = range_boundaries(str(merged))
            except ValueError:
                invalid_merges.append(f"{sheet.title}!{merged}")
                continue
            for other_name, other in merge_bounds:
                if not (
                    bounds[2] < other[0]
                    or other[2] < bounds[0]
                    or bounds[3] < other[1]
                    or other[3] < bounds[1]
                ):
                    invalid_merges.extend(
                        [f"{sheet.title}!{other_name}", f"{sheet.title}!{merged}"]
                    )
            merge_bounds.append((str(merged), bounds))
        materialized_cells += len(sheet._cells)
        if materialized_cells > MAX_TABULAR_CELLS:
            raise AerError(
                "LIMIT_EXCEEDED",
                "Workbook exceeds the validation cell limit.",
                "artifact.validate",
                str(path),
                {"cells": materialized_cells, "limit": MAX_TABULAR_CELLS},
            )
        for cell in sheet._cells.values():
            if cell.data_type not in {"b", "d", "e", "f", "inlineStr", "n", "s"}:
                invalid_cell_types.append(f"{sheet.title}!{cell.coordinate}")
            if isinstance(cell.value, str) and cell.value.startswith("="):
                formula_count += 1
                if len(cell.value) == 1 or not _formula_parentheses_balanced(cell.value):
                    invalid_formulas.append(f"{sheet.title}!{cell.coordinate}")
    if invalid_formulas:
        errors.append(
            {
                "code": "VALIDATION_FAILED",
                "message": "Basic formula syntax check failed.",
                "details": {"cells": invalid_formulas[:20]},
            }
        )
    if invalid_cell_types:
        errors.append(
            {
                "code": "VALIDATION_FAILED",
                "message": "Workbook contains invalid cell types.",
                "details": {"cells": invalid_cell_types[:20]},
            }
        )
    if invalid_merges:
        errors.append(
            {
                "code": "VALIDATION_FAILED",
                "message": "Workbook contains invalid or overlapping merged ranges.",
                "details": {"ranges": sorted(set(invalid_merges))[:20]},
            }
        )
    broken_names: list[str] = []
    for defined_name in workbook.defined_names.values():
        if "!" not in str(defined_name.attr_text):
            continue
        try:
            destinations = list(defined_name.destinations)
        except (AttributeError, TypeError, ValueError):
            broken_names.append(str(defined_name.name))
            continue
        for sheet_name, coordinate in destinations:
            # openpyxl returns doubled apostrophes from quoted sheet references
            # (for example O''Brien); undo Excel's quote escaping before lookup.
            resolved_sheet_name = unescape_sheet_name(sheet_name)
            if resolved_sheet_name not in workbook.sheetnames:
                broken_names.append(str(defined_name.name))
                continue
            try:
                range_boundaries(coordinate.replace("$", ""))
            except ValueError:
                broken_names.append(str(defined_name.name))
    if broken_names:
        errors.append(
            {
                "code": "VALIDATION_FAILED",
                "message": "Workbook contains broken named ranges.",
                "details": {"names": sorted(set(broken_names))[:20]},
            }
        )
    normalized_sheet_ids = {
        str(name).removeprefix("aer_sheet_")
        for name in workbook.defined_names
        if str(name).startswith("aer_sheet_")
    }
    actual_selectors = {f"sheet:id={sheet_id}" for sheet_id in normalized_sheet_ids}
    for raw_name in workbook.defined_names:
        name = str(raw_name)
        for sheet_id in sorted(normalized_sheet_ids, key=len, reverse=True):
            prefix = f"aer_{sheet_id}_"
            if name.startswith(prefix) and name != f"aer_sheet_{sheet_id}":
                cell_id = name.removeprefix(prefix)
                if cell_id:
                    actual_selectors.add(f"sheet:id={sheet_id}/cell:id={cell_id}")
                break
    _manifest_check(
        path,
        warnings,
        errors,
        actual_selectors=actual_selectors,
        normalize_selector=normalize_stable_selector,
    )
    return (
        errors,
        warnings,
        {
            "sheet_count": len(workbook.sheetnames),
            "sheets": workbook.sheetnames,
            "formula_count": formula_count,
            "materialized_cells": materialized_cells,
            "defined_names": len(list(workbook.defined_names.values())),
            "charts": sum(len(sheet._charts) for sheet in workbook.worksheets),
            "merged_ranges": sum(len(sheet.merged_cells.ranges) for sheet in workbook.worksheets),
        },
    )


def _pdf_checks(path: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    source = ensure_bounded_pdf_input(path, operation="artifact.validate")
    try:
        reader = PdfReader(source)
    except Exception as exc:
        raise AerError(
            "CORRUPT_FILE", f"Cannot open PDF: {exc}", "artifact.validate", str(source)
        ) from exc
    if reader.is_encrypted:
        raise AerError(
            "UNSUPPORTED_FORMAT",
            "Encrypted PDFs require decryption before validation.",
            "artifact.validate",
            str(source),
        )
    page_count = bounded_pdf_page_count(reader, path=source, operation="artifact.validate")
    checked_text_pages = min(page_count, MAX_PDF_TEXT_VALIDATION_PAGES)
    text_presence = pdf_text_presence(
        source,
        list(range(1, checked_text_pages + 1)),
        operation="artifact.validate",
    )
    for index, page in enumerate(reader.pages, start=1):
        box = page.mediabox
        if float(box.width) <= 0 or float(box.height) <= 0:
            errors.append(
                {
                    "code": "VALIDATION_FAILED",
                    "message": "PDF page has an invalid media box.",
                    "target": f"page:{index}",
                }
            )
        if index <= checked_text_pages and not text_presence[index] and not _page_has_xobject(page):
            warnings.append(
                {
                    "code": "EMPTY_PAGE_CANDIDATE",
                    "message": "PDF page has no extractable text or image.",
                    "target": f"page:{index}",
                }
            )
    if page_count > checked_text_pages:
        warnings.append(
            {
                "code": "PDF_TEXT_CHECK_LIMITED",
                "message": "PDF text and empty-page checks were limited to the bounded page prefix.",
                "details": {"pages_checked": checked_text_pages, "page_count": page_count},
            }
        )
    if page_count == 0:
        errors.append({"code": "VALIDATION_FAILED", "message": "PDF has no pages."})
    return (
        errors,
        warnings,
        {
            "page_count": page_count,
            "text_pages_checked": checked_text_pages,
            "encrypted": reader.is_encrypted,
            "metadata": {str(key): str(value) for key, value in (reader.metadata or {}).items()},
        },
    )


def _page_has_xobject(page: Any) -> bool:
    """Check for image/form XObjects without decoding their streams."""

    try:
        resources = page.get("/Resources")
        if resources is None:
            return False
        resources = resources.get_object()
        xobjects = resources.get("/XObject")
        if xobjects is None:
            return False
        return bool(xobjects.get_object())
    except Exception:
        # An unresolved resource is not proof that the page is empty.
        return True


def _image_checks(path: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    try:
        with Image.open(path) as image:
            image.verify()
        with Image.open(path) as image:
            width, height = image.size
            format_name = image.format
            mode = image.mode
    except Exception as exc:
        raise AerError(
            "CORRUPT_FILE", f"Cannot open image: {exc}", "artifact.validate", str(path)
        ) from exc
    errors = []
    expected_formats = {
        ".png": {"PNG"},
        ".jpg": {"JPEG"},
        ".jpeg": {"JPEG"},
        ".webp": {"WEBP"},
        ".gif": {"GIF"},
        ".tif": {"TIFF"},
        ".tiff": {"TIFF"},
    }
    if format_name not in expected_formats.get(path.suffix.lower(), set()):
        errors.append(
            {
                "code": "VALIDATION_FAILED",
                "message": "Image content format does not match its extension.",
                "details": {"extension": path.suffix.lower(), "detected_format": format_name},
            }
        )
    if width * height > MAX_IMAGE_PIXELS:
        errors.append(
            {"code": "LIMIT_EXCEEDED", "message": "Image pixel count exceeds the safety limit."}
        )
    return errors, [], {"width": width, "height": height, "format": format_name, "mode": mode}


def _svg_checks(path: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    if path.stat().st_size > MAX_TEXT_FILE_BYTES:
        raise AerError(
            "LIMIT_EXCEEDED",
            "SVG exceeds the validation size limit.",
            "artifact.validate",
            str(path),
            {"bytes": path.stat().st_size, "limit": MAX_TEXT_FILE_BYTES},
        )
    root_name = ""
    root_attributes: dict[str, str] = {}
    element_count = 0
    stack: list[ElementTree.Element[str]] = []
    try:
        with path.open("rb") as handle:
            for event, element in ElementTree.iterparse(handle, events=("start", "end")):
                if event == "start":
                    stack.append(element)
                    element_count += 1
                    if element_count > MAX_SVG_ELEMENTS:
                        raise AerError(
                            "LIMIT_EXCEEDED",
                            "SVG exceeds the element-count safety limit.",
                            "artifact.validate",
                            str(path),
                            {"elements": element_count, "limit": MAX_SVG_ELEMENTS},
                        )
                    if len(stack) > MAX_SVG_DEPTH:
                        raise AerError(
                            "LIMIT_EXCEEDED",
                            "SVG exceeds the nesting-depth safety limit.",
                            "artifact.validate",
                            str(path),
                            {"depth": len(stack), "limit": MAX_SVG_DEPTH},
                        )
                    if element_count == 1:
                        root_name = (
                            element.tag.rsplit("}", 1)[-1] if isinstance(element.tag, str) else ""
                        )
                        root_attributes = dict(element.attrib)
                else:
                    if len(stack) > 1:
                        stack[-2].remove(element)
                    element.clear()
                    stack.pop()
    except AerError:
        raise
    except (OSError, ElementTree.ParseError) as exc:
        raise AerError(
            "CORRUPT_FILE",
            f"Cannot parse SVG: {exc}",
            "artifact.validate",
            str(path),
        ) from exc
    if root_name != "svg":
        raise AerError(
            "CORRUPT_FILE",
            "SVG root element is missing.",
            "artifact.validate",
            str(path),
        )
    if element_count == 1:
        raise AerError(
            "VALIDATION_FAILED",
            "SVG contains no renderable structure.",
            "artifact.validate",
            str(path),
        )
    warnings: list[dict[str, Any]] = []
    errors: list[dict[str, Any]] = []
    sidecar = manifest_path(path)
    if sidecar.exists() or sidecar.is_symlink():
        _manifest_check(path, warnings, errors)
    return (
        errors,
        warnings,
        {
            "root": root_name,
            "element_count": element_count,
            "width": root_attributes.get("width"),
            "height": root_attributes.get("height"),
            "view_box": root_attributes.get("viewBox"),
        },
    )


def _raster_pdf(path: Path) -> dict[str, Any]:
    executable = shutil.which("pdftoppm")
    if not executable:
        raise AerError(
            "DEPENDENCY_MISSING",
            "pdftoppm is required for PDF raster validation.",
            "artifact.validate",
            str(path),
            {"dependency": "pdftoppm", "capability": "pdf.render_validate"},
        )
    source = ensure_bounded_pdf_input(path, operation="artifact.validate")
    reader = PdfReader(source)
    page_count = bounded_pdf_page_count(reader, path=source, operation="artifact.validate")
    pages_checked = min(3, page_count)
    if pages_checked == 0:
        return {
            "rasterized": False,
            "pages_checked": 0,
            "preview_refs": [],
            "human_visual_review_required": True,
        }
    for index, page in enumerate(reader.pages[:pages_checked], start=1):
        width = max(0, int(float(page.mediabox.width)))
        height = max(0, int(float(page.mediabox.height)))
        if width * height > MAX_IMAGE_PIXELS:
            raise AerError(
                "LIMIT_EXCEEDED",
                "PDF page exceeds the raster pixel safety limit.",
                "artifact.validate",
                f"page:{index}",
                {"pixels": width * height, "limit": MAX_IMAGE_PIXELS},
            )
    with tempfile.TemporaryDirectory(prefix="aer-pdf-raster-") as directory:
        prefix = Path(directory) / "page"
        try:
            completed = subprocess.run(
                [
                    executable,
                    "-f",
                    "1",
                    "-l",
                    str(pages_checked),
                    "-r",
                    "72",
                    "-png",
                    str(source),
                    str(prefix),
                ],
                shell=False,
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise AerError(
                "COMMAND_TIMEOUT",
                "PDF raster validation timed out.",
                "artifact.validate",
                str(source),
            ) from exc
        previews = sorted(Path(directory).glob("page-*.png"))
        if completed.returncode != 0 or len(previews) != pages_checked:
            raise AerError(
                "COMMAND_FAILED",
                "PDF raster validation failed.",
                "artifact.validate",
                str(source),
                {"exit_code": completed.returncode, "pages_created": len(previews)},
            )
        store = ObjectStore(Settings.load())
        refs: list[str] = []
        for preview in previews:
            image_errors, _, _ = _image_checks(preview)
            if image_errors:
                raise AerError(
                    "LIMIT_EXCEEDED",
                    "Raster preview exceeds the image safety limit.",
                    "artifact.validate",
                    preview.name,
                    {"errors": image_errors},
                )
            refs.append(
                store.put_file(
                    preview,
                    filename=preview.name,
                    mime_type="image/png",
                    source={"operation": "artifact.validate", "input": str(source)},
                ).ref
            )
    return {
        "rasterized": True,
        "pages_checked": pages_checked,
        "preview_refs": refs,
        "human_visual_review_required": True,
    }


def _render_office(path: Path) -> dict[str, Any]:
    active = inspect_zip_active_content(path, operation="artifact.validate", target=str(path))
    if active.external_links or active.active_parts:
        raise AerError(
            "UNSUPPORTED_FORMAT",
            "Render validation rejects external relationships, macros, and executable parts.",
            "artifact.validate",
            str(path),
            {
                "external_links": list(active.external_links),
                "active_parts": list(active.active_parts),
            },
        )
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise AerError(
            "DEPENDENCY_MISSING",
            "LibreOffice is required for Office render validation.",
            "artifact.validate",
            str(path),
            {"dependency": "libreoffice", "capability": "office.render_validate"},
        )
    with tempfile.TemporaryDirectory(prefix="aer-render-") as directory:
        profile = Path(directory) / "lo-profile"
        profile.mkdir(mode=0o700)
        command = [
            executable,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            "--norestore",
            "--safe-mode",
            f"-env:UserInstallation={profile.resolve().as_uri()}",
            "--convert-to",
            "pdf",
            "--outdir",
            directory,
            str(path),
        ]
        try:
            completed = subprocess.run(
                command,
                shell=False,
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise AerError(
                "COMMAND_TIMEOUT",
                "LibreOffice render validation timed out.",
                "artifact.validate",
                str(path),
                {"timeout_seconds": 120},
            ) from exc
        rendered = Path(directory) / f"{path.stem}.pdf"
        if completed.returncode != 0 or not rendered.is_file():
            raise AerError(
                "COMMAND_FAILED",
                "LibreOffice render validation failed.",
                "artifact.validate",
                str(path),
                {"exit_code": completed.returncode},
            )
        rendered = ensure_bounded_pdf_input(rendered, operation="artifact.validate")
        reader = PdfReader(rendered)
        page_count = bounded_pdf_page_count(reader, path=rendered, operation="artifact.validate")
        result: dict[str, Any] = {
            "rendered": True,
            "page_count": page_count,
            "human_visual_review_required": True,
        }
        if shutil.which("pdftoppm"):
            result["raster_preview"] = _raster_pdf(rendered)
        else:
            result["raster_preview"] = {
                "available": False,
                "dependency": "pdftoppm",
            }
        return result


def validate_file(path: Path, *, strict: bool = False, render: bool = False) -> dict[str, Any]:
    source = ensure_regular_input(path, operation="artifact.validate")
    if source.stat().st_size == 0:
        raise AerError("VALIDATION_FAILED", "File is empty.", "artifact.validate", str(source))
    suffix = source.suffix.lower()
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    details: dict[str, Any] = {
        "path": str(source),
        "size": source.stat().st_size,
        "sha256": sha256_file(source),
        "mime": mimetypes.guess_type(source.name)[0] or "application/octet-stream",
    }
    if suffix in OOXML_REQUIRED:
        zip_errors, zip_warnings, zip_details = _zip_checks(source, suffix)
        errors.extend(zip_errors)
        warnings.extend(zip_warnings)
        details.update(zip_details)
        specific = {".pptx": _pptx_checks, ".docx": _docx_checks, ".xlsx": _xlsx_checks}[suffix]
        format_errors, format_warnings, format_details = specific(source)
        errors.extend(format_errors)
        warnings.extend(format_warnings)
        details.update(format_details)
        if render:
            details["render"] = _render_office(source)
    elif suffix == ".pdf":
        format_errors, format_warnings, format_details = _pdf_checks(source)
        errors.extend(format_errors)
        warnings.extend(format_warnings)
        details.update(format_details)
        if render:
            details["render"] = _raster_pdf(source)
    elif suffix == ".svg":
        format_errors, format_warnings, format_details = _svg_checks(source)
        errors.extend(format_errors)
        warnings.extend(format_warnings)
        details.update(format_details)
    elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".tif", ".tiff"}:
        format_errors, format_warnings, format_details = _image_checks(source)
        errors.extend(format_errors)
        warnings.extend(format_warnings)
        details.update(format_details)
    elif suffix == ".zip":
        zip_errors, zip_warnings, zip_details = _zip_checks(source, suffix)
        errors.extend(zip_errors)
        warnings.extend(zip_warnings)
        details.update(zip_details)
    elif suffix in {".json", ".yaml", ".yml", ".jsonl", ".ndjson", ".csv", ".tsv"}:
        if source.stat().st_size > MAX_TEXT_FILE_BYTES:
            raise AerError(
                "LIMIT_EXCEEDED",
                "Structured text exceeds the validation size limit.",
                "artifact.validate",
                str(source),
            )
        try:
            text = source.read_text(encoding="utf-8-sig")
            if suffix == ".json":
                parsed = json.loads(text)
                details["root_type"] = type(parsed).__name__
            elif suffix in {".yaml", ".yml"}:
                parsed = load_yaml_safely(text, operation="artifact.validate", target=str(source))
                details["root_type"] = type(parsed).__name__
            elif suffix in {".jsonl", ".ndjson"}:
                records = [json.loads(line) for line in text.splitlines() if line.strip()]
                details["record_count"] = len(records)
            else:
                rows = list(
                    csv.reader(text.splitlines(), delimiter="\t" if suffix == ".tsv" else ",")
                )
                details["row_count"] = len(rows)
                details["columns"] = len(rows[0]) if rows else 0
        except AerError as exc:
            if exc.code == "LIMIT_EXCEEDED":
                raise
            raise AerError(
                "CORRUPT_FILE",
                exc.message,
                "artifact.validate",
                str(source),
                exc.details,
            ) from exc
        except (UnicodeError, json.JSONDecodeError, csv.Error) as exc:
            raise AerError(
                "CORRUPT_FILE",
                f"Structured text could not be reopened: {exc}",
                "artifact.validate",
                str(source),
            ) from exc
    elif suffix in {".txt", ".md", ".markdown", ".html", ".htm"}:
        if source.stat().st_size > MAX_TEXT_FILE_BYTES:
            raise AerError(
                "LIMIT_EXCEEDED",
                "Text exceeds the validation size limit.",
                "artifact.validate",
                str(source),
            )
        try:
            text = source.read_text(encoding="utf-8-sig")
        except UnicodeError as exc:
            raise AerError(
                "CORRUPT_FILE", "Text is not valid UTF-8.", "artifact.validate", str(source)
            ) from exc
        details["line_count"] = len(text.splitlines())
    else:
        raise AerError(
            "UNSUPPORTED_FORMAT",
            "No validator is available for this file format.",
            "artifact.validate",
            str(source),
        )
    if errors or (strict and warnings):
        raise AerError(
            "VALIDATION_FAILED",
            "Artifact validation failed." if errors else "Strict validation rejected warnings.",
            "artifact.validate",
            str(source),
            {"errors": errors, "warnings": warnings, "checks": details},
            "Inspect the listed targets, patch the artifact, and validate again.",
        )
    return {
        "valid": True,
        "checks": details,
        "warnings": warnings,
        "automatic_checks_only": True,
        "human_visual_review_required": suffix
        in {
            ".pptx",
            ".docx",
            ".xlsx",
            ".pdf",
            ".png",
            ".jpg",
            ".jpeg",
            ".webp",
            ".gif",
            ".tif",
            ".tiff",
            ".svg",
            ".html",
            ".htm",
        },
    }
