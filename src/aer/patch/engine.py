from __future__ import annotations

import copy
import io
import json
import os
import re
import tempfile
from collections.abc import Callable
from pathlib import Path
from typing import Any

import regex as regex_engine
import yaml
from docx import Document
from docx.oxml.ns import qn
from filelock import FileLock, Timeout
from openpyxl import load_workbook
from openpyxl.utils.cell import range_boundaries
from pptx import Presentation
from pptx.chart.data import ChartData

from aer.artifacts.common.spec import manifest_path
from aer.config import Settings
from aer.errors import AerError
from aer.hashing import normalized_hash, sha256_bytes, sha256_file
from aer.limits import (
    MAX_PATCH_TARGET_BYTES,
    MAX_SPEC_FILE_BYTES,
    MAX_TABULAR_CELLS,
    MAX_TEXT_FILE_BYTES,
)
from aer.paths import atomic_write_bytes, ensure_regular_input
from aer.yaml_safety import load_yaml_safely
from aer.zip_safety import enforce_zip_expansion_limits


def load_patch_spec(path: Path) -> dict[str, Any]:
    source = ensure_regular_input(path, operation="artifact.patch")
    if source.stat().st_size > MAX_SPEC_FILE_BYTES:
        raise AerError(
            "LIMIT_EXCEEDED",
            "Patch spec exceeds the size limit.",
            "artifact.patch",
            str(source),
            {"bytes": source.stat().st_size, "limit": MAX_SPEC_FILE_BYTES},
        )
    try:
        text = source.read_text(encoding="utf-8")
        value = (
            json.loads(text)
            if source.suffix.lower() == ".json"
            else load_yaml_safely(text, operation="artifact.patch", target=str(path))
        )
    except AerError as exc:
        if exc.code == "INVALID_SPEC":
            raise AerError(
                "INVALID_PATCH", exc.message, "artifact.patch", str(path), exc.details
            ) from exc
        raise
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        raise AerError(
            "INVALID_PATCH", f"Cannot parse patch spec: {exc}", "artifact.patch", str(path)
        ) from exc
    if not isinstance(value, dict) or value.get("version") != 1:
        raise AerError(
            "INVALID_PATCH", "Patch root must be a version 1 object.", "artifact.patch", "/version"
        )
    operations = value.get("operations")
    if not isinstance(operations, list) or not operations:
        raise AerError(
            "INVALID_PATCH",
            "Patch operations must be a non-empty array.",
            "artifact.patch",
            "/operations",
        )
    for index, operation in enumerate(operations):
        if not isinstance(operation, dict) or not isinstance(operation.get("op"), str):
            raise AerError(
                "INVALID_PATCH",
                "Each patch operation requires op.",
                "artifact.patch",
                f"/operations/{index}",
            )
    return value


def _pointer_parts(pointer: str) -> list[str]:
    if pointer == "":
        return []
    if not pointer.startswith("/"):
        raise AerError(
            "INVALID_SELECTOR",
            "JSON/YAML selector must be an RFC 6901 pointer.",
            "artifact.patch",
            pointer,
        )
    return [part.replace("~1", "/").replace("~0", "~") for part in pointer[1:].split("/")]


def _container(root: Any, pointer: str) -> tuple[Any, str]:
    parts = _pointer_parts(pointer)
    if not parts:
        raise AerError(
            "INVALID_SELECTOR",
            "The root cannot be changed by this operation.",
            "artifact.patch",
            pointer,
        )
    current = root
    for part in parts[:-1]:
        try:
            current = current[int(part)] if isinstance(current, list) else current[part]
        except (KeyError, IndexError, TypeError, ValueError) as exc:
            raise AerError(
                "INVALID_SELECTOR", "Selector does not exist.", "artifact.patch", pointer
            ) from exc
    return current, parts[-1]


def _structured_operation(root: Any, operation: dict[str, Any]) -> None:
    op = str(operation["op"]).split(".", 1)[1]
    pointer = str(operation.get("target", ""))
    container, key = _container(root, pointer)
    if op == "set":
        if isinstance(container, list):
            try:
                container[int(key)] = operation.get("value")
            except (IndexError, ValueError) as exc:
                raise AerError(
                    "INVALID_SELECTOR", "Array index is invalid.", "artifact.patch", pointer
                ) from exc
        elif isinstance(container, dict):
            container[key] = operation.get("value")
        else:
            raise AerError(
                "INVALID_SELECTOR", "Selector parent is not a container.", "artifact.patch", pointer
            )
    elif op == "remove":
        try:
            if isinstance(container, list):
                del container[int(key)]
            else:
                del container[key]
        except (KeyError, IndexError, TypeError, ValueError) as exc:
            raise AerError(
                "INVALID_SELECTOR", "Selector does not exist.", "artifact.patch", pointer
            ) from exc
    elif op == "insert":
        if not isinstance(container, list):
            raise AerError(
                "INVALID_SELECTOR",
                "Insert target parent must be an array.",
                "artifact.patch",
                pointer,
            )
        try:
            index = len(container) if key == "-" else int(key)
        except ValueError as exc:
            raise AerError(
                "INVALID_SELECTOR", "Insert array index is invalid.", "artifact.patch", pointer
            ) from exc
        if index < 0 or index > len(container):
            raise AerError(
                "INVALID_SELECTOR", "Insert array index is out of range.", "artifact.patch", pointer
            )
        container.insert(index, operation.get("value"))
    else:
        raise AerError(
            "INVALID_PATCH", f"Unsupported structured patch operation: {op}", "artifact.patch"
        )


def _safe_regex(pattern: str) -> Any:
    if len(pattern) > 500:
        raise AerError("LIMIT_EXCEEDED", "Regex pattern exceeds 500 characters.", "artifact.patch")
    if re.search(r"\([^)]*[+*][^)]*\)[+*{]", pattern):
        raise AerError(
            "INVALID_PATCH", "Nested unbounded regex quantifiers are not allowed.", "artifact.patch"
        )
    try:
        return regex_engine.compile(pattern)
    except regex_engine.error as exc:
        raise AerError(
            "INVALID_PATCH", f"Invalid regular expression: {exc}", "artifact.patch"
        ) from exc


def _patch_text(data: bytes, operations: list[dict[str, Any]]) -> bytes:
    if len(data) > MAX_TEXT_FILE_BYTES:
        raise AerError(
            "LIMIT_EXCEEDED", "Text input exceeds the patch size limit.", "artifact.patch"
        )
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise AerError(
            "UNSUPPORTED_FORMAT", "Text patch requires UTF-8 input.", "artifact.patch"
        ) from exc
    for operation in operations:
        op = operation["op"]
        old = str(operation.get("old", operation.get("target", "")))
        if op == "text.replace":
            if old not in text:
                raise AerError(
                    "INVALID_SELECTOR", "Replacement text was not found.", "artifact.patch", old
                )
            count = int(operation.get("count", -1))
            text = text.replace(old, str(operation.get("value", operation.get("new", ""))), count)
        elif op == "text.regex_replace":
            pattern = _safe_regex(str(operation.get("pattern", old)))
            try:
                text, replaced = pattern.subn(
                    str(operation.get("value", operation.get("replacement", ""))),
                    text,
                    count=max(0, int(operation.get("count", 0))),
                    timeout=1.0,
                )
            except TimeoutError as exc:
                raise AerError(
                    "LIMIT_EXCEEDED",
                    "Regular expression execution exceeded the safety timeout.",
                    "artifact.patch",
                    pattern.pattern,
                    suggested_action="Use a literal replacement or a simpler bounded regex.",
                ) from exc
            if replaced == 0:
                raise AerError(
                    "INVALID_SELECTOR",
                    "Regular expression matched no text.",
                    "artifact.patch",
                    pattern.pattern,
                )
        else:
            raise AerError("INVALID_PATCH", f"Unsupported text operation: {op}", "artifact.patch")
    return text.encode("utf-8")


def _split_pptx_selector(selector: str) -> tuple[str, str]:
    match = re.fullmatch(r"slide:id=([^/]+)/shape:id=(.+)", selector)
    if not match:
        raise AerError(
            "INVALID_SELECTOR",
            "PPTX selector must be slide:id=.../shape:id=...",
            "presentation.patch",
            selector,
        )
    return match.group(1), match.group(2)


def _pptx_shape(presentation: Any, selector: str) -> Any:
    slide_id, shape_id = _split_pptx_selector(selector)
    for slide in presentation.slides:
        if slide.element.cSld.get("name") == f"aer:{slide_id}":
            for shape in slide.shapes:
                if shape.name == f"aer:{slide_id}/{shape_id}":
                    return shape
            raise AerError(
                "INVALID_SELECTOR", "PPTX shape ID was not found.", "presentation.patch", selector
            )
    raise AerError(
        "INVALID_SELECTOR", "PPTX slide ID was not found.", "presentation.patch", selector
    )


def _replace_paragraph_runs(paragraph: Any, old: str, new: str) -> int:
    if not old:
        return 0
    runs = list(paragraph.runs)
    texts = [run.text for run in runs]
    combined = "".join(texts)
    count = combined.count(old)
    if count == 0:
        return 0
    if not runs:
        paragraph.text = combined.replace(old, new)
        return count

    starts: list[int] = []
    ends: list[int] = []
    cursor = 0
    for text in texts:
        starts.append(cursor)
        cursor += len(text)
        ends.append(cursor)
    positions: list[int] = []
    position = combined.find(old)
    while position >= 0:
        positions.append(position)
        position = combined.find(old, position + len(old))
    for start in reversed(positions):
        end = start + len(old)
        first = next(index for index, boundary in enumerate(ends) if boundary > start)
        last = next(index for index, boundary in enumerate(ends) if boundary >= end)
        first_offset = start - starts[first]
        last_offset = end - starts[last]
        if first == last:
            texts[first] = texts[first][:first_offset] + new + texts[first][last_offset:]
            continue
        texts[first] = texts[first][:first_offset] + new
        for index in range(first + 1, last):
            texts[index] = ""
        texts[last] = texts[last][last_offset:]
    for run, text in zip(runs, texts, strict=True):
        run.text = text
    return count


def _pptx_replace_shape(shape: Any, old: str, new: str) -> int:
    replaced = 0
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            replaced += _replace_paragraph_runs(paragraph, old, new)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    replaced += _replace_paragraph_runs(paragraph, old, new)
    return replaced


def _patch_pptx(data: bytes, operations: list[dict[str, Any]]) -> bytes:
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise AerError("CORRUPT_FILE", f"Cannot open PPTX: {exc}", "presentation.patch") from exc
    for operation in operations:
        op = operation["op"]
        shape = _pptx_shape(presentation, str(operation.get("target", "")))
        if op == "pptx.set_text":
            if not getattr(shape, "has_text_frame", False):
                raise AerError(
                    "UNSUPPORTED_FORMAT",
                    "Target shape has no editable text frame.",
                    "presentation.patch",
                    str(operation.get("target")),
                )
            shape.text = str(operation.get("value", ""))
        elif op == "pptx.replace_text":
            old = str(operation.get("old", ""))
            if (
                not old
                or _pptx_replace_shape(
                    shape, old, str(operation.get("value", operation.get("new", "")))
                )
                == 0
            ):
                raise AerError(
                    "INVALID_SELECTOR",
                    "Replacement text was not found in target shape.",
                    "presentation.patch",
                    str(operation.get("target")),
                )
        elif op == "pptx.remove_shape":
            shape._element.getparent().remove(shape._element)
        elif op == "pptx.update_chart_data":
            if not getattr(shape, "has_chart", False):
                raise AerError(
                    "UNSUPPORTED_FORMAT",
                    "Target shape is not a chart.",
                    "presentation.patch",
                    str(operation.get("target")),
                )
            categories = list(operation.get("categories", []))
            series = list(operation.get("series", []))
            if not categories or not series:
                raise AerError(
                    "INVALID_PATCH",
                    "Chart update requires categories and series.",
                    "presentation.patch",
                )
            chart_data = ChartData()
            chart_data.categories = [str(value) for value in categories]
            for item in series:
                chart_data.add_series(
                    str(item.get("name", "Series")),
                    [float(value) for value in item.get("values", [])],
                )
            shape.chart.replace_data(chart_data)
        else:
            raise AerError(
                "INVALID_PATCH", f"Unsupported PPTX operation: {op}", "presentation.patch"
            )
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _docx_selector_key(selector: str) -> str | None:
    match = re.fullmatch(r"(?:block|section):id=(.+)", selector)
    return match.group(1).replace("-", "_") if match else None


def _docx_block(document: Any, selector: str) -> tuple[str, list[Any]]:
    match = re.fullmatch(r"(?:block|section):id=(.+)", selector)
    if not match:
        raise AerError(
            "INVALID_SELECTOR", "DOCX selector must be block:id=...", "document.patch", selector
        )
    expected = f"aer_{match.group(1).replace('-', '_')}"
    paragraphs = list(document.paragraphs)
    for index, paragraph in enumerate(paragraphs):
        for bookmark in paragraph._p.findall(qn("w:bookmarkStart")):
            if bookmark.get(qn("w:name")) == expected:
                bookmark_id = bookmark.get(qn("w:id"))
                if bookmark_id is None:
                    raise AerError(
                        "CORRUPT_FILE",
                        "DOCX stable block bookmark has no ID.",
                        "document.patch",
                        selector,
                    )
                block: list[Any] = []
                for part in paragraphs[index:]:
                    block.append(part)
                    if any(
                        end.get(qn("w:id")) == bookmark_id
                        for end in part._p.findall(qn("w:bookmarkEnd"))
                    ):
                        return "paragraph", block
                raise AerError(
                    "CORRUPT_FILE",
                    "DOCX stable block bookmark has no matching end marker.",
                    "document.patch",
                    selector,
                )
    for table in document.tables:
        caption = table._tbl.tblPr.find(qn("w:tblCaption"))
        if caption is not None and caption.get(qn("w:val")) == f"aer:{match.group(1)}":
            return "table", [table]
    raise AerError("INVALID_SELECTOR", "DOCX block ID was not found.", "document.patch", selector)


def _set_docx_paragraph(paragraph: Any, value: str) -> None:
    for run in list(paragraph._p.findall(qn("w:r"))):
        paragraph._p.remove(run)
    paragraph.add_run(value)


def _set_docx_block(paragraphs: list[Any], value: str, selector: str) -> None:
    first = paragraphs[0]
    selector_key = _docx_selector_key(selector)
    start = next(
        (
            marker
            for marker in first._p.findall(qn("w:bookmarkStart"))
            if marker.get(qn("w:name")) == f"aer_{selector_key}"
        ),
        None,
    )
    if start is None:
        raise AerError(
            "CORRUPT_FILE",
            "DOCX stable block bookmark start marker is missing.",
            "document.patch",
            selector,
        )
    bookmark_id = start.get(qn("w:id"))
    end = next(
        (
            marker
            for paragraph in paragraphs
            for marker in paragraph._p.findall(qn("w:bookmarkEnd"))
            if marker.get(qn("w:id")) == bookmark_id
        ),
        None,
    )
    if end is None:
        raise AerError(
            "CORRUPT_FILE",
            "DOCX stable block bookmark end marker is missing.",
            "document.patch",
            selector,
        )
    end_parent = end.getparent()
    if end_parent is not None:
        end_parent.remove(end)
    _set_docx_paragraph(first, value)
    first._p.append(end)
    for paragraph in paragraphs[1:]:
        parent = paragraph._p.getparent()
        if parent is not None:
            parent.remove(paragraph._p)


def _all_docx_paragraphs(document: Any) -> list[Any]:
    paragraphs = list(document.paragraphs)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                paragraphs.extend(cell.paragraphs)
    return paragraphs


def _patch_docx(data: bytes, operations: list[dict[str, Any]]) -> bytes:
    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:
        raise AerError("CORRUPT_FILE", f"Cannot open DOCX: {exc}", "document.patch") from exc
    for operation in operations:
        op = operation["op"]
        if op == "docx.replace_text":
            old = str(operation.get("old", ""))
            new = str(operation.get("value", operation.get("new", "")))
            target = operation.get("target")
            paragraphs: list[Any]
            if target:
                kind, blocks = _docx_block(document, str(target))
                if kind != "paragraph":
                    raise AerError(
                        "UNSUPPORTED_FORMAT",
                        "Text replacement target must be a paragraph block.",
                        "document.patch",
                        str(target),
                    )
                paragraphs = blocks
            else:
                paragraphs = _all_docx_paragraphs(document)
            if (
                not old
                or sum(_replace_paragraph_runs(paragraph, old, new) for paragraph in paragraphs)
                == 0
            ):
                raise AerError(
                    "INVALID_SELECTOR", "Replacement text was not found.", "document.patch", old
                )
        elif op == "docx.set_block":
            target = str(operation.get("target", ""))
            kind, blocks = _docx_block(document, target)
            if kind != "paragraph":
                raise AerError(
                    "UNSUPPORTED_FORMAT",
                    "set_block currently supports paragraph blocks.",
                    "document.patch",
                    target,
                )
            _set_docx_block(blocks, str(operation.get("value", "")), target)
        elif op == "docx.remove_block":
            _, blocks = _docx_block(document, str(operation.get("target", "")))
            for block in blocks:
                element = block._p if hasattr(block, "_p") else block._tbl
                parent = element.getparent()
                if parent is not None:
                    parent.remove(element)
        else:
            raise AerError("INVALID_PATCH", f"Unsupported DOCX operation: {op}", "document.patch")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _xlsx_stable_name(sheet_id: str, cell_id: str) -> str:
    sheet = "".join(character if character.isalnum() else "_" for character in sheet_id)
    cell = "".join(character if character.isalnum() else "_" for character in cell_id)
    return f"aer_{sheet}_{cell}"


def _xlsx_target(workbook: Any, target: str) -> tuple[Any, str]:
    native = re.fullmatch(r"(.+)!([A-Za-z]+[0-9]+(?::[A-Za-z]+[0-9]+)?)", target)
    stable_address = re.fullmatch(r"sheet:id=([^/]+)/(?:(?:cell|range)=)(.+)", target)
    stable_cell = re.fullmatch(r"sheet:id=([^/]+)/cell:id=(.+)", target)
    if native:
        sheet_name, address = native.groups()
    elif stable_address or stable_cell:
        matched = stable_address or stable_cell
        assert matched is not None
        stable_id, selected_value = matched.groups()
        address = selected_value if stable_address else ""
        sheet_name = stable_id if stable_id in workbook.sheetnames else None
        defined_name = "aer_sheet_" + "".join(
            character if character.isalnum() else "_" for character in stable_id
        )
        if sheet_name is None and defined_name in workbook.defined_names:
            destinations = list(workbook.defined_names[defined_name].destinations)
            if destinations:
                sheet_name = destinations[0][0]
        if sheet_name is None:
            for name in workbook.sheetnames:
                if name.lower().replace(" ", "-") == stable_id.lower():
                    sheet_name = name
                    break
        if sheet_name is None:
            raise AerError(
                "INVALID_SELECTOR", "Workbook sheet ID was not found.", "workbook.patch", target
            )
        if stable_cell:
            cell_name = _xlsx_stable_name(stable_id, selected_value)
            if cell_name not in workbook.defined_names:
                raise AerError(
                    "INVALID_SELECTOR",
                    "Workbook cell ID was not found.",
                    "workbook.patch",
                    target,
                )
            destinations = list(workbook.defined_names[cell_name].destinations)
            address = next(
                (
                    coordinate.replace("$", "")
                    for destination_sheet, coordinate in destinations
                    if destination_sheet == sheet_name
                ),
                "",
            )
            if not address:
                raise AerError(
                    "INVALID_SELECTOR",
                    "Workbook cell ID does not resolve within the selected sheet.",
                    "workbook.patch",
                    target,
                )
    else:
        raise AerError(
            "INVALID_SELECTOR",
            "XLSX target must be Sheet!A1, sheet:id=.../cell=A1, or sheet:id=.../cell:id=....",
            "workbook.patch",
            target,
        )
    if sheet_name not in workbook.sheetnames:
        raise AerError(
            "INVALID_SELECTOR", "Workbook sheet was not found.", "workbook.patch", target
        )
    return workbook[sheet_name], address


def _xlsx_range_bounds(address: str) -> tuple[int, int, int, int]:
    try:
        bounds = range_boundaries(address)
    except (TypeError, ValueError) as exc:
        raise AerError(
            "INVALID_SELECTOR", "Workbook cell range is invalid.", "workbook.patch", address
        ) from exc
    if any(not isinstance(value, int) for value in bounds):
        raise AerError(
            "INVALID_SELECTOR", "Workbook cell range is invalid.", "workbook.patch", address
        )
    min_col, min_row, max_col, max_row = bounds
    if (
        min_col < 1
        or min_row < 1
        or max_col > 16_384
        or max_row > 1_048_576
        or min_col > max_col
        or min_row > max_row
    ):
        raise AerError(
            "INVALID_SELECTOR",
            "Workbook cell range is outside Excel worksheet bounds.",
            "workbook.patch",
            address,
        )
    cell_count = (max_col - min_col + 1) * (max_row - min_row + 1)
    if cell_count > MAX_TABULAR_CELLS:
        raise AerError(
            "LIMIT_EXCEEDED",
            "Workbook patch range exceeds the cell safety limit.",
            "workbook.patch",
            address,
            {"cells": cell_count, "limit": MAX_TABULAR_CELLS},
        )
    return min_col, min_row, max_col, max_row


def _patch_xlsx(data: bytes, operations: list[dict[str, Any]]) -> bytes:
    try:
        workbook = load_workbook(io.BytesIO(data))
    except Exception as exc:
        raise AerError("CORRUPT_FILE", f"Cannot open XLSX: {exc}", "workbook.patch") from exc
    for operation in operations:
        op = operation["op"]
        sheet, address = _xlsx_target(workbook, str(operation.get("target", "")))
        bounds = _xlsx_range_bounds(address)
        if op == "xlsx.set_cell":
            if ":" in address:
                raise AerError(
                    "INVALID_SELECTOR", "set_cell requires one cell.", "workbook.patch", address
                )
            sheet[address] = operation.get("value")
        elif op == "xlsx.set_range":
            if ":" not in address:
                raise AerError(
                    "INVALID_SELECTOR", "set_range requires a range.", "workbook.patch", address
                )
            values = operation.get("values", operation.get("value"))
            if not isinstance(values, list) or any(not isinstance(row, list) for row in values):
                raise AerError(
                    "INVALID_PATCH",
                    "set_range values must be a two-dimensional array.",
                    "workbook.patch",
                    address,
                )
            min_col, min_row, max_col, max_row = bounds
            if len(values) != max_row - min_row + 1 or any(
                len(row) != max_col - min_col + 1 for row in values
            ):
                raise AerError(
                    "INVALID_PATCH",
                    "set_range dimensions do not match the target range.",
                    "workbook.patch",
                    address,
                )
            for row_offset, row in enumerate(values):
                for column_offset, value in enumerate(row):
                    sheet.cell(min_row + row_offset, min_col + column_offset, value)
        elif op == "xlsx.replace_text":
            old, new = (
                str(operation.get("old", "")),
                str(operation.get("value", operation.get("new", ""))),
            )
            cells = sheet[address] if ":" in address else ((sheet[address],),)
            count = 0
            for row in cells:
                for cell in row:
                    if isinstance(cell.value, str) and old in cell.value:
                        cell.value = cell.value.replace(old, new)
                        count += 1
            if not old or count == 0:
                raise AerError(
                    "INVALID_SELECTOR", "Replacement text was not found.", "workbook.patch", address
                )
        elif op == "xlsx.clear_range":
            cells = sheet[address] if ":" in address else ((sheet[address],),)
            for row in cells:
                for cell in row:
                    cell.value = None
        else:
            raise AerError("INVALID_PATCH", f"Unsupported XLSX operation: {op}", "workbook.patch")
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _patch_structured(data: bytes, operations: list[dict[str, Any]], *, yaml_format: bool) -> bytes:
    try:
        text = data.decode("utf-8")
        root = (
            load_yaml_safely(text, operation="artifact.patch", target="yaml target")
            if yaml_format
            else json.loads(text)
        )
    except AerError as exc:
        if exc.code == "INVALID_SPEC":
            raise AerError(
                "CORRUPT_FILE", exc.message, "artifact.patch", details=exc.details
            ) from exc
        raise
    except (UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise AerError(
            "CORRUPT_FILE", f"Cannot parse structured input: {exc}", "artifact.patch"
        ) from exc
    working = copy.deepcopy(root)
    for operation in operations:
        _structured_operation(working, operation)
    if yaml_format:
        return yaml.safe_dump(working, allow_unicode=True, sort_keys=False).encode("utf-8")
    return (json.dumps(working, ensure_ascii=False, indent=2) + "\n").encode("utf-8")


def _select_patcher(
    suffix: str, operations: list[dict[str, Any]]
) -> Callable[[bytes, list[dict[str, Any]]], bytes]:
    prefixes = {str(operation["op"]).split(".", 1)[0] for operation in operations}
    expected = {
        ".json": {"json"},
        ".yaml": {"yaml"},
        ".yml": {"yaml"},
        ".pptx": {"pptx"},
        ".docx": {"docx"},
        ".xlsx": {"xlsx"},
    }.get(suffix)
    if expected is None:
        expected = {"text"}
    if prefixes != expected:
        raise AerError(
            "INVALID_PATCH",
            "Patch operations do not match the target format.",
            "artifact.patch",
            details={"expected_prefix": sorted(expected), "actual_prefix": sorted(prefixes)},
        )
    if suffix == ".json":
        return lambda data, ops: _patch_structured(data, ops, yaml_format=False)
    if suffix in {".yaml", ".yml"}:
        return lambda data, ops: _patch_structured(data, ops, yaml_format=True)
    if suffix == ".pptx":
        return _patch_pptx
    if suffix == ".docx":
        return _patch_docx
    if suffix == ".xlsx":
        return _patch_xlsx
    return _patch_text


def _load_manifest(target: Path) -> tuple[Path, bytes, dict[str, Any]] | None:
    path = manifest_path(target)
    if path.is_symlink():
        raise AerError(
            "INVALID_ARGUMENT",
            "Symbolic-link manifests are not accepted for patch operations.",
            "artifact.patch",
            str(path),
        )
    if not path.exists():
        return None
    source = ensure_regular_input(path, operation="artifact.patch")
    try:
        original = source.read_bytes()
        payload = json.loads(original.decode("utf-8"))
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        raise AerError(
            "CORRUPT_FILE",
            "AER mapping manifest cannot be parsed.",
            "artifact.patch",
            str(path),
        ) from exc
    if not isinstance(payload, dict) or not isinstance(payload.get("elements"), list):
        raise AerError(
            "CORRUPT_FILE",
            "AER mapping manifest must contain an elements array.",
            "artifact.patch",
            str(path),
        )
    return path, original, payload


def _removed_manifest_selector(selector: str, operations: list[dict[str, Any]]) -> bool:
    for operation in operations:
        op = operation["op"]
        target = str(operation.get("target", ""))
        if op == "pptx.remove_shape" and selector == target:
            return True
        if op == "docx.remove_block":
            selector_key = _docx_selector_key(selector)
            target_key = _docx_selector_key(target)
            if selector_key is not None and selector_key == target_key:
                return True
    return False


def _prune_manifest_elements(elements: list[Any], operations: list[dict[str, Any]]) -> list[Any]:
    retained: list[Any] = []
    for element in elements:
        if not isinstance(element, dict):
            retained.append(element)
            continue
        selector = element.get("selector")
        if isinstance(selector, str) and _removed_manifest_selector(selector, operations):
            continue
        children = element.get("children")
        if isinstance(children, list):
            element["children"] = _prune_manifest_elements(children, operations)
        retained.append(element)
    return retained


def _updated_manifest_payload(
    state: tuple[Path, bytes, dict[str, Any]] | None,
    *,
    digest: str,
    operations: list[dict[str, Any]],
    artifact_name: str,
) -> dict[str, Any] | None:
    if state is None:
        return None
    payload = copy.deepcopy(state[2])
    payload["artifact"] = artifact_name
    payload["artifact_sha256"] = digest
    payload["elements"] = _prune_manifest_elements(payload["elements"], operations)
    return payload


def _manifest_bytes(payload: dict[str, Any]) -> bytes:
    return (json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True) + "\n").encode(
        "utf-8"
    )


def apply_patch(
    target: Path,
    patch_spec: Path,
    *,
    dry_run: bool = False,
    backup: bool = False,
    expected_sha256: str | None = None,
    validate: bool = False,
) -> dict[str, Any]:
    source = ensure_regular_input(target, operation="artifact.patch")
    spec = load_patch_spec(patch_spec)
    operations = spec["operations"]
    settings = Settings.load()
    lock_directory = settings.cache_dir / "patch-locks"
    lock_directory.mkdir(parents=True, exist_ok=True)
    lock = FileLock(lock_directory / f"{normalized_hash(str(source))}.lock", timeout=30)
    try:
        with lock:
            source = ensure_regular_input(source, operation="artifact.patch")
            if source.stat().st_size > MAX_PATCH_TARGET_BYTES:
                raise AerError(
                    "LIMIT_EXCEEDED",
                    "Patch target exceeds the in-memory safety limit.",
                    "artifact.patch",
                    str(source),
                    {"bytes": source.stat().st_size, "limit": MAX_PATCH_TARGET_BYTES},
                )
            original = source.read_bytes()
            before = sha256_bytes(original)
            if expected_sha256 and expected_sha256.lower() != before:
                raise AerError(
                    "HASH_MISMATCH",
                    "Target changed since the expected hash was recorded.",
                    "artifact.patch",
                    str(source),
                    {"expected": expected_sha256.lower(), "actual": before},
                    "Inspect the current file and regenerate the patch.",
                )
            manifest_state = _load_manifest(source)
            patcher = _select_patcher(source.suffix.lower(), operations)
            if source.suffix.lower() in {".pptx", ".docx", ".xlsx"}:
                enforce_zip_expansion_limits(
                    original, operation="artifact.patch", target=str(source)
                )
            changed = patcher(original, operations)
            if changed == original:
                raise AerError(
                    "INVALID_PATCH", "Patch produced no change.", "artifact.patch", str(source)
                )
            changed_digest = sha256_bytes(changed)
            updated_manifest = _updated_manifest_payload(
                manifest_state,
                digest=changed_digest,
                operations=operations,
                artifact_name=source.name,
            )
            summary = [
                {
                    "op": operation["op"],
                    "target": operation.get("target", operation.get("old")),
                }
                for operation in operations
            ]
            validation: dict[str, Any] | None = None
            if validate:
                from aer.validation import validate_file

                descriptor, temporary_name = tempfile.mkstemp(
                    prefix=f".{source.stem}.validate-", suffix=source.suffix, dir=source.parent
                )
                temporary = Path(temporary_name)
                temporary_manifest = manifest_path(temporary)
                try:
                    os.fchmod(descriptor, 0o600)
                    with os.fdopen(descriptor, "wb") as handle:
                        handle.write(changed)
                        handle.flush()
                        os.fsync(handle.fileno())
                    if updated_manifest is not None:
                        validation_manifest = copy.deepcopy(updated_manifest)
                        validation_manifest["artifact"] = temporary.name
                        atomic_write_bytes(temporary_manifest, _manifest_bytes(validation_manifest))
                    validation = validate_file(temporary)
                finally:
                    temporary_manifest.unlink(missing_ok=True)
                    temporary.unlink(missing_ok=True)
            if dry_run:
                return {
                    "target": str(source),
                    "dry_run": True,
                    "before_sha256": before,
                    "planned_operations": summary,
                    "would_change_bytes": len(changed),
                    "validation": validation,
                }
            current = sha256_file(source)
            if current != before:
                raise AerError(
                    "HASH_MISMATCH" if expected_sha256 else "CONFLICT",
                    "Target changed while the patch was being prepared.",
                    "artifact.patch",
                    str(source),
                    {"before": before, "current": current},
                    "Inspect the current file and retry with its current SHA-256.",
                )
            current_manifest_path = manifest_path(source)
            if manifest_state is None:
                if current_manifest_path.exists() or current_manifest_path.is_symlink():
                    raise AerError(
                        "CONFLICT",
                        "The target manifest changed while the patch was being prepared.",
                        "artifact.patch",
                        str(current_manifest_path),
                    )
            else:
                stored_manifest_path, original_manifest, _ = manifest_state
                try:
                    manifest_unchanged = stored_manifest_path.read_bytes() == original_manifest
                except OSError:
                    manifest_unchanged = False
                if not manifest_unchanged:
                    raise AerError(
                        "CONFLICT",
                        "The target manifest changed while the patch was being prepared.",
                        "artifact.patch",
                        str(stored_manifest_path),
                    )
            backup_path: Path | None = None
            if backup:
                backup_path = source.with_name(source.name + ".bak")
                atomic_write_bytes(backup_path, original)
            try:
                atomic_write_bytes(source, changed)
                if manifest_state is not None and updated_manifest is not None:
                    atomic_write_bytes(manifest_state[0], _manifest_bytes(updated_manifest))
            except BaseException:
                atomic_write_bytes(source, original)
                if manifest_state is not None:
                    atomic_write_bytes(manifest_state[0], manifest_state[1])
                raise
            after = sha256_file(source)
            return {
                "target": str(source),
                "before_sha256": before,
                "after_sha256": after,
                "operations": summary,
                "backup": str(backup_path) if backup_path else None,
                "validation": validation,
            }
    except Timeout as exc:
        raise AerError(
            "CONFLICT",
            "Another patch operation is holding the target lock.",
            "artifact.patch",
            str(source),
        ) from exc
