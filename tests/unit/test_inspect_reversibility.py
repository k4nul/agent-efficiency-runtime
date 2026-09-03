from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

import aer.inspect.office as office_inspect
import aer.inspect.repository as repository_inspect
from aer.inspect import inspect_target


class _RawCapture:
    def __init__(self) -> None:
        self.values: dict[str, bytes] = {}

    def __call__(self, data: bytes, name: str) -> str:
        self.values[name] = data
        return "aer://sha256/" + "a" * 64


def test_long_text_range_and_search_are_exactly_recoverable(tmp_path: Path) -> None:
    end_marker = "::EXACT_TEXT_END::"
    exact_line = "needle:" + "x" * 900 + end_marker
    target = tmp_path / "long.log"
    target.write_text(exact_line + "\n", encoding="utf-8")

    range_capture = _RawCapture()
    selected = inspect_target(
        target,
        start_line=1,
        end_line=1,
        raw_sink=range_capture,
    )

    assert selected["truncated"] is True
    assert selected["preview"][0]["text_truncated"] is True
    assert end_marker not in selected["preview"][0]["text"]
    preserved_range = json.loads(range_capture.values["long.log.lines.json"])
    assert preserved_range == [{"line": 1, "text": exact_line}]

    query_capture = _RawCapture()
    searched = inspect_target(target, query="needle", raw_sink=query_capture)
    assert searched["matches"][0]["text_truncated"] is True
    assert json.loads(query_capture.values["long.log.matches.json"])[0]["text"].endswith(end_marker)

    full = inspect_target(target, start_line=1, end_line=1, full=True)
    assert full["preview"] == [{"line": 1, "text": exact_line}]


def test_long_repository_match_preview_has_exact_raw_and_full_forms(
    monkeypatch: Any, tmp_path: Path
) -> None:
    end_marker = "::EXACT_REPOSITORY_END::"
    exact_line = "needle:" + "r" * 900 + end_marker
    (tmp_path / "source.txt").write_text(exact_line + "\n", encoding="utf-8")
    monkeypatch.setattr(repository_inspect.shutil, "which", lambda _name: None)
    capture = _RawCapture()

    result = inspect_target(tmp_path, query="needle", raw_sink=capture)

    assert result["truncated"] is True
    assert result["matches"][0]["text_truncated"] is True
    assert end_marker not in result["matches"][0]["text"]
    preserved = json.loads(capture.values[f"{tmp_path.name}.matches.json"])
    assert preserved[0]["text"] == exact_line

    full = inspect_target(tmp_path, query="needle", full=True)
    assert full["matches"][0]["text"] == exact_line
    assert "text_truncated" not in full["matches"][0]


def test_long_pptx_shape_text_has_exact_raw_and_full_forms(tmp_path: Path) -> None:
    end_marker = "::EXACT_OFFICE_END::"
    exact_text = "needle:" + "p" * 900 + end_marker
    target = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.element.cSld.set("name", "aer:sample")
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    shape.name = "aer:sample/body"
    shape.text = exact_text
    presentation.save(target)

    capture = _RawCapture()
    result = inspect_target(
        target,
        selector="slide:id=sample/shape:id=body",
        raw_sink=capture,
    )

    assert result["selection"]["text_truncated"] is True
    assert end_marker not in result["selection"]["text"]
    preserved = json.loads(capture.values["deck.pptx.slide-1-shape.json"])
    assert preserved[0]["text"] == exact_text

    full = inspect_target(
        target,
        selector="slide:id=sample/shape:id=body",
        full=True,
    )
    assert full["selection"]["text"] == exact_text
    assert "text_truncated" not in full["selection"]


def test_long_pdf_text_records_are_recoverable(monkeypatch: Any, tmp_path: Path) -> None:
    end_marker = "::EXACT_PDF_END::"
    exact_line = "needle:" + "d" * 900 + end_marker
    target = tmp_path / "sample.pdf"
    target.write_bytes(b"fake")

    class _Box:
        width = 612
        height = 792

    class _Page:
        mediabox = _Box()

    class _Reader:
        def __init__(self, _path: Path, *, strict: bool) -> None:
            assert strict is False
            self.is_encrypted = False
            self.pages = [_Page()]
            self.trailer = {"/Root": {"/Pages": {"/Count": 1}}}
            self.metadata: dict[str, str] = {}
            self.attachments: dict[str, bytes] = {}

    monkeypatch.setattr(office_inspect, "PdfReader", _Reader)
    monkeypatch.setattr(
        office_inspect,
        "extract_pdf_page_text",
        lambda *_args, **_kwargs: {
            "page": 1,
            "text": exact_line,
            "line_count": 1,
            "text_bytes": len(exact_line),
            "extraction_truncated": False,
        },
    )
    monkeypatch.setattr(
        office_inspect,
        "search_pdf_text",
        lambda *_args, **_kwargs: {
            "query": "needle",
            "match_count": 1,
            "matches": [{"page": 1, "line": 1, "text": exact_line}],
            "truncated": False,
        },
    )
    capture = _RawCapture()

    result = inspect_target(target, page=1, query="needle", raw_sink=capture)

    assert result["selection"]["text_truncated"] is True
    assert result["matches"][0]["text_truncated"] is True
    assert capture.values["sample.pdf.page-1.txt"].decode("utf-8") == exact_line
    assert json.loads(capture.values["sample.pdf.matches.json"])[0]["text"] == exact_line

    full = inspect_target(target, page=1, query="needle", full=True)
    assert full["selection"]["text"] == [exact_line]
    assert full["matches"][0]["text"] == exact_line
