from __future__ import annotations

import json
import warnings
import xml.etree.ElementTree as ET
from importlib.resources import as_file, files
from pathlib import Path

import pytest
import yaml
from docx import Document
from docx.oxml.ns import qn
from matplotlib.ft2font import FT2Font
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

import aer.artifacts.build as artifact_build
import aer.artifacts.document.builder as document_builder
import aer.artifacts.presentation.builder as presentation_builder
import aer.artifacts.workbook.builder as workbook_builder
from aer.artifacts import build_artifact
from aer.errors import AerError
from aer.hashing import sha256_file


def _spec(tmp_path: Path, name: str, value: dict[str, object]) -> Path:
    path = tmp_path / name
    path.write_text(yaml.safe_dump(value, allow_unicode=True, sort_keys=False), encoding="utf-8")
    return path


def _manifest(output: Path) -> dict[str, object]:
    return json.loads(output.with_name(output.name + ".aer.json").read_text(encoding="utf-8"))


def test_presentation_build_reopens_with_stable_ids_and_manifest(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "deck.yaml",
        {
            "version": 1,
            "kind": "presentation",
            "theme": "business-clean",
            "metadata": {"title": "AER 테스트", "locale": "ko-KR", "ratio": "16:9"},
            "footer": "Agent Efficiency Runtime",
            "content": [
                {
                    "id": "cover",
                    "layout": "title",
                    "title": "토큰 최적화",
                    "subtitle": "로컬 런타임",
                },
                {
                    "id": "metrics",
                    "layout": "metrics",
                    "title": "목표",
                    "metrics": [
                        {"id": "token-reduction", "value": "50%+", "label": "토큰 감소"},
                        {"id": "calls", "value": "3", "label": "호출"},
                    ],
                },
            ],
        },
    )
    output = tmp_path / "deck.pptx"
    result = build_artifact(spec, output)

    presentation = Presentation(output)
    assert len(presentation.slides) == 2
    assert presentation.slide_width * 9 == presentation.slide_height * 16
    assert [slide.element.cSld.get("name") for slide in presentation.slides] == [
        "aer:cover",
        "aer:metrics",
    ]
    names = {shape.name for slide in presentation.slides for shape in slide.shapes}
    assert "aer:cover/title" in names
    assert "aer:metrics/token-reduction-value" in names
    assert "aer:metrics/footer" in names
    assert "aer:metrics/slide-number" in names

    manifest = _manifest(output)
    assert manifest["kind"] == "presentation"
    assert manifest["artifact_sha256"] == sha256_file(output) == result["sha256"]
    slide_elements = manifest["elements"]
    assert isinstance(slide_elements, list)
    assert {item["id"] for item in slide_elements} == {"cover", "metrics"}


def test_document_build_reopens_with_bookmarks_table_id_and_manifest(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "report.yaml",
        {
            "version": 1,
            "kind": "document",
            "metadata": {"title": "AER 보고서", "locale": "ko-KR"},
            "header": "AER",
            "footer": "검증 문서",
            "content": [
                {"id": "title", "type": "title", "text": "Agent Efficiency Runtime"},
                {"id": "overview", "type": "heading", "level": 1, "text": "개요"},
                {"id": "body", "type": "paragraph", "text": "작은 컨텍스트로 처리합니다."},
                {
                    "id": "summary-table",
                    "type": "table",
                    "headers": ["기능", "효과"],
                    "rows": [["Patch", "부분 수정"]],
                },
            ],
        },
    )
    output = tmp_path / "report.docx"
    build_artifact(spec, output)

    document = Document(output)
    assert document.core_properties.title == "AER 보고서"
    bookmark_names = {
        bookmark.get(qn("w:name"))
        for paragraph in document.paragraphs
        for bookmark in paragraph._p.findall(qn("w:bookmarkStart"))
    }
    assert {"aer_title", "aer_overview", "aer_body"} <= bookmark_names
    caption = document.tables[0]._tbl.tblPr.find(qn("w:tblCaption"))
    assert caption is not None
    assert caption.get(qn("w:val")) == "aer:summary-table"
    assert _manifest(output)["artifact_sha256"] == sha256_file(output)


def test_document_table_preserves_cells_beyond_header_width(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "ragged-table.yaml",
        {
            "version": 1,
            "kind": "document",
            "content": [
                {
                    "id": "ragged",
                    "type": "table",
                    "headers": ["first"],
                    "rows": [["one", "must-not-be-dropped"], ["two"]],
                }
            ],
        },
    )
    output = tmp_path / "ragged-table.docx"

    build_artifact(spec, output)

    table = Document(output).tables[0]
    assert [[cell.text for cell in row.cells] for row in table.rows] == [
        ["first", ""],
        ["one", "must-not-be-dropped"],
        ["two", ""],
    ]


def test_workbook_build_reopens_with_formula_defined_name_and_manifest(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "book.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {
                    "id": "summary",
                    "name": "Summary",
                    "columns": ["workflow", "tokens"],
                    "rows": [["direct", 100], ["aer", 40]],
                    "freeze": "A2",
                    "auto_filter": True,
                    "table": {"name": "MetricsTable"},
                    "cells": [
                        {
                            "id": "ratio",
                            "address": "D2",
                            "formula": "=B3/B2",
                            "number_format": "0.0%",
                        }
                    ],
                    "charts": [
                        {
                            "type": "bar",
                            "min_col": 2,
                            "max_col": 2,
                            "categories_col": 1,
                            "min_row": 1,
                            "max_row": 3,
                        }
                    ],
                }
            ],
        },
    )
    output = tmp_path / "metrics.xlsx"
    result = build_artifact(spec, output)

    workbook = load_workbook(output, data_only=False)
    sheet = workbook["Summary"]
    assert sheet["D2"].value == "=B3/B2"
    assert sheet.freeze_panes == "A2"
    assert sheet.auto_filter.ref == "A1:D3"
    assert len(sheet.tables) == 1
    assert len(sheet._charts) == 1
    defined = workbook.defined_names["aer_summary_ratio"]
    assert defined.attr_text == "'Summary'!D2"
    assert result["warnings"][0]["code"] == "FORMULAS_NOT_CALCULATED"
    manifest = _manifest(output)
    assert manifest["artifact_sha256"] == sha256_file(output)
    assert manifest["elements"][0]["children"] == [
        {
            "id": "ratio",
            "type": "cell",
            "address": "D2",
            "selector": "sheet:id=summary/cell:id=ratio",
        }
    ]


def test_workbook_quotes_apostrophe_sheet_names_and_rejects_casefold_duplicates(
    tmp_path: Path,
) -> None:
    quoted_spec = _spec(
        tmp_path,
        "quoted-sheet.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {
                    "id": "orders",
                    "name": "O'Brien",
                    "cells": [{"id": "total", "address": "B2", "value": 7}],
                    "named_range": "order_data",
                }
            ],
        },
    )
    quoted_output = tmp_path / "quoted-sheet.xlsx"

    build_artifact(quoted_spec, quoted_output, validate=True)

    workbook = load_workbook(quoted_output, data_only=False)
    assert workbook.sheetnames == ["O'Brien"]
    assert workbook.defined_names["aer_sheet_orders"].attr_text == "'O''Brien'!$A$1"
    assert workbook.defined_names["aer_orders_total"].attr_text == "'O''Brien'!B2"
    assert workbook.defined_names["order_data"].attr_text == "'O''Brien'!A1:B2"

    duplicate_spec = _spec(
        tmp_path,
        "duplicate-sheet.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {"id": "upper", "name": "Data", "rows": [[1]]},
                {"id": "lower", "name": "data", "rows": [[2]]},
            ],
        },
    )
    duplicate_output = tmp_path / "duplicate-sheet.xlsx"

    with pytest.raises(AerError) as captured:
        build_artifact(duplicate_spec, duplicate_output)

    assert captured.value.code == "INVALID_SPEC"
    assert captured.value.target == "/sheets/1/name"
    assert not duplicate_output.exists()


@pytest.mark.parametrize(
    "cell",
    [
        {"address": "A1", "formula": "SUM(B1:B2)"},
        {"address": "A1", "formula": 7},
        {"address": "A1", "formula": "=A2", "value": 7},
    ],
)
def test_workbook_build_rejects_invalid_or_ambiguous_formula(
    tmp_path: Path, cell: dict[str, object]
) -> None:
    spec = _spec(
        tmp_path,
        "invalid-formula.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [{"id": "data", "name": "Data", "cells": [cell]}],
        },
    )
    output = tmp_path / "invalid-formula.xlsx"

    with pytest.raises(AerError) as captured:
        build_artifact(spec, output)

    assert captured.value.code == "INVALID_SPEC"
    assert not output.exists()


def test_chart_build_creates_reopenable_png_at_requested_dimensions(
    tmp_path: Path, caplog: pytest.LogCaptureFixture
) -> None:
    data = tmp_path / "metrics.csv"
    data.write_text("workflow,tokens\ndirect,100\naer,40\n", encoding="utf-8")
    spec = _spec(
        tmp_path,
        "chart.yaml",
        {
            "version": 1,
            "kind": "chart",
            "type": "bar",
            "source": "metrics.csv",
            "x": "workflow",
            "y": "tokens",
            "title": "토큰 사용량",
            "output": {"width": 800, "height": 500, "dpi": 100},
        },
    )
    output = tmp_path / "chart.png"
    font_resource = files("aer").joinpath("resources/fonts/NanumGothic.ttf")
    assert font_resource.is_file()
    with as_file(font_resource) as font_path:
        assert sha256_file(font_path) == (
            "48a28e97b34fc8e5b157657633670cd1b7de126cfc414da65ce9c3d5bc8be733"
        )
        charmap = FT2Font(str(font_path)).get_charmap()
        assert all(ord(character) in charmap for character in "토큰사용량")

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        build_artifact(spec, output)

    assert not [warning for warning in caught if "Glyph" in str(warning.message)]
    assert "findfont" not in caplog.text.casefold()

    with Image.open(output) as image:
        image.verify()
    with Image.open(output) as image:
        assert image.size == (800, 500)
        assert image.format == "PNG"
    assert _manifest(output)["elements"] == [
        {"id": "chart", "selector": "chart:id=chart", "type": "bar"}
    ]


def test_chart_build_creates_reopenable_svg(tmp_path: Path) -> None:
    data = tmp_path / "metrics.csv"
    data.write_text("workflow,tokens\ndirect,100\naer,40\n", encoding="utf-8")
    spec = _spec(
        tmp_path,
        "chart-svg.yaml",
        {
            "version": 1,
            "kind": "chart",
            "type": "bar",
            "source": "metrics.csv",
            "x": "workflow",
            "y": "tokens",
        },
    )
    output = tmp_path / "chart.svg"

    build_artifact(spec, output)

    root = ET.parse(output).getroot()
    assert root.tag == "{http://www.w3.org/2000/svg}svg"


def test_single_series_presentation_chart_builds(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "single-series.yaml",
        {
            "version": 1,
            "kind": "presentation",
            "content": [
                {
                    "id": "chart",
                    "layout": "chart",
                    "title": "One series",
                    "categories": ["A", "B"],
                    "series": [{"name": "Values", "values": [1, 2]}],
                }
            ],
        },
    )
    output = tmp_path / "single-series.pptx"

    build_artifact(spec, output)

    presentation = Presentation(output)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) >= 2


def test_presentation_scatter_chart_uses_numeric_xy_data(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "scatter.yaml",
        {
            "version": 1,
            "kind": "presentation",
            "content": [
                {
                    "id": "scatter",
                    "layout": "chart",
                    "chart_type": "scatter",
                    "categories": [1, 3],
                    "series": [{"name": "Values", "values": [2, 4]}],
                }
            ],
        },
    )
    output = tmp_path / "scatter.pptx"

    build_artifact(spec, output)

    presentation = Presentation(output)
    chart = next(shape.chart for shape in presentation.slides[0].shapes if shape.has_chart)
    assert chart.chart_type == XL_CHART_TYPE.XY_SCATTER
    assert tuple(chart.series[0].values) == (2.0, 4.0)


def test_presentation_table_preserves_cells_beyond_header_width(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "wide-table.yaml",
        {
            "version": 1,
            "kind": "presentation",
            "content": [
                {
                    "id": "wide",
                    "layout": "table",
                    "headers": ["first"],
                    "rows": [["one", "must-not-be-dropped"], ["two"]],
                }
            ],
        },
    )
    output = tmp_path / "wide-table.pptx"

    build_artifact(spec, output)

    presentation = Presentation(output)
    table = next(shape.table for shape in presentation.slides[0].shapes if shape.has_table)
    assert [[cell.text for cell in row.cells] for row in table.rows] == [
        ["first", ""],
        ["one", "must-not-be-dropped"],
        ["two", ""],
    ]


@pytest.mark.parametrize(
    ("overrides", "target"),
    [
        ({"theme": "unknown"}, "/theme"),
        ({"metadata": {"ratio": "4:3"}}, "/metadata/ratio"),
    ],
)
def test_presentation_rejects_unsupported_theme_or_ratio(
    tmp_path: Path, overrides: dict[str, object], target: str
) -> None:
    spec_value: dict[str, object] = {
        "version": 1,
        "kind": "presentation",
        "content": [{"id": "cover", "layout": "title", "title": "AER"}],
        **overrides,
    }
    spec = _spec(tmp_path, "unsupported-presentation-option.yaml", spec_value)
    output = tmp_path / "unsupported.pptx"

    with pytest.raises(AerError) as captured:
        build_artifact(spec, output)

    assert captured.value.code == "INVALID_SPEC"
    assert captured.value.target == target
    assert captured.value.details == {
        "supported": ["business-clean" if target == "/theme" else "16:9"]
    }
    assert not output.exists()


@pytest.mark.parametrize("kind", ["presentation", "workbook"])
def test_office_build_rejects_unknown_chart_type(tmp_path: Path, kind: str) -> None:
    if kind == "presentation":
        body: dict[str, object] = {
            "content": [
                {
                    "id": "chart",
                    "layout": "chart",
                    "categories": ["A"],
                    "series": [{"name": "Values", "values": [1]}],
                    "chart_type": "radar",
                }
            ]
        }
        extension = ".pptx"
    else:
        body = {
            "sheets": [
                {
                    "id": "data",
                    "columns": ["label", "value"],
                    "rows": [["A", 1]],
                    "charts": [{"type": "radar"}],
                }
            ]
        }
        extension = ".xlsx"
    spec = _spec(
        tmp_path,
        f"unknown-{kind}-chart.yaml",
        {"version": 1, "kind": kind, **body},
    )
    output = tmp_path / f"unknown-{kind}{extension}"

    with pytest.raises(AerError) as captured:
        build_artifact(spec, output)

    assert captured.value.code == "INVALID_SPEC"
    assert "chart type" in captured.value.message.casefold()
    assert not output.exists()


def test_artifact_builders_enforce_materialization_count_limits(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    monkeypatch.setattr(workbook_builder, "MAX_TABULAR_CELLS", 5)
    workbook_spec = _spec(
        tmp_path,
        "large-workbook.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [{"id": "data", "rows": [[1, 2, 3], [4, 5, 6]]}],
        },
    )
    with pytest.raises(AerError) as workbook_limit:
        build_artifact(workbook_spec, tmp_path / "large.xlsx")
    assert workbook_limit.value.code == "LIMIT_EXCEEDED"
    assert not (tmp_path / "large.xlsx").exists()

    monkeypatch.setattr(document_builder, "MAX_ARTIFACT_ELEMENTS", 2)
    document_spec = _spec(
        tmp_path,
        "large-document.yaml",
        {
            "version": 1,
            "kind": "document",
            "content": [{"id": "list", "type": "bullets", "items": ["a", "b", "c"]}],
        },
    )
    with pytest.raises(AerError) as document_limit:
        build_artifact(document_spec, tmp_path / "large.docx")
    assert document_limit.value.code == "LIMIT_EXCEEDED"
    assert not (tmp_path / "large.docx").exists()

    monkeypatch.setattr(presentation_builder, "MAX_PRESENTATION_SLIDES", 1)
    presentation_spec = _spec(
        tmp_path,
        "large-presentation.yaml",
        {
            "version": 1,
            "kind": "presentation",
            "content": [
                {"id": "one", "layout": "title", "title": "One"},
                {"id": "two", "layout": "title", "title": "Two"},
            ],
        },
    )
    with pytest.raises(AerError) as presentation_limit:
        build_artifact(presentation_spec, tmp_path / "large.pptx")
    assert presentation_limit.value.code == "LIMIT_EXCEEDED"
    assert not (tmp_path / "large.pptx").exists()


def test_presentation_enforces_aggregate_materialization_limit(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    monkeypatch.setattr(artifact_build, "MAX_PRESENTATION_ELEMENTS", 5)
    spec = _spec(
        tmp_path,
        "aggregate-presentation.yaml",
        {
            "version": 1,
            "kind": "presentation",
            "content": [
                {"id": "one", "layout": "bullets", "items": ["a", "b"]},
                {"id": "two", "layout": "bullets", "items": ["c", "d"]},
            ],
        },
    )
    output = tmp_path / "aggregate.pptx"

    with pytest.raises(AerError) as aggregate_limit:
        build_artifact(spec, output)

    assert aggregate_limit.value.code == "LIMIT_EXCEEDED"
    assert aggregate_limit.value.operation == "presentation.build"
    assert aggregate_limit.value.target == "/content"
    assert aggregate_limit.value.details == {"materialized_items": 6, "limit": 5}
    assert not output.exists()
    assert not output.with_name(output.name + ".aer.json").exists()


def test_dry_run_writes_nothing_and_invalid_spec_is_actionable(tmp_path: Path) -> None:
    valid = _spec(
        tmp_path,
        "valid.yaml",
        {
            "version": 1,
            "kind": "document",
            "content": [{"id": "body", "type": "paragraph", "text": "hello"}],
        },
    )
    output = tmp_path / "report.docx"
    result = build_artifact(valid, output, dry_run=True)
    assert result["dry_run"] is True
    assert not output.exists()
    assert not output.with_name(output.name + ".aer.json").exists()

    invalid_nested = _spec(
        tmp_path,
        "invalid-nested.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [{"id": "data", "cells": [{"value": 1}]}],
        },
    )
    with pytest.raises(AerError) as nested:
        build_artifact(invalid_nested, tmp_path / "nested.xlsx", dry_run=True)
    assert nested.value.code == "INVALID_SPEC"
    assert nested.value.target == "/address"
    assert not (tmp_path / "nested.xlsx").exists()

    invalid = _spec(tmp_path, "invalid.yaml", {"version": 2, "kind": "document"})
    with pytest.raises(AerError) as caught:
        build_artifact(invalid, tmp_path / "invalid.docx")
    assert caught.value.code == "INVALID_SPEC"
    assert caught.value.target == "/version"
    assert caught.value.suggested_action == "Set version: 1"


def test_build_validation_happens_before_publish_and_preserves_existing_files(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    spec = _spec(
        tmp_path,
        "validated.yaml",
        {
            "version": 1,
            "kind": "document",
            "content": [{"id": "body", "type": "paragraph", "text": "new content"}],
        },
    )
    output = tmp_path / "report.docx"
    manifest = output.with_name(output.name + ".aer.json")
    output.write_bytes(b"existing artifact")
    manifest.write_text("existing manifest", encoding="utf-8")

    staged_path: Path | None = None

    def reject(path: Path) -> dict[str, object]:
        nonlocal staged_path
        staged_path = path
        raise AerError(
            "VALIDATION_FAILED",
            "simulated rejection",
            "artifact.validate",
            str(path),
            {"checks": {"path": str(path)}},
        )

    monkeypatch.setattr("aer.validation.validate_file", reject)
    with pytest.raises(AerError) as rejected:
        build_artifact(spec, output, validate=True)

    assert rejected.value.code == "VALIDATION_FAILED"
    assert rejected.value.target == str(output.resolve())
    assert rejected.value.details["checks"]["path"] == str(output.resolve())
    assert staged_path is not None
    assert not staged_path.exists()
    assert ".aer-build-" not in str(rejected.value.target)
    assert output.read_bytes() == b"existing artifact"
    assert manifest.read_text(encoding="utf-8") == "existing manifest"


def test_successful_build_validation_reports_published_path(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "validated-success.yaml",
        {
            "version": 1,
            "kind": "document",
            "content": [{"id": "body", "type": "paragraph", "text": "valid content"}],
        },
    )
    output = tmp_path / "published.docx"

    result = build_artifact(spec, output, validate=True)

    assert result["validation"]["valid"] is True
    assert result["validation"]["checks"]["path"] == str(output.resolve())


def test_cyclic_yaml_spec_is_rejected_without_partial_artifact(tmp_path: Path) -> None:
    spec = tmp_path / "cycle.yaml"
    spec.write_text(
        "version: 1\n"
        "kind: presentation\n"
        "content:\n"
        "  - &slide\n"
        "    id: cycle\n"
        "    layout: bullets\n"
        "    title: Cycle\n"
        "    items: [*slide]\n",
        encoding="utf-8",
    )
    output = tmp_path / "cycle.pptx"

    with pytest.raises(AerError, match="Cyclic YAML aliases") as captured:
        build_artifact(spec, output)

    assert captured.value.code == "INVALID_SPEC"
    assert not output.exists()
    assert not output.with_name(output.name + ".aer.json").exists()


def test_html_build_escapes_semantic_content(tmp_path: Path) -> None:
    spec = _spec(
        tmp_path,
        "safe-html.yaml",
        {
            "version": 1,
            "kind": "html",
            "metadata": {"title": "<unsafe>"},
            "content": [{"type": "paragraph", "text": "<script>alert(1)</script>"}],
        },
    )
    output = tmp_path / "safe.html"

    build_artifact(spec, output)

    rendered = output.read_text(encoding="utf-8")
    assert "<script>" not in rendered
    assert "&lt;script&gt;alert(1)&lt;/script&gt;" in rendered
    assert "<title>&lt;unsafe&gt;</title>" in rendered


@pytest.mark.parametrize(
    ("name", "extension", "spec"),
    [
        (
            "duplicate-sheets",
            ".xlsx",
            {
                "version": 1,
                "kind": "workbook",
                "sheets": [
                    {"id": "duplicate", "name": "First", "rows": [[1]]},
                    {"id": "duplicate", "name": "Second", "rows": [[2]]},
                ],
            },
        ),
        (
            "normalized-cells",
            ".xlsx",
            {
                "version": 1,
                "kind": "workbook",
                "sheets": [
                    {
                        "id": "summary",
                        "cells": [
                            {"id": "token-value", "address": "A1", "value": 1},
                            {"id": "token_value", "address": "A2", "value": 2},
                        ],
                    }
                ],
            },
        ),
        (
            "duplicate-metrics",
            ".pptx",
            {
                "version": 1,
                "kind": "presentation",
                "content": [
                    {
                        "id": "metrics",
                        "layout": "metrics",
                        "title": "Metrics",
                        "metrics": [
                            {"id": "duplicate", "value": "1", "label": "First"},
                            {"id": "duplicate", "value": "2", "label": "Second"},
                        ],
                    }
                ],
            },
        ),
        (
            "normalized-bookmarks",
            ".docx",
            {
                "version": 1,
                "kind": "document",
                "content": [
                    {"id": "same-id", "type": "paragraph", "text": "First"},
                    {"id": "same_id", "type": "paragraph", "text": "Second"},
                ],
            },
        ),
    ],
)
def test_build_rejects_ambiguous_stable_ids(
    tmp_path: Path, name: str, extension: str, spec: dict[str, object]
) -> None:
    source = _spec(tmp_path, f"{name}.yaml", spec)
    output = tmp_path / f"{name}{extension}"

    with pytest.raises(AerError) as caught:
        build_artifact(source, output)

    assert caught.value.code == "INVALID_SPEC"
    assert not output.exists()
