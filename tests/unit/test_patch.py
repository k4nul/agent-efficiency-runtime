from __future__ import annotations

import json
import os
import stat
from pathlib import Path

import pytest
import yaml
from docx import Document
from docx.oxml.ns import qn
from openpyxl import load_workbook
from openpyxl.workbook.defined_name import DefinedName
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

import aer.patch.engine as patch_engine
from aer.artifacts import build_artifact
from aer.config import Settings
from aer.errors import AerError
from aer.hashing import normalized_hash, sha256_file
from aer.patch import apply_patch
from aer.validation import validate_file


@pytest.fixture(autouse=True)
def _isolated_aer_home(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    monkeypatch.setenv("AER_HOME", str(tmp_path / "aer-home"))


def _write_yaml(path: Path, value: dict[str, object]) -> Path:
    path.write_text(yaml.safe_dump(value, allow_unicode=True, sort_keys=False), encoding="utf-8")
    return path


def _build_deck(tmp_path: Path) -> Path:
    spec = _write_yaml(
        tmp_path / "deck.yaml",
        {
            "version": 1,
            "kind": "presentation",
            "content": [
                {"id": "cover", "layout": "title", "title": "Keep me", "subtitle": "Stable"},
                {
                    "id": "metrics",
                    "layout": "metrics",
                    "title": "Metrics",
                    "metrics": [{"id": "token", "value": "50%", "label": "Reduction"}],
                },
            ],
        },
    )
    output = tmp_path / "deck.pptx"
    build_artifact(spec, output)
    return output


def _shape_snapshot(path: Path) -> dict[str, tuple[str, int, int, int, int]]:
    presentation = Presentation(path)
    return {
        shape.name: (getattr(shape, "text", ""), shape.left, shape.top, shape.width, shape.height)
        for slide in presentation.slides
        for shape in slide.shapes
    }


def test_cyclic_yaml_patch_is_rejected_without_mutating_target(tmp_path: Path) -> None:
    target = tmp_path / "target.json"
    target.write_text('{"status":"old"}\n', encoding="utf-8")
    before = target.read_bytes()
    patch = tmp_path / "cycle.yaml"
    patch.write_text(
        "version: 1\n"
        "operations:\n"
        "  - &operation\n"
        "    op: json.set\n"
        "    target: /status\n"
        "    value: new\n"
        "    self: *operation\n",
        encoding="utf-8",
    )

    with pytest.raises(AerError, match="Cyclic YAML aliases") as captured:
        apply_patch(target, patch)

    assert captured.value.code == "INVALID_PATCH"
    assert target.read_bytes() == before


def test_pptx_patch_changes_only_selected_shape_and_updates_manifest(tmp_path: Path) -> None:
    target = _build_deck(tmp_path)
    before = _shape_snapshot(target)
    patch = _write_yaml(
        tmp_path / "patch.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "pptx.set_text",
                    "target": "slide:id=metrics/shape:id=token-value",
                    "value": "63%",
                }
            ],
        },
    )
    result = apply_patch(target, patch)
    after = _shape_snapshot(target)

    assert after["aer:metrics/token-value"][0] == "63%"
    assert {name: value for name, value in after.items() if name != "aer:metrics/token-value"} == {
        name: value for name, value in before.items() if name != "aer:metrics/token-value"
    }
    manifest = json.loads(target.with_name(target.name + ".aer.json").read_text(encoding="utf-8"))
    assert manifest["artifact_sha256"] == result["after_sha256"] == sha256_file(target)


def test_pptx_set_text_preserves_first_run_formatting_and_collapses_split_runs(
    tmp_path: Path,
) -> None:
    target = _build_deck(tmp_path)
    presentation = Presentation(target)
    shape = next(
        shape for shape in presentation.slides[1].shapes if shape.name == "aer:metrics/token-value"
    )
    paragraph = shape.text_frame.paragraphs[0]
    first = paragraph.runs[0]
    first.text = "5"
    first.font.name = "Aptos Display"
    first.font.size = Pt(31)
    first.font.bold = True
    first.font.italic = True
    first.font.color.rgb = RGBColor(12, 34, 56)
    paragraph.add_run().text = "0%"
    presentation.save(target)
    patch = _write_yaml(
        tmp_path / "set-text.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "pptx.set_text",
                    "target": "slide:id=metrics/shape:id=token-value",
                    "value": "63%",
                }
            ],
        },
    )

    apply_patch(target, patch)

    reopened = Presentation(target)
    patched = next(
        shape for shape in reopened.slides[1].shapes if shape.name == "aer:metrics/token-value"
    )
    assert patched.text == "63%"
    assert len(patched.text_frame.paragraphs) == 1
    runs = patched.text_frame.paragraphs[0].runs
    assert len(runs) == 1
    assert runs[0].text == "63%"
    assert runs[0].font.name == "Aptos Display"
    assert runs[0].font.size.pt == 31
    assert runs[0].font.bold is True
    assert runs[0].font.italic is True
    assert runs[0].font.color.rgb == RGBColor(12, 34, 56)


def test_pptx_replace_text_handles_text_split_across_runs(tmp_path: Path) -> None:
    target = _build_deck(tmp_path)
    presentation = Presentation(target)
    shape = next(
        shape for shape in presentation.slides[0].shapes if shape.name == "aer:cover/title"
    )
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.runs[0].text = "Split"
    paragraph.add_run().text = "AcrossRuns"
    presentation.save(target)
    patch = _write_yaml(
        tmp_path / "split.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "pptx.replace_text",
                    "target": "slide:id=cover/shape:id=title",
                    "old": "SplitAcrossRuns",
                    "value": "Combined",
                }
            ],
        },
    )
    apply_patch(target, patch)
    assert _shape_snapshot(target)["aer:cover/title"][0] == "Combined"


def test_docx_replace_text_handles_run_split_and_preserves_other_blocks(tmp_path: Path) -> None:
    spec = _write_yaml(
        tmp_path / "document.yaml",
        {
            "version": 1,
            "kind": "document",
            "content": [
                {"id": "target", "type": "paragraph", "text": "SplitAcrossRuns"},
                {"id": "untouched", "type": "paragraph", "text": "Do not change"},
            ],
        },
    )
    target = tmp_path / "document.docx"
    build_artifact(spec, target)
    document = Document(target)
    paragraph = document.paragraphs[0]
    paragraph.runs[0].text = "Split"
    paragraph.add_run("AcrossRuns")
    document.save(target)
    patch = _write_yaml(
        tmp_path / "docx-patch.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "docx.replace_text",
                    "target": "block:id=target",
                    "old": "SplitAcrossRuns",
                    "value": "Combined",
                }
            ],
        },
    )
    apply_patch(target, patch)
    reopened = Document(target)
    assert reopened.paragraphs[0].text == "Combined"
    assert reopened.paragraphs[1].text == "Do not change"


def test_docx_replace_text_preserves_unaffected_run_formatting(tmp_path: Path) -> None:
    document = Document()
    paragraph = document.add_paragraph()
    first = paragraph.add_run("Hello ")
    first.bold = True
    second = paragraph.add_run("world")
    second.italic = True
    target = tmp_path / "formatted.docx"
    document.save(target)
    patch = _write_yaml(
        tmp_path / "format-patch.yaml",
        {
            "version": 1,
            "operations": [{"op": "docx.replace_text", "old": "world", "value": "earth"}],
        },
    )

    apply_patch(target, patch)

    reopened = Document(target)
    runs = reopened.paragraphs[0].runs
    assert [(run.text, run.bold, run.italic) for run in runs] == [
        ("Hello ", True, None),
        ("earth", None, True),
    ]


@pytest.mark.parametrize("kind", ["bullets", "numbered-list", "source-list"])
def test_docx_set_block_replaces_entire_multi_paragraph_logical_block(
    tmp_path: Path, kind: str
) -> None:
    spec = _write_yaml(
        tmp_path / f"{kind}.yaml",
        {
            "version": 1,
            "kind": "document",
            "content": [
                {"id": "before", "type": "paragraph", "text": "Keep before"},
                {"id": "logical-list", "type": kind, "items": ["one", "two", "three"]},
                {"id": "after", "type": "paragraph", "text": "Keep after"},
            ],
        },
    )
    target = tmp_path / f"{kind}.docx"
    build_artifact(spec, target)

    built = Document(target)
    starts = [
        paragraph_index
        for paragraph_index, paragraph in enumerate(built.paragraphs)
        for bookmark in paragraph._p.findall(qn("w:bookmarkStart"))
        if bookmark.get(qn("w:name")) == "aer_logical_list"
    ]
    ends = [
        paragraph_index
        for paragraph_index, paragraph in enumerate(built.paragraphs)
        for bookmark in paragraph._p.findall(qn("w:bookmarkEnd"))
        if bookmark.get(qn("w:id")) == "2"
    ]
    assert starts == [1]
    assert ends == [3]
    built.paragraphs[0].runs[0].bold = True
    built.paragraphs[-1].runs[0].italic = True
    built.save(target)

    patch = _write_yaml(
        tmp_path / f"set-{kind}.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "docx.set_block",
                    "target": "block:id=logical-list",
                    "value": "Replacement",
                }
            ],
        },
    )
    result = apply_patch(target, patch, validate=True)

    reopened = Document(target)
    assert [paragraph.text for paragraph in reopened.paragraphs] == [
        "Keep before",
        "Replacement",
        "Keep after",
    ]
    assert reopened.paragraphs[0].runs[0].bold is True
    assert reopened.paragraphs[-1].runs[0].italic is True
    replacement = reopened.paragraphs[1]
    assert any(
        bookmark.get(qn("w:name")) == "aer_logical_list"
        for bookmark in replacement._p.findall(qn("w:bookmarkStart"))
    )
    assert any(
        bookmark.get(qn("w:id")) == "2" for bookmark in replacement._p.findall(qn("w:bookmarkEnd"))
    )
    assert result["validation"]["valid"] is True
    assert result["validation"]["checks"]["path"] == str(target)
    assert Path(result["validation"]["checks"]["path"]).is_file()
    assert validate_file(target)["valid"] is True


@pytest.mark.parametrize("kind", ["bullets", "numbered-list", "source-list"])
def test_docx_remove_block_removes_all_parts_and_manifest_selector(
    tmp_path: Path, kind: str
) -> None:
    spec = _write_yaml(
        tmp_path / f"remove-{kind}.yaml",
        {
            "version": 1,
            "kind": "document",
            "content": [
                {"id": "before", "type": "paragraph", "text": "Keep before"},
                {"id": "logical-list", "type": kind, "items": ["one", "two", "three"]},
                {"id": "after", "type": "paragraph", "text": "Keep after"},
            ],
        },
    )
    target = tmp_path / f"remove-{kind}.docx"
    build_artifact(spec, target)
    patch = _write_yaml(
        tmp_path / f"remove-{kind}-patch.yaml",
        {
            "version": 1,
            "operations": [{"op": "docx.remove_block", "target": "block:id=logical-list"}],
        },
    )

    result = apply_patch(target, patch, validate=True)

    reopened = Document(target)
    assert [paragraph.text for paragraph in reopened.paragraphs] == [
        "Keep before",
        "Keep after",
    ]
    manifest_path = target.with_name(target.name + ".aer.json")
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert [element["selector"] for element in manifest["elements"]] == [
        "block:id=before",
        "block:id=after",
    ]
    assert manifest["artifact_sha256"] == result["after_sha256"] == sha256_file(target)
    assert result["validation"]["valid"] is True
    assert validate_file(target)["valid"] is True


def test_pptx_remove_shape_prunes_manifest_and_preserves_other_shapes(tmp_path: Path) -> None:
    target = _build_deck(tmp_path)
    before = _shape_snapshot(target)
    patch = _write_yaml(
        tmp_path / "remove-shape.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "pptx.remove_shape",
                    "target": "slide:id=metrics/shape:id=token-value",
                }
            ],
        },
    )

    result = apply_patch(target, patch, validate=True)

    after = _shape_snapshot(target)
    assert "aer:metrics/token-value" not in after
    assert after == {
        name: value for name, value in before.items() if name != "aer:metrics/token-value"
    }
    manifest = json.loads(target.with_name(target.name + ".aer.json").read_text(encoding="utf-8"))
    selectors = {
        child["selector"] for slide in manifest["elements"] for child in slide.get("children", [])
    }
    assert "slide:id=metrics/shape:id=token-value" not in selectors
    assert manifest["artifact_sha256"] == result["after_sha256"] == sha256_file(target)
    assert result["validation"]["valid"] is True
    assert validate_file(target)["valid"] is True


def test_manifest_write_failure_rolls_back_artifact(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    target = _build_deck(tmp_path)
    artifact_before = target.read_bytes()
    sidecar = target.with_name(target.name + ".aer.json")
    manifest_before = sidecar.read_bytes()
    patch = _write_yaml(
        tmp_path / "rollback.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "pptx.remove_shape",
                    "target": "slide:id=metrics/shape:id=token-value",
                }
            ],
        },
    )
    real_write = patch_engine.atomic_write_bytes
    failed = False

    def fail_first_manifest_write(destination: Path, data: bytes) -> None:
        nonlocal failed
        if destination == sidecar and not failed:
            failed = True
            raise OSError("simulated manifest write failure")
        real_write(destination, data)

    monkeypatch.setattr(patch_engine, "atomic_write_bytes", fail_first_manifest_write)

    with pytest.raises(OSError, match="simulated manifest write failure"):
        apply_patch(target, patch)

    assert target.read_bytes() == artifact_before
    assert sidecar.read_bytes() == manifest_before


def test_xlsx_cell_patch_preserves_formula_style_and_other_cells(tmp_path: Path) -> None:
    spec = _write_yaml(
        tmp_path / "workbook.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {
                    "id": "summary",
                    "name": "Summary",
                    "columns": ["name", "value"],
                    "rows": [["before", 10], ["other", 20]],
                    "cells": [
                        {
                            "id": "ratio",
                            "address": "D2",
                            "formula": "=B3/B2",
                            "number_format": "0.0%",
                        }
                    ],
                }
            ],
        },
    )
    target = tmp_path / "book.xlsx"
    build_artifact(spec, target)
    before = load_workbook(target, data_only=False)
    formula, number_format, other = (
        before["Summary"]["D2"].value,
        before["Summary"]["D2"].number_format,
        before["Summary"]["A3"].value,
    )
    patch = _write_yaml(
        tmp_path / "xlsx-patch.yaml",
        {
            "version": 1,
            "operations": [
                {"op": "xlsx.set_cell", "target": "sheet:id=summary/cell=B2", "value": 15}
            ],
        },
    )
    apply_patch(target, patch)
    after = load_workbook(target, data_only=False)
    assert after["Summary"]["B2"].value == 15
    assert after["Summary"]["D2"].value == formula == "=B3/B2"
    assert after["Summary"]["D2"].number_format == number_format == "0.0%"
    assert after["Summary"]["A3"].value == other == "other"


def test_xlsx_set_cell_accepts_formula_and_preserves_cell_style(tmp_path: Path) -> None:
    spec = _write_yaml(
        tmp_path / "formula-workbook.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {
                    "id": "summary",
                    "name": "Summary",
                    "columns": ["value"],
                    "rows": [[10]],
                    "cells": [
                        {
                            "id": "ratio",
                            "address": "B2",
                            "value": 1,
                            "number_format": "0.0%",
                        }
                    ],
                }
            ],
        },
    )
    target = tmp_path / "formula.xlsx"
    build_artifact(spec, target)
    patch = _write_yaml(
        tmp_path / "formula-patch.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "xlsx.set_cell",
                    "target": "sheet:id=summary/cell:id=ratio",
                    "formula": "=A2/100",
                }
            ],
        },
    )

    apply_patch(target, patch)

    workbook = load_workbook(target, data_only=False)
    try:
        assert workbook["Summary"]["B2"].value == "=A2/100"
        assert workbook["Summary"]["B2"].number_format == "0.0%"
    finally:
        workbook.close()


@pytest.mark.parametrize(
    "fields",
    [
        {},
        {"value": 1, "formula": "=A1"},
        {"formula": "A1"},
        {"formula": 1},
    ],
)
def test_xlsx_set_cell_rejects_ambiguous_or_invalid_values(
    tmp_path: Path, fields: dict[str, object]
) -> None:
    spec = _write_yaml(
        tmp_path / "invalid-formula-workbook.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [{"id": "data", "name": "Data", "columns": ["value"], "rows": [[1]]}],
        },
    )
    target = tmp_path / "invalid-formula.xlsx"
    build_artifact(spec, target)
    before = target.read_bytes()
    patch = _write_yaml(
        tmp_path / "invalid-formula-patch.yaml",
        {
            "version": 1,
            "operations": [{"op": "xlsx.set_cell", "target": "sheet:id=data/cell=A2", **fields}],
        },
    )

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "INVALID_PATCH"
    assert target.read_bytes() == before


def test_xlsx_stable_sheet_id_is_independent_from_display_name(tmp_path: Path) -> None:
    spec = _write_yaml(
        tmp_path / "stable-workbook.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {
                    "id": "financial-summary",
                    "name": "KPI Dashboard",
                    "columns": ["metric", "value"],
                    "rows": [["tokens", 100]],
                    "cells": [{"id": "token-value", "address": "B2", "value": 100}],
                }
            ],
        },
    )
    target = tmp_path / "stable.xlsx"
    build_artifact(spec, target)
    patch = _write_yaml(
        tmp_path / "stable-patch.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "xlsx.set_cell",
                    "target": "sheet:id=financial-summary/cell:id=token-value",
                    "value": 63,
                }
            ],
        },
    )

    apply_patch(target, patch)

    workbook = load_workbook(target, data_only=False)
    try:
        assert workbook["KPI Dashboard"]["B2"].value == 63
        assert "aer_sheet_financial_summary" in workbook.defined_names
        assert "aer_financial_summary_token_value" in workbook.defined_names
    finally:
        workbook.close()


def test_xlsx_stable_cell_selector_handles_apostrophe_sheet_name(tmp_path: Path) -> None:
    spec = _write_yaml(
        tmp_path / "apostrophe-workbook.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {
                    "id": "orders",
                    "name": "O'Brien",
                    "cells": [{"id": "total", "address": "B2", "value": 7}],
                }
            ],
        },
    )
    target = tmp_path / "apostrophe.xlsx"
    build_artifact(spec, target)
    patch = _write_yaml(
        tmp_path / "apostrophe-patch.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "xlsx.set_cell",
                    "target": "sheet:id=orders/cell:id=total",
                    "value": 11,
                }
            ],
        },
    )

    apply_patch(target, patch, validate=True)

    workbook = load_workbook(target, data_only=False)
    try:
        assert workbook["O'Brien"]["B2"].value == 11
    finally:
        workbook.close()


def test_xlsx_stable_sheet_id_wins_over_conflicting_display_name(tmp_path: Path) -> None:
    spec = _write_yaml(
        tmp_path / "conflicting-workbook.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {
                    "id": "stable",
                    "name": "Dashboard",
                    "columns": ["value"],
                    "rows": [[1]],
                },
                {
                    "id": "other",
                    "name": "stable",
                    "columns": ["value"],
                    "rows": [[2]],
                },
            ],
        },
    )
    target = tmp_path / "conflicting.xlsx"
    build_artifact(spec, target)
    patch = _write_yaml(
        tmp_path / "conflicting-patch.yaml",
        {
            "version": 1,
            "operations": [
                {"op": "xlsx.set_cell", "target": "sheet:id=stable/cell=A2", "value": 99}
            ],
        },
    )

    apply_patch(target, patch)

    workbook = load_workbook(target, data_only=False)
    try:
        assert workbook["Dashboard"]["A2"].value == 99
        assert workbook["stable"]["A2"].value == 2
    finally:
        workbook.close()


def test_xlsx_stable_cell_id_rejects_multiple_destinations_without_mutation(
    tmp_path: Path,
) -> None:
    spec = _write_yaml(
        tmp_path / "ambiguous-cell.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [
                {
                    "id": "dashboard",
                    "name": "Dashboard",
                    "cells": [{"id": "metric", "address": "A2", "value": 1}],
                }
            ],
        },
    )
    target = tmp_path / "ambiguous-cell.xlsx"
    build_artifact(spec, target)
    workbook = load_workbook(target)
    del workbook.defined_names["aer_dashboard_metric"]
    workbook.defined_names.add(
        DefinedName(
            "aer_dashboard_metric",
            attr_text="'Dashboard'!$A$2,'Dashboard'!$B$2",
        )
    )
    workbook.save(target)
    workbook.close()
    before = target.read_bytes()
    patch = _write_yaml(
        tmp_path / "ambiguous-cell-patch.yaml",
        {
            "version": 1,
            "operations": [
                {
                    "op": "xlsx.set_cell",
                    "target": "sheet:id=dashboard/cell:id=metric",
                    "value": 99,
                }
            ],
        },
    )

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "INVALID_SELECTOR"
    assert target.read_bytes() == before


@pytest.mark.parametrize("op", ["xlsx.set_range", "xlsx.replace_text", "xlsx.clear_range"])
def test_xlsx_patch_rejects_oversized_range_before_materializing_cells(
    tmp_path: Path, op: str
) -> None:
    spec = _write_yaml(
        tmp_path / "bounded-workbook.yaml",
        {
            "version": 1,
            "kind": "workbook",
            "sheets": [{"id": "data", "name": "Data", "columns": ["value"], "rows": [[1]]}],
        },
    )
    target = tmp_path / "bounded.xlsx"
    build_artifact(spec, target)
    before = target.read_bytes()
    operation: dict[str, object] = {
        "op": op,
        "target": "Data!A1:XFD1048576",
    }
    if op == "xlsx.set_range":
        operation["values"] = [["never materialized"]]
    elif op == "xlsx.replace_text":
        operation.update({"old": "value", "value": "other"})
    patch = _write_yaml(
        tmp_path / f"{op}.yaml",
        {"version": 1, "operations": [operation]},
    )

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert captured.value.target == "A1:XFD1048576"
    assert captured.value.details == {
        "cells": 16_384 * 1_048_576,
        "limit": 1_000_000,
    }
    assert target.read_bytes() == before


def test_multi_operation_failure_is_atomic(tmp_path: Path) -> None:
    target = tmp_path / "notes.txt"
    target.write_text("alpha beta", encoding="utf-8")
    before = target.read_bytes()
    patch = _write_yaml(
        tmp_path / "atomic.yaml",
        {
            "version": 1,
            "operations": [
                {"op": "text.replace", "old": "alpha", "value": "changed"},
                {"op": "text.replace", "old": "missing", "value": "never"},
            ],
        },
    )
    with pytest.raises(AerError) as caught:
        apply_patch(target, patch)
    assert caught.value.code == "INVALID_SELECTOR"
    assert target.read_bytes() == before


def test_stale_hash_and_dry_run_leave_original_unchanged(tmp_path: Path) -> None:
    target = tmp_path / "notes.txt"
    target.write_text("before", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "patch.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.replace", "old": "before", "value": "after"}],
        },
    )
    with pytest.raises(AerError) as stale:
        apply_patch(target, patch, expected_sha256="0" * 64)
    assert stale.value.code == "HASH_MISMATCH"
    assert target.read_text() == "before"

    result = apply_patch(target, patch, dry_run=True, expected_sha256=sha256_file(target))
    assert result["dry_run"] is True
    assert target.read_text() == "before"


def test_concurrent_change_during_patch_is_detected(monkeypatch, tmp_path: Path) -> None:
    target = tmp_path / "notes.txt"
    target.write_text("before", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "race.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.replace", "old": "before", "value": "after"}],
        },
    )
    real_selector = patch_engine._select_patcher

    def racing_selector(suffix, operations):
        patcher = real_selector(suffix, operations)

        def race(data, selected_operations):
            changed = patcher(data, selected_operations)
            target.write_text("external-change", encoding="utf-8")
            return changed

        return race

    monkeypatch.setattr(patch_engine, "_select_patcher", racing_selector)

    with pytest.raises(AerError) as conflict:
        apply_patch(target, patch, expected_sha256=sha256_file(target))

    assert conflict.value.code == "HASH_MISMATCH"
    assert target.read_text(encoding="utf-8") == "external-change"


def test_validation_failure_leaves_patch_target_unchanged(tmp_path: Path) -> None:
    target = tmp_path / "notes.unsupported"
    target.write_text("before", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "validate.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.replace", "old": "before", "value": "after"}],
        },
    )

    with pytest.raises(AerError) as invalid:
        apply_patch(target, patch, validate=True)

    assert invalid.value.code == "UNSUPPORTED_FORMAT"
    assert invalid.value.target == str(target)
    assert target.read_text(encoding="utf-8") == "before"
    assert list(tmp_path.glob(".notes.validate-*")) == []


@pytest.mark.parametrize("suffix", [".json", ".yaml"])
def test_structured_set_insert_remove_operations(tmp_path: Path, suffix: str) -> None:
    target = tmp_path / f"data{suffix}"
    original = {"items": [{"name": "first"}], "obsolete": True}
    if suffix == ".json":
        target.write_text(json.dumps(original), encoding="utf-8")
        prefix = "json"
    else:
        target.write_text(yaml.safe_dump(original), encoding="utf-8")
        prefix = "yaml"
    patch = _write_yaml(
        tmp_path / f"structured-{prefix}.yaml",
        {
            "version": 1,
            "operations": [
                {"op": f"{prefix}.set", "target": "/items/0/name", "value": "updated"},
                {"op": f"{prefix}.insert", "target": "/items/-", "value": {"name": "second"}},
                {"op": f"{prefix}.remove", "target": "/obsolete"},
            ],
        },
    )
    apply_patch(target, patch)
    result = (
        json.loads(target.read_text()) if suffix == ".json" else yaml.safe_load(target.read_text())
    )
    assert result == {"items": [{"name": "updated"}, {"name": "second"}]}


@pytest.mark.parametrize(("suffix", "prefix"), [(".json", "json"), (".yaml", "yaml")])
def test_structured_patch_rejects_negative_array_index_without_mutation(
    tmp_path: Path, suffix: str, prefix: str
) -> None:
    target = tmp_path / f"array{suffix}"
    original = [1, 2, 3]
    if suffix == ".json":
        target.write_text(json.dumps(original), encoding="utf-8")
    else:
        target.write_text(yaml.safe_dump(original), encoding="utf-8")
    before = target.read_bytes()
    patch = _write_yaml(
        tmp_path / f"negative-{prefix}.yaml",
        {
            "version": 1,
            "operations": [{"op": f"{prefix}.set", "target": "/-1", "value": 9}],
        },
    )

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "INVALID_SELECTOR"
    assert target.read_bytes() == before


def test_patch_rejects_more_than_one_thousand_operations_without_mutation(tmp_path: Path) -> None:
    target = tmp_path / "bounded.json"
    target.write_text('{"value":0}\n', encoding="utf-8")
    before = target.read_bytes()
    patch = _write_yaml(
        tmp_path / "too-many-operations.yaml",
        {
            "version": 1,
            "operations": [
                {"op": "json.set", "target": "/value", "value": index} for index in range(1001)
            ],
        },
    )

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert captured.value.target == "/operations"
    assert captured.value.details == {"operations": 1001, "limit": 1000}
    assert target.read_bytes() == before


def test_nested_unbounded_regex_is_rejected_without_touching_file(tmp_path: Path) -> None:
    target = tmp_path / "text.txt"
    target.write_text("a" * 100, encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "regex.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.regex_replace", "pattern": "(a+)+$", "value": "x"}],
        },
    )
    with pytest.raises(AerError) as caught:
        apply_patch(target, patch)
    assert caught.value.code == "INVALID_PATCH"
    assert target.read_text() == "a" * 100


@pytest.mark.skipif(os.name != "posix", reason="POSIX permission bits")
def test_patch_initializes_private_runtime_directories(tmp_path: Path) -> None:
    target = tmp_path / "private-runtime.txt"
    target.write_text("before", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "private-runtime.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.replace", "old": "before", "value": "after"}],
        },
    )

    apply_patch(target, patch, dry_run=True)

    home = Path(os.environ["AER_HOME"])
    assert stat.S_IMODE(home.stat().st_mode) == 0o700
    assert stat.S_IMODE((home / "cache").stat().st_mode) == 0o700
    assert stat.S_IMODE((home / "cache" / "patch-locks").stat().st_mode) == 0o700
    lock_files = list((home / "cache" / "patch-locks").glob("*.lock"))
    assert len(lock_files) == 1
    assert stat.S_IMODE(lock_files[0].stat().st_mode) == 0o600


@pytest.mark.skipif(os.name != "posix", reason="POSIX symlink behavior")
def test_patch_rejects_cache_symlink_without_writing_outside(tmp_path: Path) -> None:
    home = Path(os.environ["AER_HOME"])
    home.mkdir(mode=0o700)
    if os.name == "posix":
        home.chmod(0o700)
    outside = tmp_path / "outside-cache"
    outside.mkdir()
    (home / "cache").symlink_to(outside, target_is_directory=True)
    target = tmp_path / "symlink.txt"
    target.write_text("before", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "symlink.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.replace", "old": "before", "value": "after"}],
        },
    )

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "INVALID_ARGUMENT"
    assert captured.value.target == str(home / "cache")
    assert list(outside.iterdir()) == []
    assert target.read_text(encoding="utf-8") == "before"


@pytest.mark.skipif(os.name != "posix", reason="POSIX permission bits")
def test_patch_preserves_legacy_lock_mode_beneath_private_cache(tmp_path: Path) -> None:
    settings = Settings.load()
    settings.ensure()
    lock_directory = settings.cache_dir / "patch-locks"
    lock_directory.mkdir(mode=0o770)
    lock_directory.chmod(0o770)
    target = tmp_path / "legacy-lock.txt"
    target.write_text("before", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "legacy-lock.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.replace", "old": "before", "value": "after"}],
        },
    )

    result = apply_patch(target, patch, dry_run=True)

    assert result["dry_run"] is True
    assert stat.S_IMODE(lock_directory.stat().st_mode) == 0o770


@pytest.mark.skipif(os.name != "posix", reason="POSIX permission bits")
def test_patch_rejects_shared_lock_reachable_through_cache(tmp_path: Path) -> None:
    settings = Settings.load()
    settings.ensure()
    settings.cache_dir.chmod(0o750)
    lock_directory = settings.cache_dir / "patch-locks"
    lock_directory.mkdir(mode=0o770)
    lock_directory.chmod(0o770)
    target = tmp_path / "reachable-lock.txt"
    target.write_text("before", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "reachable-lock.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.replace", "old": "before", "value": "after"}],
        },
    )

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "INVALID_ARGUMENT"
    assert captured.value.target == str(lock_directory)
    assert target.read_text(encoding="utf-8") == "before"


@pytest.mark.skipif(os.name != "posix", reason="POSIX symlink behavior")
def test_patch_rejects_hashed_lock_symlink_in_legacy_directory(tmp_path: Path) -> None:
    settings = Settings.load()
    settings.ensure()
    lock_directory = settings.cache_dir / "patch-locks"
    lock_directory.mkdir(mode=0o770)
    lock_directory.chmod(0o770)
    target = (tmp_path / "hashed-lock.txt").resolve()
    target.write_text("before", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "hashed-lock.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.replace", "old": "before", "value": "after"}],
        },
    )
    outside = tmp_path / "outside-lock-target"
    outside.write_bytes(b"unchanged")
    lock_path = lock_directory / f"{normalized_hash(str(target))}.lock"
    lock_path.symlink_to(outside)

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "INVALID_ARGUMENT"
    assert captured.value.target == str(lock_path)
    assert outside.read_bytes() == b"unchanged"
    assert target.read_text(encoding="utf-8") == "before"


def test_ambiguous_regex_is_bounded_by_timeout(tmp_path: Path) -> None:
    target = tmp_path / "ambiguous.txt"
    target.write_text("a" * 100_000 + "!", encoding="utf-8")
    patch = _write_yaml(
        tmp_path / "ambiguous.yaml",
        {
            "version": 1,
            "operations": [{"op": "text.regex_replace", "pattern": "(a|aa)+$", "value": "x"}],
        },
    )

    with pytest.raises(AerError) as captured:
        apply_patch(target, patch)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert target.read_text(encoding="utf-8").endswith("!")
