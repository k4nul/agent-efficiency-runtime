from __future__ import annotations

import json
import shutil
import subprocess
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from openpyxl.workbook.defined_name import DefinedName
from pptx import Presentation
from pypdf import PdfWriter

import aer.inspect.repository as repository_inspect
import aer.inspect.structured as structured_inspect
import aer.inspect.text as text_inspect
import aer.pdf.safety as pdf_safety
from aer.artifacts import build_artifact
from aer.errors import AerError
from aer.inspect import Inspector, inspect_target
from aer.validation import validate_file


def test_text_query_context_and_overflow_reference(tmp_path: Path) -> None:
    target = tmp_path / "app.log"
    target.write_text("\n".join(f"ERROR item {number}" for number in range(30)), encoding="utf-8")
    stored: list[bytes] = []

    def sink(data: bytes, name: str) -> str:
        assert name.endswith(".matches.json")
        stored.append(data)
        return "aer://sha256/" + "a" * 64

    result = inspect_target(target, query="ERROR", context=1, max_items=3, raw_sink=sink)

    assert result["match_count"] == 30
    assert len(result["matches"]) == 3
    assert result["raw_ref"].startswith("aer://sha256/")
    assert len(json.loads(stored[0])) == 30
    assert result["matches"][1]["before"][0]["line"] == 1


def test_default_inspection_result_stays_within_output_budget(tmp_path: Path) -> None:
    target = tmp_path / "large.log"
    long_value = "x" * 500
    target.write_text(
        "\n".join(f"ERROR {number} {long_value}" for number in range(100)), encoding="utf-8"
    )
    stored: list[bytes] = []

    def sink(data: bytes, name: str) -> str:
        stored.append(data)
        return "aer://sha256/" + "c" * 64

    result = inspect_target(target, query="ERROR", context=3, raw_sink=sink)
    encoded = json.dumps(result, ensure_ascii=False, separators=(",", ":")).encode()

    assert len(encoded) <= 16 * 1024
    assert result["truncated"] is True
    assert result["raw_ref"].startswith("aer://sha256/")
    assert stored


def test_text_line_range_and_regex_safety(tmp_path: Path) -> None:
    target = tmp_path / "notes.txt"
    target.write_text("one\ntwo\nthree\nfour\n", encoding="utf-8")

    result = inspect_target(target, start_line=2, end_line=3)
    assert [item["text"] for item in result["preview"]] == ["two", "three"]

    with pytest.raises(AerError, match="high-cost") as captured:
        inspect_target(target, query="(a+)+$", regex=True)
    assert captured.value.code == "INVALID_ARGUMENT"

    with pytest.raises(AerError, match="high-cost"):
        inspect_target(target, query="(a|aa)+$", regex=True)


def test_ambiguous_inspect_regex_is_stopped_by_runtime_timeout(tmp_path: Path) -> None:
    target = tmp_path / "hostile.txt"
    target.write_text("a" * 20_000 + "!\n", encoding="utf-8")

    with pytest.raises(AerError) as captured:
        inspect_target(target, query="(?:a{1,2})+$", regex=True)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert "timeout" in captured.value.message


@pytest.mark.parametrize(
    ("suffix", "payload"),
    [
        (".json", "[" * 150 + "0" + "]" * 150),
        (".yaml", "[" * 150 + "0" + "]" * 150),
    ],
)
def test_structured_inspection_rejects_excessive_nesting(
    tmp_path: Path, suffix: str, payload: str
) -> None:
    target = tmp_path / f"deep{suffix}"
    target.write_text(payload, encoding="utf-8")

    with pytest.raises(AerError) as captured:
        inspect_target(target)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert captured.value.details["limit"] == structured_inspect._MAX_STRUCTURED_DEPTH


@pytest.mark.parametrize(
    ("suffix", "payload"),
    [
        (".json", json.dumps(list(range(30)))),
        (".yaml", "items: [" + ",".join(str(value) for value in range(30)) + "]\n"),
    ],
)
def test_structured_inspection_rejects_excessive_nodes_before_traversal(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    suffix: str,
    payload: str,
) -> None:
    monkeypatch.setattr(structured_inspect, "_MAX_STRUCTURED_NODES", 20)
    target = tmp_path / f"wide{suffix}"
    target.write_text(payload, encoding="utf-8")

    with pytest.raises(AerError) as captured:
        inspect_target(target)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert captured.value.details["limit"] == 20


def test_text_query_streams_only_preview_records_and_preserves_exact_raw(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    target = tmp_path / "many.log"
    target.write_text("\n".join(f"ERROR item {index}" for index in range(5_000)), encoding="utf-8")
    stored: list[bytes] = []
    compact_calls = 0
    original_compact = text_inspect._compact_match_record

    def compact(record: dict[str, object]) -> dict[str, object]:
        nonlocal compact_calls
        compact_calls += 1
        return original_compact(record)

    def sink(data: bytes, _name: str) -> str:
        stored.append(data)
        return "aer://sha256/" + "e" * 64

    monkeypatch.setattr(text_inspect, "_compact_match_record", compact)
    result = inspect_target(target, query="ERROR", max_items=3, raw_sink=sink)

    assert result["match_count"] == 5_000
    assert len(result["matches"]) == 3
    assert compact_calls == 3
    assert len(json.loads(stored[0])) == 5_000


def test_text_query_hard_limit_returns_bounded_partial_raw(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    monkeypatch.setattr(text_inspect, "_MAX_TEXT_MATCHES", 5)
    target = tmp_path / "limited.log"
    target.write_text("\n".join(f"ERROR item {index}" for index in range(10)), encoding="utf-8")
    stored: list[bytes] = []

    def sink(data: bytes, _name: str) -> str:
        stored.append(data)
        return "aer://sha256/" + "f" * 64

    with pytest.raises(AerError) as captured:
        inspect_target(target, query="ERROR", max_items=2, raw_sink=sink)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert captured.value.details["observed_at_least"] == 6
    assert captured.value.details["partial_results"] == 5
    assert captured.value.raw_ref == "aer://sha256/" + "f" * 64
    assert len(json.loads(stored[0])) == 5


def test_text_query_encoded_byte_limit_returns_valid_partial_raw(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    monkeypatch.setattr(text_inspect, "_MAX_TEXT_MATCH_BYTES", 250)
    target = tmp_path / "large-matches.log"
    target.write_text("\n".join(["ERROR " + "x" * 100] * 3), encoding="utf-8")
    stored: list[bytes] = []

    def sink(data: bytes, _name: str) -> str:
        stored.append(data)
        return "aer://sha256/" + "a" * 64

    with pytest.raises(AerError) as captured:
        inspect_target(target, query="ERROR", max_items=2, raw_sink=sink)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert captured.value.details["encoded_bytes_limit"] == 250
    assert captured.value.details["partial_results"] == 1
    assert captured.value.raw_ref == "aer://sha256/" + "a" * 64
    assert len(json.loads(stored[0])) == 1


def test_json_pointer_slice_outline_and_query(tmp_path: Path) -> None:
    target = tmp_path / "result.json"
    target.write_text(
        json.dumps({"items": [{"status": "ok"}, {"status": "failed"}, {"status": "ok"}]}),
        encoding="utf-8",
    )

    scalar = inspect_target(target, selector="/items/1/status")
    sliced = inspect_target(target, selector="/items/1:3", max_items=5)
    searched = inspect_target(target, query="failed", outline=True, max_items=5)

    assert scalar["selection"] == "failed"
    assert sliced["selection"]["total"] == 2
    assert searched["matches"] == [{"path": "/items/1/status", "match": "value", "value": "failed"}]
    assert any(item["path"] == "/items" for item in searched["outline"])


def test_json_escaped_pointer_key_query_and_overflow(tmp_path: Path) -> None:
    target = tmp_path / "escaped.json"
    target.write_text(
        json.dumps({"a/b": {"~key": list(range(30))}, "needle-key": "value"}),
        encoding="utf-8",
    )
    stored: list[bytes] = []

    def sink(data: bytes, name: str) -> str:
        stored.append(data)
        return "aer://sha256/" + "d" * 64

    result = inspect_target(
        target,
        selector="/a~1b/~0key",
        query="needle",
        outline=True,
        max_items=2,
        raw_sink=sink,
    )

    assert result["selection"]["items"] == [0, 1]
    assert result["matches"][0]["match"] == "key"
    assert result["truncated"] is True
    assert stored


@pytest.mark.parametrize(
    "selector", ["items/0", "/items/nope", "/items/2:1", "/items/01", "/~2bad"]
)
def test_invalid_json_pointer_is_actionable(tmp_path: Path, selector: str) -> None:
    target = tmp_path / "data.json"
    target.write_text('{"items":[1,2]}', encoding="utf-8")

    with pytest.raises(AerError) as captured:
        inspect_target(target, selector=selector)

    assert captured.value.code == "INVALID_SELECTOR"


def test_yaml_uses_safe_loader_and_alias_limit(tmp_path: Path) -> None:
    unsafe = tmp_path / "unsafe.yaml"
    unsafe.write_text("value: !!python/object/apply:os.system ['echo bad']\n", encoding="utf-8")
    aliases = tmp_path / "aliases.yaml"
    aliases.write_text("base: &x [1]\nitems: [" + ",".join("*x" for _ in range(101)) + "]\n")

    with pytest.raises(AerError) as unsafe_error:
        inspect_target(unsafe)
    with pytest.raises(AerError) as alias_error:
        inspect_target(aliases)

    assert unsafe_error.value.code == "CORRUPT_FILE"
    assert alias_error.value.code == "LIMIT_EXCEEDED"


def test_yaml_alias_cycle_is_bounded(tmp_path: Path) -> None:
    target = tmp_path / "cycle.yaml"
    target.write_text("root: &root\n  child: *root\n", encoding="utf-8")

    result = inspect_target(target, selector="/root", outline=True)

    assert result["type"] == "yaml"
    assert result["selection"]["items"]["child"]["type"] == "alias"


def test_csv_and_jsonl_are_bounded_and_selectable(tmp_path: Path) -> None:
    csv_path = tmp_path / "rows.csv"
    csv_path.write_text(
        "id,status,total\n"
        + "\n".join(
            f"{number},{'pending' if number % 2 else 'done'},{number * 10}" for number in range(30)
        ),
        encoding="utf-8",
    )
    jsonl_path = tmp_path / "rows.jsonl"
    jsonl_path.write_text('{"id":1,"state":"ok"}\n{"id":2,"state":"failed"}\n')

    csv_result = inspect_target(
        csv_path,
        selector="id,total",
        query="pending",
        rows="5:29",
    )
    jsonl_result = inspect_target(jsonl_path, query="failed")

    assert csv_result["row_count"] == 30
    assert len(csv_result["preview"]) <= 20
    assert set(csv_result["preview"][0]) == {"id", "total"}
    assert jsonl_result["preview"] == [{"id": 2, "state": "failed"}]


def test_tabular_duplicate_headers_tsv_and_invalid_jsonl(tmp_path: Path) -> None:
    tsv = tmp_path / "rows.tsv"
    tsv.write_text("id\tid\tflag\n1\t2\ttrue\n", encoding="utf-8")
    invalid = tmp_path / "bad.jsonl"
    invalid.write_text('{"ok":1}\nnot-json\n', encoding="utf-8")

    result = inspect_target(tsv, selector="id#2,flag")

    assert result["warnings"]
    assert result["preview"] == [{"id#2": "2", "flag": "true"}]
    assert result["columns"][2]["types"] == ["boolean"]
    with pytest.raises(AerError) as captured:
        inspect_target(invalid)
    assert captured.value.code == "CORRUPT_FILE"


def test_tabular_missing_column_is_rejected(tmp_path: Path) -> None:
    target = tmp_path / "rows.csv"
    target.write_text("id\n1\n", encoding="utf-8")

    with pytest.raises(AerError) as captured:
        inspect_target(target, selector="missing")

    assert captured.value.code == "INVALID_SELECTOR"
    assert captured.value.details["available"] == ["id"]


def test_sparse_jsonl_inspection_enforces_rectangular_cell_limit(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    target = tmp_path / "sparse.jsonl"
    target.write_text('{"a":1}\n{"b":2}\n{"c":3}\n', encoding="utf-8")
    monkeypatch.setattr("aer.inspect.tabular.MAX_TABULAR_CELLS", 8)

    with pytest.raises(AerError) as captured:
        inspect_target(target)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert captured.value.details == {
        "rows": 3,
        "columns": 3,
        "cells": 9,
        "limit": 8,
    }


def test_repository_excludes_ignored_secrets_and_binary(tmp_path: Path) -> None:
    (tmp_path / "keep.py").write_text("before\nneedle\nafter\n", encoding="utf-8")
    (tmp_path / "ignored.txt").write_text("needle\n", encoding="utf-8")
    (tmp_path / "git-ignored.txt").write_text("needle\n", encoding="utf-8")
    (tmp_path / ".env").write_text("TOKEN=needle\n", encoding="utf-8")
    (tmp_path / ".netrc").write_text("password needle\n", encoding="utf-8")
    (tmp_path / "secrets.yml").write_text("token: needle\n", encoding="utf-8")
    (tmp_path / "secrets.production.yaml").write_text("token: needle\n", encoding="utf-8")
    (tmp_path / "blob.bin").write_bytes(b"\x00needle")
    (tmp_path / ".aerignore").write_text("ignored.txt\n", encoding="utf-8")
    (tmp_path / ".gitignore").write_text("git-ignored.txt\n", encoding="utf-8")

    result = inspect_target(tmp_path, query="needle", context=1)

    assert result["file_match_counts"] == {"keep.py": 1}
    assert result["skipped_binary_or_undecodable"] == 1
    assert result["matches"][0]["before"][0]["text"] == "before"
    assert result["matches"][0]["after"][0]["text"] == "after"


@pytest.mark.skipif(shutil.which("rg") is None, reason="ripgrep unavailable")
def test_repository_uses_streaming_ripgrep_without_a_shell(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    (tmp_path / "one.txt").write_text("before\nneedle\nafter\n", encoding="utf-8")
    (tmp_path / "--literal.txt").write_text("needle\n", encoding="utf-8")
    original_popen = repository_inspect.subprocess.Popen
    calls: list[tuple[list[str], dict[str, object]]] = []

    def recorded_popen(command: list[str], **kwargs: object) -> subprocess.Popen[bytes]:
        calls.append((command, kwargs))
        return original_popen(command, **kwargs)  # type: ignore[arg-type, return-value]

    monkeypatch.setattr(repository_inspect.subprocess, "Popen", recorded_popen)

    result = inspect_target(tmp_path, query="needle", context=1)

    assert result["search_engine"] == "ripgrep"
    assert result["file_match_counts"] == {"--literal.txt": 1, "one.txt": 1}
    rg_calls = [call for call in calls if "--json" in call[0]]
    assert len(rg_calls) == 1
    assert rg_calls[0][0][0] == shutil.which("rg")
    assert rg_calls[0][0].index("--") < rg_calls[0][0].index("--literal.txt")
    assert rg_calls[0][1]["shell"] is False


def test_repository_falls_back_to_bounded_python_search(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    (tmp_path / "one.txt").write_text("needle\n", encoding="utf-8")
    monkeypatch.setattr(repository_inspect.shutil, "which", lambda _name: None)

    result = inspect_target(tmp_path, query="needle")

    assert result["search_engine"] == "python"
    assert result["match_count"] == 1


def test_repository_match_hard_limit_returns_partial_raw_reference(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    (tmp_path / "many.txt").write_text("needle\n" * 5, encoding="utf-8")
    monkeypatch.setattr(repository_inspect.shutil, "which", lambda _name: None)
    monkeypatch.setattr(repository_inspect, "_MAX_REPOSITORY_MATCHES", 3)
    stored: list[bytes] = []

    def sink(data: bytes, _name: str) -> str:
        stored.append(data)
        return "aer://sha256/" + "f" * 64

    with pytest.raises(AerError) as captured:
        inspect_target(tmp_path, query="needle", raw_sink=sink)

    assert captured.value.code == "LIMIT_EXCEEDED"
    assert captured.value.raw_ref == "aer://sha256/" + "f" * 64
    assert captured.value.details["observed_at_least"] == 4
    assert len(json.loads(stored[0])) == 3


def test_repository_overflow_preview_preserves_all_bounded_matches(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    (tmp_path / "many.txt").write_text("needle\n" * 30, encoding="utf-8")
    monkeypatch.setattr(repository_inspect.shutil, "which", lambda _name: None)
    stored: list[bytes] = []

    def sink(data: bytes, name: str) -> str:
        if name.endswith(".matches.json"):
            stored.append(data)
        return "aer://sha256/" + "b" * 64

    result = inspect_target(tmp_path, query="needle", max_items=2, raw_sink=sink)

    assert result["match_count"] == 30
    assert len(result["matches"]) == 2
    assert result["truncated"] is True
    assert len(json.loads(stored[0])) == 30


def test_git_repository_includes_tracked_then_untracked(tmp_path: Path) -> None:
    if subprocess.run(["git", "--version"], capture_output=True, check=False).returncode:
        pytest.skip("git unavailable")
    subprocess.run(["git", "init", "-q", "-b", "main"], cwd=tmp_path, check=True)
    (tmp_path / "tracked.txt").write_text("tracked\n", encoding="utf-8")
    subprocess.run(["git", "add", "tracked.txt"], cwd=tmp_path, check=True)
    (tmp_path / "untracked.txt").write_text("untracked\n", encoding="utf-8")
    (tmp_path / "linked.txt").symlink_to(tmp_path / "tracked.txt")
    subprocess.run(["git", "add", "linked.txt"], cwd=tmp_path, check=True)

    result = inspect_target(tmp_path, outline=True)

    assert result["git"] is True
    assert result["outline"].index("tracked.txt") < result["outline"].index("untracked.txt")
    assert "linked.txt" not in result["outline"]


def test_git_repository_changed_glob_and_regex(tmp_path: Path) -> None:
    if subprocess.run(["git", "--version"], capture_output=True, check=False).returncode:
        pytest.skip("git unavailable")
    subprocess.run(["git", "init", "-q", "-b", "main"], cwd=tmp_path, check=True)
    (tmp_path / "match.py").write_text("value = 42\n", encoding="utf-8")
    (tmp_path / "skip.txt").write_text("value = 42\n", encoding="utf-8")

    result = inspect_target(
        tmp_path,
        changed=True,
        glob="*.py",
        query=r"value\s*=\s*\d+",
        regex=True,
    )

    assert result["changed_only"] is True
    assert result["file_match_counts"] == {"match.py": 1}


def test_non_git_repository_rejects_changed_filter(tmp_path: Path) -> None:
    (tmp_path / "ordinary.txt").write_text("not tracked\n", encoding="utf-8")

    with pytest.raises(AerError) as captured:
        inspect_target(tmp_path, changed=True)

    assert captured.value.code == "INVALID_ARGUMENT"
    assert captured.value.target == str(tmp_path)
    assert captured.value.suggested_action == (
        "Run without --changed or initialize a Git repository."
    )


def test_repository_path_batches_preserve_each_path_once() -> None:
    paths = [f"folder/{index:03d}-{'x' * 300}.txt" for index in range(400)]

    batches = repository_inspect._path_batches(paths)

    assert len(batches) > 1
    assert [path for batch in batches for path in batch] == paths


def test_xlsx_summary_range_and_formula(tmp_path: Path) -> None:
    target = tmp_path / "book.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet["A1"] = "value"
    sheet["B1"] = "formula"
    sheet["A2"] = 3
    sheet["B2"] = "=A2*2"
    sheet.merge_cells("C1:D1")
    workbook.defined_names.add(DefinedName("sample", attr_text="'Summary'!$A$1:$B$2"))
    workbook.save(target)
    workbook.close()

    result = inspect_target(target, sheet="Summary", cell_range="A1:B2", formulas=True)

    assert result["sheets"][0]["formula_count"] == 1
    assert result["formulas"][0]["formula"] == "=A2*2"
    assert result["selection"]["rows"][1]["cells"][1]["cell"] == "B2"


def test_xlsx_selector_rows_and_invalid_sheet(tmp_path: Path) -> None:
    target = tmp_path / "book.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    for row in range(1, 6):
        sheet.cell(row, 1, row)
    workbook.save(target)
    workbook.close()

    selected = inspect_target(target, selector="Data!A2:A4")
    row_selected = inspect_target(target, sheet="Data", rows="3:5")

    assert selected["selection"]["rows"][0]["cells"][0]["value"] == 2
    assert row_selected["selection"]["range"] == "A3:A5"
    with pytest.raises(AerError) as captured:
        inspect_target(target, sheet="Missing", cell_range="A1")
    assert captured.value.code == "INVALID_SELECTOR"


def test_xlsx_selector_resolves_builder_stable_sheet_id(tmp_path: Path) -> None:
    target = tmp_path / "stable.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "KPI Dashboard"
    sheet["A1"] = "tokens"
    workbook.defined_names.add(
        DefinedName("aer_sheet_financial_summary", attr_text="'KPI Dashboard'!$A$1")
    )
    workbook.save(target)
    workbook.close()

    selected = inspect_target(target, selector="sheet:id=financial-summary/cell:A1")

    assert selected["selection"]["sheet"] == "KPI Dashboard"
    assert selected["selection"]["rows"][0]["cells"][0]["value"] == "tokens"


def test_xlsx_stable_selector_wins_over_conflicting_display_name(tmp_path: Path) -> None:
    spec = tmp_path / "conflicting-workbook.json"
    spec.write_text(
        json.dumps(
            {
                "version": 1,
                "kind": "workbook",
                "sheets": [
                    {"id": "stable", "name": "Dashboard", "columns": ["value"], "rows": [[1]]},
                    {"id": "other", "name": "stable", "columns": ["value"], "rows": [[2]]},
                ],
            }
        ),
        encoding="utf-8",
    )
    target = tmp_path / "conflicting.xlsx"
    build_artifact(spec, target)

    selected = inspect_target(target, selector="sheet:id=stable/cell:A2")

    assert selected["selection"]["sheet"] == "Dashboard"
    assert selected["selection"]["rows"][0]["cells"][0]["value"] == 1


def test_built_xlsx_stable_cell_selector_and_encoded_ids_validate(tmp_path: Path) -> None:
    spec = tmp_path / "workbook.json"
    spec.write_text(
        json.dumps(
            {
                "version": 1,
                "kind": "workbook",
                "sheets": [
                    {
                        "id": "financial summary.v1",
                        "name": "KPI Dashboard",
                        "cells": [
                            {
                                "id": "token.value current",
                                "address": "C4",
                                "value": 63,
                            }
                        ],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    target = tmp_path / "stable.xlsx"
    build_artifact(spec, target)

    selected = inspect_target(
        target,
        selector="sheet:id=financial summary.v1/cell:id=token.value current",
    )

    assert selected["selection"]["sheet"] == "KPI Dashboard"
    assert selected["selection"]["range"] == "C4:C4"
    assert selected["selection"]["rows"][0]["cells"][0]["value"] == 63
    assert validate_file(target)["valid"] is True


def test_pptx_summary_stable_ids_query_and_selector(tmp_path: Path) -> None:
    target = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.element.cSld.set("name", "aer:market")
    title = slide.shapes.title
    assert title is not None
    title.text = "Market Analysis"
    title.name = "aer:market/title"
    presentation.save(target)

    result = inspect_target(target, selector="slide:id=market/shape:id=title", query="Market")

    assert result["slides"][0]["id"] == "market"
    assert result["selection"]["id"] == "title"
    assert result["match_count"] == 1


def test_pptx_slide_number_title_and_invalid_shape(tmp_path: Path) -> None:
    target = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    assert slide.shapes.title is not None
    slide.shapes.title.text = "Title"
    presentation.save(target)

    selected = inspect_target(target, slide=1)
    title = inspect_target(target, selector="slide:1/title")

    assert selected["selection"]["shapes"]
    assert title["selection"]["title"] == "Title"
    with pytest.raises(AerError) as captured:
        inspect_target(target, selector="slide:1/shape:id=missing")
    assert captured.value.code == "INVALID_SELECTOR"


def test_docx_summary_outline_query_and_selector(tmp_path: Path) -> None:
    target = tmp_path / "report.docx"
    document = Document()
    document.add_heading("Conclusion", level=1)
    document.add_paragraph("The exact decision is retain.")
    document.add_table(rows=1, cols=2).cell(0, 0).text = "Evidence"
    document.save(target)

    result = inspect_target(
        target,
        outline=True,
        selector="section:id=conclusion",
        query="retain",
    )

    assert result["paragraph_count"] == 2
    assert result["outline"][0]["text"] == "Conclusion"
    assert result["selection"]["paragraphs"][1]["text"].endswith("retain.")
    assert result["match_count"] == 1


def test_docx_paragraph_table_and_block_selectors(tmp_path: Path) -> None:
    target = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("Alpha")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Needle"
    table.cell(0, 1).text = "Value"
    document.save(target)

    paragraph = inspect_target(target, selector="paragraph:1")
    block_id = paragraph["selection"]["id"]
    block = inspect_target(target, selector=f"block:id={block_id}")
    table_result = inspect_target(target, selector="table:1", query="Needle")

    assert block["selection"]["text"] == "Alpha"
    assert table_result["selection"]["rows"] == [["Needle", "Value"]]
    assert table_result["matches"][0]["table"] == 1


def test_built_docx_table_stable_block_selector_resolves(tmp_path: Path) -> None:
    spec = tmp_path / "document.json"
    spec.write_text(
        json.dumps(
            {
                "version": 1,
                "kind": "document",
                "content": [
                    {
                        "id": "summary-table",
                        "type": "table",
                        "headers": ["metric", "value"],
                        "rows": [["tokens", "63%"]],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    target = tmp_path / "report.docx"
    build_artifact(spec, target)

    selected = inspect_target(target, selector="block:id=summary-table")

    assert selected["selection"]["id"] == "summary-table"
    assert selected["selection"]["id_source"] == "stable"
    assert selected["selection"]["rows"] == [["metric", "value"], ["tokens", "63%"]]


def test_pdf_summary_and_selected_page(tmp_path: Path) -> None:
    target = tmp_path / "sample.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with target.open("wb") as handle:
        writer.write(handle)

    result = inspect_target(target, page=1)

    assert result["page_count"] == 1
    assert result["pages"][0]["width_points"] == 612.0
    assert result["selection"]["page"] == 1


def test_pdf_inspect_rejects_oversized_input_before_opening(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    target = tmp_path / "oversized.pdf"
    target.write_bytes(b"%PDF-1.7\n")
    monkeypatch.setattr(pdf_safety, "MAX_PDF_INPUT_BYTES", target.stat().st_size - 1)

    with pytest.raises(AerError) as caught:
        inspect_target(target, page=1)

    assert caught.value.code == "LIMIT_EXCEEDED"
    assert caught.value.operation == "inspect"


def test_encrypted_pdf_omits_content(tmp_path: Path) -> None:
    target = tmp_path / "encrypted.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    writer.encrypt("password")
    with target.open("wb") as handle:
        writer.write(handle)

    result = inspect_target(target)

    assert result["encrypted"] is True
    assert result["page_count"] is None
    assert result["content_omitted"] is True


def test_aer_reference_uses_resolver_hook(tmp_path: Path) -> None:
    target = tmp_path / "stored.txt"
    target.write_text("stored content\n", encoding="utf-8")
    reference = "aer://sha256/" + "b" * 64

    def resolver(value: str) -> Path:
        return target if value == reference else Path("missing")

    result = inspect_target(reference, resolver=resolver)

    assert result["target"] == reference
    assert result["preview"][0]["text"] == "stored content"

    inspector = Inspector(resolver=resolver)
    assert inspector.inspect(reference)["line_count"] == 1


def test_dispatch_limits_and_unsupported_format_are_stable(tmp_path: Path) -> None:
    target = tmp_path / "payload.unknown"
    target.write_bytes(b"payload")

    with pytest.raises(AerError) as unsupported:
        inspect_target(target)
    with pytest.raises(AerError) as no_resolver:
        inspect_target("aer://sha256/" + "e" * 64)
    with pytest.raises(AerError) as invalid_depth:
        inspect_target(target, max_depth=0)

    assert unsupported.value.code == "UNSUPPORTED_FORMAT"
    assert no_resolver.value.code == "INVALID_ARGUMENT"
    assert invalid_depth.value.code == "INVALID_ARGUMENT"


@pytest.mark.parametrize("suffix", [".xlsx", ".pptx", ".docx", ".pdf"])
def test_corrupt_structured_artifact_is_rejected(tmp_path: Path, suffix: str) -> None:
    target = tmp_path / f"bad{suffix}"
    target.write_bytes(b"not-a-valid-artifact")

    with pytest.raises(AerError) as captured:
        inspect_target(target)

    assert captured.value.code == "CORRUPT_FILE"


def test_sensitive_file_and_corrupt_input_do_not_leak(tmp_path: Path) -> None:
    secret = tmp_path / ".env"
    secret.write_text("OPENAI_API_KEY=sk-secret-value\n", encoding="utf-8")
    corrupt = tmp_path / "bad.json"
    corrupt.write_text("{not-json", encoding="utf-8")

    secret_result = inspect_target(secret)
    assert secret_result["content_omitted"] is True
    assert "sk-secret-value" not in json.dumps(secret_result)

    with pytest.raises(AerError) as captured:
        inspect_target(corrupt)
    assert captured.value.code == "CORRUPT_FILE"
