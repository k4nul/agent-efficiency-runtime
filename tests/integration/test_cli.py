from __future__ import annotations

import hashlib
import json
import os
import subprocess
import sys
from collections.abc import Callable
from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader, PdfWriter

PROJECT_ROOT = Path(__file__).resolve().parents[2]
EXAMPLES = PROJECT_ROOT / "examples"
Cli = Callable[..., subprocess.CompletedProcess[str]]


@pytest.fixture
def cli(tmp_path: Path) -> Cli:
    home = tmp_path / "aer-home"

    def run(
        *args: str,
        cwd: Path | None = None,
        stdin: str | None = None,
        env: dict[str, str] | None = None,
    ) -> subprocess.CompletedProcess[str]:
        environment = os.environ.copy()
        environment.update(
            {
                "AER_HOME": str(home),
                "NO_COLOR": "1",
                "PYTHONIOENCODING": "utf-8",
            }
        )
        if env:
            environment.update(env)
        completed = subprocess.run(
            [sys.executable, "-m", "aer", *args],
            cwd=cwd or tmp_path,
            env=environment,
            input=stdin,
            capture_output=True,
            text=True,
            shell=False,
            timeout=60,
            check=False,
        )
        assert "\x1b" not in completed.stdout
        assert "\x1b" not in completed.stderr
        return completed

    return run


def payload(completed: subprocess.CompletedProcess[str]) -> dict[str, object]:
    assert completed.stdout, completed.stderr
    parsed = json.loads(completed.stdout)
    assert isinstance(parsed, dict)
    return parsed


def successful(completed: subprocess.CompletedProcess[str]) -> dict[str, object]:
    parsed = payload(completed)
    assert completed.returncode == 0, parsed
    assert parsed["ok"] is True
    return parsed


def result(completed: subprocess.CompletedProcess[str]) -> dict[str, object]:
    parsed = successful(completed)
    value = parsed["result"]
    assert isinstance(value, dict)
    return value


def test_version_discovery_schema_and_invalid_args_are_compact(cli: Cli) -> None:
    version = successful(cli("--version"))
    assert version["operation"] == "version"
    assert version["result"] == {"version": "0.1.1"}

    discovered = cli("discover", "ppt patch")
    discovered_payload = successful(discovered)
    assert len(discovered.stdout.encode()) <= 2 * 1024
    assert discovered_payload["result"]["capabilities"][0]["name"] == "presentation.patch"

    compact_schema = cli("schema", "presentation.patch", "--compact", "--example")
    schema_payload = successful(compact_schema)
    assert len(compact_schema.stdout.encode()) <= 4 * 1024
    assert "input_schema" not in schema_payload["result"]

    invalid = cli("schema")
    invalid_payload = payload(invalid)
    assert invalid.returncode == 2
    assert invalid_payload["ok"] is False
    assert invalid_payload["code"] == "INVALID_ARGUMENT"
    assert "Traceback" not in invalid.stdout
    assert "Traceback" not in invalid.stderr


def test_store_put_deduplicates_cat_get_verify_and_inspect_ref(cli: Cli, tmp_path: Path) -> None:
    source = tmp_path / "source.json"
    source.write_text('{"items":[{"status":"ok"},{"status":"failed"}]}\n', encoding="utf-8")

    first = result(cli("store", "put", str(source)))
    second = result(cli("store", "put", str(source)))
    reference = str(first["ref"])

    assert reference == second["ref"]
    assert reference.startswith("aer://sha256/")
    listed = result(cli("store", "list"))
    assert len(listed["objects"]) == 1

    selected = result(cli("store", "cat", reference, "--start-line", "1", "--end-line", "1"))
    assert "failed" in selected["text"]
    assert result(cli("store", "verify", reference))["verified"] is True

    restored = tmp_path / "restored.json"
    result(cli("store", "get", reference, "-o", str(restored)))
    assert restored.read_bytes() == source.read_bytes()

    local = result(cli("inspect", str(source), "--selector", "/items/1/status"))
    stored = result(cli("inspect", reference, "--selector", "/items/1/status"))
    assert local["selection"] == "failed"
    assert stored["selection"] == "failed"
    assert stored["target"] == reference


def test_compact_runner_success_and_failure_exit_codes(cli: Cli) -> None:
    success_script = "for i in range(5000): print(f'ok {i}')"
    success = cli("run", "--", sys.executable, "-c", success_script)
    success_payload = successful(success)
    assert len(success.stdout.encode()) <= 16 * 1024
    assert success_payload["result"]["summary"] == "ok 4999"
    assert success_payload["result"]["raw_ref"].startswith("aer://sha256/")
    assert "ok 0" not in success.stdout

    failure_script = (
        "import sys; "
        "[print(f'noise {i}') for i in range(5000)]; "
        "print('ASSERTION_ERROR exact failure', file=sys.stderr); "
        "sys.exit(7)"
    )
    failed = cli("run", "--", sys.executable, "-c", failure_script)
    failed_payload = payload(failed)
    assert failed.returncode == 8
    assert failed_payload["ok"] is False
    assert failed_payload["code"] == "COMMAND_FAILED"
    assert len(failed.stdout.encode()) <= 16 * 1024
    assert failed_payload["raw_ref"].startswith("aer://sha256/")
    details = failed_payload["details"]
    assert isinstance(details, dict)
    assert any("ASSERTION_ERROR exact failure" in line for line in details["failure_context"])
    assert len(details["failure_context"]) <= 80


def test_data_query_ten_thousand_rows_has_twenty_row_preview(cli: Cli, tmp_path: Path) -> None:
    source = tmp_path / "large.csv"
    source.write_text(
        "id,status,total\n"
        + "\n".join(f"{index},pending,{index * 10}" for index in range(10_000))
        + "\n",
        encoding="utf-8",
    )

    queried = cli(
        "data",
        "query",
        str(source),
        "--where",
        "status == pending",
        "--select",
        "id,total",
    )
    queried_payload = successful(queried)
    queried_result = queried_payload["result"]

    assert queried_result["matched_rows"] == 10_000
    assert len(queried_result["preview"]) == 20
    assert queried_result["result_ref"].startswith("aer://sha256/")
    assert len(queried.stdout.encode()) <= 16 * 1024


@pytest.mark.parametrize(
    ("kind", "extension"),
    [("presentation", ".pptx"), ("document", ".docx"), ("workbook", ".xlsx")],
)
def test_build_reopen_patch_preserve_and_validate_office(
    cli: Cli, tmp_path: Path, kind: str, extension: str
) -> None:
    output = tmp_path / f"artifact{extension}"
    build = cli("build", str(EXAMPLES / f"{kind}.yaml"), "-o", str(output))
    build_payload = successful(build)
    build_result = build_payload["result"]

    assert build_payload["operation"] == f"{kind}.build"
    assert isinstance(build_result, dict)
    assert "operation" not in build_result
    assert output.is_file()
    assert Path(str(build_result["manifest"])).is_file()
    assert len(build.stdout.encode()) <= 16 * 1024

    if kind == "presentation":
        before = Presentation(str(output))
        before_text = {
            shape.name: shape.text
            for slide in before.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        }
        assert "aer:metrics/token-reduction-value" in before_text
    elif kind == "document":
        before_doc = Document(str(output))
        preserved_title = before_doc.paragraphs[0].text
        preserved_tables = len(before_doc.tables)
    else:
        before_book = load_workbook(output, data_only=False)
        preserved_formula = before_book["Summary"]["E2"].value
        before_book.close()

    patched = cli(
        "patch",
        str(output),
        "--spec",
        str(EXAMPLES / "patches" / f"{kind}.yaml"),
    )
    patch_result = result(patched)
    assert patch_result["before_sha256"] != patch_result["after_sha256"]
    assert len(patched.stdout.encode()) <= 16 * 1024

    if kind == "presentation":
        after = Presentation(str(output))
        after_text = {
            shape.name: shape.text
            for slide in after.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        }
        assert after_text["aer:metrics/token-reduction-value"] == "63%"
        assert {
            key: value
            for key, value in before_text.items()
            if key != "aer:metrics/token-reduction-value"
        } == {
            key: value
            for key, value in after_text.items()
            if key != "aer:metrics/token-reduction-value"
        }
    elif kind == "document":
        after_doc = Document(str(output))
        assert preserved_title == after_doc.paragraphs[0].text
        assert preserved_tables == len(after_doc.tables)
        assert any("더 작은 컨텍스트" in paragraph.text for paragraph in after_doc.paragraphs)
    else:
        after_book = load_workbook(output, data_only=False)
        assert after_book["Summary"]["B3"].value == 16000
        assert after_book["Summary"]["E2"].value == preserved_formula == "=B3/B2"
        after_book.close()

    validated = result(cli("validate", str(output)))
    assert validated["valid"] is True
    assert validated["automatic_checks_only"] is True


def test_build_svg_chart_with_validation(cli: Cli, tmp_path: Path) -> None:
    data = tmp_path / "chart.csv"
    data.write_text("label,value\nA,1\nB,2\n", encoding="utf-8")
    spec = tmp_path / "chart.yaml"
    spec.write_text(
        "version: 1\nkind: chart\ntype: bar\nsource: chart.csv\nx: label\ny: value\n",
        encoding="utf-8",
    )
    output = tmp_path / "chart.svg"

    built = result(cli("build", str(spec), "--validate", "-o", str(output)))

    assert output.is_file()
    assert built["validation"]["valid"] is True
    assert built["validation"]["checks"]["root"] == "svg"


def test_docx_remove_block_validate_cli_keeps_manifest_consistent(cli: Cli, tmp_path: Path) -> None:
    spec = tmp_path / "logical-list.yaml"
    spec.write_text(
        """version: 1
kind: document
content:
  - id: before
    type: paragraph
    text: Keep before
  - id: items
    type: bullets
    items: [one, two, three]
  - id: after
    type: paragraph
    text: Keep after
""",
        encoding="utf-8",
    )
    target = tmp_path / "logical-list.docx"
    result(cli("build", str(spec), "-o", str(target)))
    patch = tmp_path / "remove-list.yaml"
    patch.write_text(
        """version: 1
operations:
  - op: docx.remove_block
    target: block:id=items
""",
        encoding="utf-8",
    )

    patched = result(cli("patch", str(target), "--spec", str(patch), "--validate"))
    validated = result(cli("validate", str(target)))

    assert patched["validation"]["valid"] is True
    assert validated["valid"] is True
    assert [paragraph.text for paragraph in Document(target).paragraphs] == [
        "Keep before",
        "Keep after",
    ]
    manifest = json.loads(target.with_name(target.name + ".aer.json").read_text(encoding="utf-8"))
    assert [element["selector"] for element in manifest["elements"]] == [
        "block:id=before",
        "block:id=after",
    ]


def test_stale_expected_hash_rejects_patch_without_mutation(cli: Cli, tmp_path: Path) -> None:
    target = tmp_path / "data.json"
    target.write_text('{"status":"old"}\n', encoding="utf-8")
    patch = tmp_path / "patch.yaml"
    patch.write_text(
        "version: 1\noperations:\n  - op: json.set\n    target: /status\n    value: new\n",
        encoding="utf-8",
    )
    before = target.read_bytes()

    failed = cli(
        "patch",
        str(target),
        "--spec",
        str(patch),
        "--expected-sha256",
        "0" * 64,
    )
    failed_payload = payload(failed)

    assert failed.returncode == 6
    assert failed_payload["code"] == "HASH_MISMATCH"
    assert target.read_bytes() == before
    assert "Traceback" not in failed.stdout + failed.stderr


def test_state_lifecycle_persists_exact_values(cli: Cli, tmp_path: Path) -> None:
    successful(cli("state", "init", "task-1", "--goal", "Ship v0.1.0"))
    successful(
        cli(
            "state",
            "update",
            "task-1",
            "--complete",
            "build",
            "--complete",
            "build",
            "--remaining",
            "publish",
            "--decision",
            "provider=TradingView",
            "--artifact",
            "repo=/exact/path",
        )
    )
    successful(cli("state", "checkpoint", "task-1"))

    state = result(cli("state", "show", "task-1"))
    assert state["task_id"] == "task-1"
    assert state["goal"] == "Ship v0.1.0"
    assert state["completed"] == ["build"]
    assert state["remaining"] == ["publish"]
    assert state["decisions"] == {"provider": "TradingView"}
    assert state["artifacts"] == {"repo": "/exact/path"}
    assert len(state["checkpoints"]) == 1

    exported = tmp_path / "state.yaml"
    result(cli("state", "export", "task-1", "-o", str(exported)))
    assert "TradingView" in exported.read_text(encoding="utf-8")


def test_recipe_project_package_dry_run_and_execution(cli: Cli, tmp_path: Path) -> None:
    source = tmp_path / "project"
    source.mkdir()
    (source / "README.md").write_text("project\n", encoding="utf-8")
    output = tmp_path / "project.zip"
    variables = ("--var", f"source={source}", "--var", f"output={output}")

    dry_run = result(cli("recipe", "run", "project-package", *variables, "--dry-run"))
    assert dry_run["dry_run"] is True
    assert len(dry_run["steps"]) == 2
    assert not output.exists()

    executed = result(cli("recipe", "run", "project-package", *variables))
    assert executed["success"] is True
    assert executed["cache_hit"] is False
    assert executed["log_ref"].startswith("aer://sha256/")
    assert output.is_file()
    assert executed["steps"][-1]["result"]["valid"] is True

    cached_payload = successful(cli("recipe", "run", "project-package", *variables))
    assert cached_payload["result"]["cache_hit"] is True
    assert cached_payload["metrics"]["cache_hit"] is True


def test_doctor_treats_missing_optional_tools_as_nonfatal(cli: Cli, tmp_path: Path) -> None:
    empty_path = tmp_path / "empty-bin"
    empty_path.mkdir()

    checked = cli("doctor", env={"PATH": str(empty_path)})
    checked_payload = successful(checked)
    checks = checked_payload["result"]["checks"]
    unavailable = [item for item in checks if not item["required"] and not item["ok"]]

    assert unavailable
    assert all(item["status"] == "unavailable" for item in unavailable)
    assert len(checked.stdout.encode()) <= 16 * 1024


def test_image_inspect_resize_and_fit(cli: Cli, tmp_path: Path) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (80, 40), "navy").save(source)

    inspected = result(cli("image", "inspect", str(source)))
    assert (inspected["width"], inspected["height"]) == (80, 40)

    resized = tmp_path / "resized.png"
    resized_result = result(
        cli("image", "resize", str(source), "--width", "40", "-o", str(resized))
    )
    assert (resized_result["width"], resized_result["height"]) == (40, 20)
    with Image.open(resized) as image:
        assert image.size == (40, 20)

    fitted = tmp_path / "fitted.webp"
    result(cli("image", "fit", str(source), "--ratio", "1:1", "--width", "32", "-o", str(fitted)))
    with Image.open(fitted) as image:
        assert image.size == (32, 32)


def test_pdf_inspect_merge_extract_and_split(cli: Cli, tmp_path: Path) -> None:
    first = tmp_path / "first.pdf"
    second = tmp_path / "second.pdf"
    for path, size in ((first, (200, 300)), (second, (400, 500))):
        writer = PdfWriter()
        writer.add_blank_page(width=size[0], height=size[1])
        with path.open("wb") as handle:
            writer.write(handle)

    inspected = result(cli("pdf", "inspect", str(first), "--page", "1"))
    assert inspected["page_count"] == 1
    assert inspected["page"]["width"] == 200.0

    merged = tmp_path / "merged.pdf"
    merged_result = result(cli("pdf", "merge", str(first), str(second), "-o", str(merged)))
    assert merged_result["page_count"] == 2
    assert len(PdfReader(merged).pages) == 2

    excerpt = tmp_path / "excerpt.pdf"
    extract_result = result(cli("pdf", "extract", str(merged), "--pages", "2", "-o", str(excerpt)))
    assert extract_result["pages"] == [2]
    assert len(PdfReader(excerpt).pages) == 1

    split_dir = tmp_path / "pages"
    split_result = result(cli("pdf", "split", str(merged), "--out-dir", str(split_dir)))
    assert split_result["count"] == 2
    assert len(list(split_dir.glob("*.pdf"))) == 2


def test_archive_create_list_verify_and_byte_determinism(cli: Cli, tmp_path: Path) -> None:
    source = tmp_path / "delivery"
    (source / "nested").mkdir(parents=True)
    (source / "a.txt").write_text("alpha\n", encoding="utf-8")
    (source / "nested" / "b.txt").write_text("beta\n", encoding="utf-8")
    first = tmp_path / "first.zip"
    second = tmp_path / "second.zip"

    first_result = result(cli("archive", "create", str(source), "-o", str(first)))
    second_result = result(cli("archive", "create", str(source), "-o", str(second)))

    assert first.read_bytes() == second.read_bytes()
    assert first_result["sha256"] == second_result["sha256"]
    assert hashlib.sha256(first.read_bytes()).hexdigest() == first_result["sha256"]

    listed = result(cli("archive", "list", str(first)))
    verified = result(cli("archive", "verify", str(first)))
    assert listed["entry_count"] == 3
    assert {item["path"] for item in listed["entries"]} == {
        "a.txt",
        "nested/b.txt",
        "manifest.json",
    }
    assert verified["valid"] is True


def test_profile_and_benchmark_cli_persist_export_and_measure(cli: Cli, tmp_path: Path) -> None:
    for variant, input_tokens in (("direct", "1000"), ("aer", "300")):
        recorded = result(
            cli(
                "profile",
                "record",
                "--task",
                "ppt-generation",
                "--variant",
                variant,
                "--model-calls",
                "2",
                "--tool-calls",
                "3",
                "--input-tokens",
                input_tokens,
                "--output-tokens",
                "100",
                "--success",
                "true",
            )
        )
        assert recorded["measurement_source"] == "caller_supplied_unverified"
        assert recorded["estimate_classification"] == "not_recorded"
        assert recorded["provider_billed_tokens_known"] is False

    comparison = result(cli("profile", "compare", "--task", "ppt-generation"))
    assert set(comparison["variants"]) == {"direct", "aer"}
    profile_csv = tmp_path / "profiles.csv"
    exported = result(cli("profile", "export", "-o", str(profile_csv)))
    assert exported["rows"] == 2
    assert profile_csv.is_file()

    benchmark_json = tmp_path / "benchmark.json"
    benchmark = result(
        cli(
            "benchmark",
            "run",
            "--scenario",
            "json-patch",
            "-o",
            str(benchmark_json),
        )
    )
    assert benchmark["success"] is True
    assert benchmark["not_provider_billed_tokens"] is True
    assert benchmark_json.is_file()
    report_json = tmp_path / "benchmark-report.json"
    report = result(cli("benchmark", "report", "--limit", "1", "-o", str(report_json)))
    assert len(report["runs"]) == 1
    assert report_json.is_file()


def test_cli_parse_errors_are_compact_and_recoverable(cli: Cli) -> None:
    completed = cli("--" + "x" * 20_000)

    assert completed.returncode == 2
    assert len(completed.stdout.encode("utf-8")) <= 16 * 1024
    payload = json.loads(completed.stdout)
    assert payload["code"] == "INVALID_ARGUMENT"
    assert payload["raw_ref"].startswith("aer://sha256/")
